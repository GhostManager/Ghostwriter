
# Standard Libraries
import os

# Django Imports
from django.conf import settings

# 3rd Party Libraries
from pptx.dml.color import RGBColor as PptxRGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

# Ghostwriter Libraries
from ghostwriter.modules.reportwriter.extensions import IMAGE_EXTENSIONS, TEXT_EXTENSIONS
from ghostwriter.modules.reportwriter.richtext.ooxml import BaseHtmlToOOXML


class HtmlToPptx(BaseHtmlToOOXML):
    """
    Converts HTML to a powerpoint document
    """

    @staticmethod
    def delete_extra_paragraph(shape):
        """
        Shapes are created with an empty paragraph at the start, since they need at least one paragraph to be valid.
        The conversion code ignores it and simply appends.

        So remove it.
        """
        if len(shape.text_frame.paragraphs) <= 1:
            return
        prefix_par = shape.text_frame.paragraphs[0]._p
        prefix_par.getparent().remove(prefix_par)

    def __init__(self, slide, shape):
        super().__init__()
        self.slide = slide
        self.shape = shape

    def style_run(self, run, style):
        super().style_run(run, style)
        if style.get("hyperlink_url"):
            run.hyperlink.address = style["hyperlink_url"]
        if style.get("inline_code"):
            run.font.name = "Courier New"
        if style.get("highlight"):
            rPr = run._r.get_or_add_rPr()
            highlight = OxmlElement("a:highlight")
            srgbClr = OxmlElement("a:srgbClr")
            srgbClr.set("val", "FFFF00")
            highlight.append(srgbClr)
            rPr.append(highlight)
        if style.get("strikethrough"):
            run.font._element.set("strike", "sngStrike")
        if style.get("subscript"):
            run.font._element.set("baseline", "-25000")
        if style.get("superscript"):
            run.font._element.set("baseline", "30000")
        if style.get("font_size"):
            pass  # TODO: what's the pptx equivalent?
        if "font_color" in style:
            run.font.color.rgb = PptxRGBColor(*style["font_color"])

    def tag_br(self, el, *, par=None, **kwargs):
        self.text_tracking.new_block()
        if par is not None:
            par.add_line_break()

    def _tag_h(self, el, **kwargs):
        self.text_tracking.new_block()
        par = self.shape.text_frame.add_paragraph()
        run = par.add_run()
        run.text = el.text
        run.font.bold = True
    tag_h1 = _tag_h
    tag_h2 = _tag_h
    tag_h3 = _tag_h
    tag_h4 = _tag_h
    tag_h5 = _tag_h
    tag_h6 = _tag_h

    def tag_p(self, el, **kwargs):
        self.text_tracking.new_block()
        par = self.shape.text_frame.add_paragraph()

        par_classes = set(el.attrs.get("class", []))
        if "left" in par_classes:
            par.alignment = PP_ALIGN.LEFT
        if "center" in par_classes:
            par.alignment = PP_ALIGN.CENTER
        if "right" in par_classes:
            par.alignment = PP_ALIGN.RIGHT
        if "justify" in par_classes:
            par.alignment = PP_ALIGN.JUSTIFY

        self.process_children(el, par=par, **kwargs)

    def tag_pre(self, el, **kwargs):
        top = Inches(1.65)
        left = Inches(8)
        width = Inches(4.5)
        height = Inches(3)
        textbox = self.slide.shapes.add_textbox(left, top, width, height)
        self.text_tracking.new_block()
        par = textbox.text_frame.add_paragraph()
        par.alignment = PP_ALIGN.LEFT

        if len(el.contents) == 1 and el.contents[0].name == "code":
            content_el = el.contents[0]
        else:
            content_el = el

        self.text_tracking.in_pre = True
        try:
            self.process_children(content_el, par=par, **kwargs)
        finally:
            self.text_tracking.in_pre = False

        for run in par.runs:
            run.font.size = Pt(11)
            run.font.name = "Courier New"

    def tag_blockquote(self, el, **kwargs):
        par = self.shape.text_frame.add_paragraph()
        self.text_tracking.new_block()
        self.process_children(el.children, par=par, **kwargs)

    def tag_ul(self, el, *, par=None, list_level=None, **kwargs):
        if list_level is None:
            next_list_level = 1
        else:
            next_list_level = list_level + 1

        for child in el.children:
            if child.name != "li":
                # TODO: log
                continue
            self.text_tracking.new_block()
            par = self.shape.text_frame.add_paragraph()
            par.level = next_list_level
            self.process_children(child.children, par=par, list_level=next_list_level, **kwargs)

    tag_ol = tag_ul

    def create_table(self, rows, cols, **kwargs):
        return self.slide.shapes.add_table(rows=rows, cols=cols, left=Inches(10), top=Inches(5), width=Inches(3), height=Inches(1)).table

    def paragraph_for_table_cell(self, cell, td_el):
        return next(iter(cell.text_frame.paragraphs))


class HtmlToPptxWithEvidence(HtmlToPptx):
    """
    Augments `HtmlToPptx`, replacing marked spans with evidence figures, captions,
    and references.
    """
    def __init__(self, slide, shape, evidences):
        super().__init__(slide, shape)
        self.evidences = evidences

    def tag_span(self, el, *, par, **kwargs):
        if "data-gw-evidence" in el.attrs:
            try:
                evidence = self.evidences[int(el.attrs["data-gw-evidence"])]
            except (KeyError, ValueError):
                return
            par._p.getparent().remove(par._p)
            self.make_evidence(self.slide, evidence)
        elif "data-gw-caption" in el.attrs:
            ref_name = el.attrs["data-gw-caption"]
            if not ref_name:
                return
            run = par.add_run()
            run.text = f"See {ref_name}"
            run.font.italic = True
        elif "data-gw-ref" in el.attrs:
            ref_name = el.attrs["data-gw-ref"]
            if not ref_name:
                return
            run = par.add_run()
            run.text = f"See {ref_name}"
            run.font.italic = True
        else:
            super().tag_span(el, par=par, **kwargs)

    @staticmethod
    def make_evidence(slide, evidence):
        file_path = settings.MEDIA_ROOT + "/" + evidence["path"]
        if not os.path.exists(file_path):
            raise FileNotFoundError(file_path)

        extension = file_path.split(".")[-1].lower()
        if extension in TEXT_EXTENSIONS:
            with open(file_path, "r") as evidence_file:
                evidence_text = evidence_file.read()

            top = Inches(1.65)
            left = Inches(8)
            width = Inches(4.5)
            height = Inches(3)
            # Create new textbox, textframe, paragraph, and run
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            p = text_frame.paragraphs[0]
            run = p.add_run()
            # Insert evidence and apply formatting
            run.text = evidence_text
            font = run.font
            font.size = Pt(11)
            font.name = "Courier New"
        elif extension in IMAGE_EXTENSIONS:
            # Place new textbox to the mid-right
            top = Inches(1.65)
            left = Inches(8)
            width = Inches(4.5)
            slide.shapes.add_picture(file_path, left, top, width=width)
