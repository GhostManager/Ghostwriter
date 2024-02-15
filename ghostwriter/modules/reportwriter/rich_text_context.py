
# Standard Libraries
import html
import logging
import re
from typing import Any, Callable, Iterable

# 3rd Party Libraries
import docx
import jinja2
from bs4 import BeautifulSoup, NavigableString, Tag
from docx.enum.dml import MSO_THEME_COLOR_INDEX
from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_COLOR_INDEX
from docx.oxml.shared import OxmlElement, qn
from docx.shared import Inches, Pt, RGBColor
from markupsafe import Markup
from pptx.dml.color import RGBColor as PptxRGBColor
from pptx.enum.text import PP_ALIGN

logger = logging.getLogger(__name__)

# Allowlist for HTML tags expected to come from the WYSIWYG
TAG_ALLOWLIST = [
    "code",
    "span",
    "p",
    "ul",
    "ol",
    "li",
    "a",
    "em",
    "strong",
    "u",
    "b",
    "pre",
    "sub",
    "sup",
    "del",
    "blockquote",
]


class ReportConstants:
    """Constant values used for report generation."""

    DEFAULT_STYLE_VALUES = {
        "bold": False,
        "underline": False,
        "italic": False,
        "inline_code": False,
        "strikethrough": False,
        "font_family": None,
        "font_size": None,
        "font_color": None,
        "highlight": None,
        "superscript": False,
        "subscript": False,
        "hyperlink": False,
        "hyperlink_url": None,
    }


class RichTextConverter:
    """
    Runs jinja templating over HTML rich text and converts the result to OOXML.

    The HTML rich text comes from the TinyMCE WYSIWYG editors on various pages.

    This class encapsulate state for the entire report. Methods on this
    return sub-contexts annotated for various situations.
    """
    jinja_env: jinja2.Environment
    report_type: str
    report_json: dict
    base_context: dict
    base_evidences: dict
    make_figure: Callable[[Any, Any], Any]
    make_evidence: Any
    p_style: Any
    word_doc: Any

    def __init__(
        self,
        jinja_env: jinja2.Environment,
        report_type: str,
        report_json: dict,
        make_figure: Callable[[Any, Any], Any],
        make_evidence: Any,
        p_style=None,
        word_doc=None,
    ):
        self.jinja_env = jinja_env
        self.report_type = report_type
        self.report_json = report_json
        self.make_figure = make_figure
        self.make_evidence = make_evidence
        self.word_doc = word_doc
        self.p_style = p_style

        self.base_context = {
            # `{{.foo}}` converts to `{{obsolete.foo}}`
            "_old_dot_vars": {
                "client": report_json["client"]["short_name"] or report_json["client"]["name"],
                "project_start": report_json["project"]["start_date"],
                "project_end": report_json["project"]["end_date"],
                "project_type": report_json["project"]["type"].lower(),
            },
            "report": report_json,

            "evidence": _jinja_evidence,
            "caption": _jinja_caption,
            "ref": _jinja_ref,
        }
        self.base_evidences = {}
        for evidence in self.report_json["evidence"]:
            if evidence.get("friendly_name"):
                self.base_context["_old_dot_vars"][evidence["friendly_name"]] = _jinja_evidence(evidence["friendly_name"])
                self.base_evidences[evidence["friendly_name"]] = evidence

    def for_base(self):
        """
        Returns a `RichTextConverterWithContext` with a basic context, for report-level rich text
        """
        return RichTextConverterWithContext(self, self.base_context, self.base_evidences)

    def for_finding(self, finding):
        """
        Returns a `RichTextConverterWithContext` with a context suitable for rendering rich text for a finding
        """
        jinja_context = self.base_context | {
            "finding": finding,
            "_old_dot_vars": self.base_context["_old_dot_vars"].copy(),
        }
        evidences = self.base_evidences.copy()
        for evidence in finding["evidence"]:
            if evidence.get("friendly_name"):
                jinja_context["_old_dot_vars"][evidence["friendly_name"]] = _jinja_evidence(evidence["friendly_name"])
                evidences[evidence["friendly_name"]] = evidence
        return RichTextConverterWithContext(self, jinja_context, evidences)

    def render(
        self,
        *args,
        **kwargs,
    ):
        """
        Convenience for `self.for_base().render(...)`
        """
        return self.for_base().render(
            *args,
            **kwargs,
        )


class RichTextConverterWithContext:
    """
    `RichTextConverter` with a bound jinja context and evidences.
    """
    rtc: RichTextConverter
    jinja_context: dict
    evidences: dict

    def __init__(
        self,
        rtc: RichTextConverter,
        jinja_context: dict,
        evidences: dict,
    ):
        self.rtc = rtc
        self.jinja_context = jinja_context
        self.evidences = evidences

    def render(
        self,
        text,
        # Powerpoint doc parameters
        finding_slide=None,
        finding_body_shape=None,
    ):
        """
        Performs the template rendering and conversion.

        For word docs, the content is returned as a subdocument. For
        powerpoint docs, the content is added to the passed-in slide.
        """
        return RichTextConverterRenderer(
            self,
            finding_slide=finding_slide,
            finding_body_shape=finding_body_shape,
        ).render(text)

    @property
    def report_type(self) -> str:
        return self.rtc.report_type


class RichTextConverterRenderer:
    """
    Ephemeral class used by `RichTextConverterWithContext.render` containing
    the state for a single conversion.
    """
    rtcwc: RichTextConverterWithContext

    def __init__(
        self,
        rtcwc: RichTextConverterWithContext,
        # Powerpoint doc parameters
        finding_slide=None,
        finding_body_shape=None,
    ):
        self.rtcwc = rtcwc
        self.finding_slide = finding_slide
        self.finding_body_shape = finding_body_shape
        if rtcwc.rtc.word_doc is not None:
            self.sacrificial_doc = rtcwc.rtc.word_doc.new_subdoc()
        else:
            self.sacrificial_doc = None

    @property
    def report_type(self) -> str:
        return self.rtcwc.report_type

    @property
    def p_style(self):
        return self.rtcwc.rtc.p_style

    def render(
        self,
        text: str,
    ):
        if not text.strip():
            return ""

        # Replace old `{{.item}}`` syntax with jinja templates or elements to replace
        def replace_old_tag(match):
            contents = match.group(1).strip()
            # These will be swapped out when parsing the HTML
            if contents.startswith("ref "):
                return _jinja_ref(contents[4:].strip())
            elif contents.startswith("caption "):
                return _jinja_caption(contents[8:].strip())
            else:
                return "{{ _old_dot_vars[" + repr(contents.strip()) + "]}}"
        text_old_dot_subbed = re.sub(r"\{\{\.(.*?)\}\}", replace_old_tag, text)

        # Run rich text as template
        template = self.rtcwc.rtc.jinja_env.from_string(text_old_dot_subbed)
        text_rendered = template.render(self.rtcwc.jinja_context)

        # Filter out XML-incompatible characters
        text_char_filtered = "".join(c for c in text_rendered if _valid_xml_char_ordinal(c))

        # Parse the HTML into a BS4 soup object
        soup = BeautifulSoup(text_char_filtered, "lxml")

        # Each WYSIWYG field begins with ``<html><body>`` so get the contents of ``body``
        self._convert_top_level(soup.find("body"), self.p_style)

        return self.sacrificial_doc

    def _convert_top_level(self, body: Tag, p_style):
        """
        Converts the root element of the HTML.
        """
        for tag in body.children:
            # Get the HTML tag's name to determine next steps
            tag_name = tag.name

            # Hn – Headings
            if tag_name in ["h1", "h2", "h3", "h4", "h5", "h6"]:
                if self.report_type == "pptx":
                    # No headings in PPTX, so add a new line and bold it as a pseudo-heading
                    p = self.finding_body_shape.text_frame.add_paragraph()
                    run = p.add_run()
                    run.text = tag.text
                    font = run.font
                    font.bold = True
                else:
                    heading_num = int(tag_name[1])
                    # Add the heading to the document
                    # This discards any inline formatting, but that should be managed
                    # by editing the style in the template
                    p = self.sacrificial_doc.add_heading(tag.text, heading_num)

            # P – Paragraphs
            elif tag_name == "p":

                # Add a paragraph to the document based on doc type
                if self.report_type == "pptx":
                    p = self.finding_body_shape.text_frame.add_paragraph()
                    ALIGNMENT = PP_ALIGN
                else:
                    p = self.sacrificial_doc.add_paragraph(style=p_style)
                    ALIGNMENT = WD_ALIGN_PARAGRAPH

                # Check for alignment classes on the ``p`` tag
                if "class" in tag.attrs:
                    tag_attrs = tag.attrs["class"]
                    if "left" in tag_attrs:
                        p.alignment = ALIGNMENT.LEFT
                    if "center" in tag_attrs:
                        p.alignment = ALIGNMENT.CENTER
                    if "right" in tag_attrs:
                        p.alignment = ALIGNMENT.RIGHT
                    if "justify" in tag_attrs:
                        p.alignment = ALIGNMENT.JUSTIFY

                # Pass the contents and new paragraph on to drill down into nested formatting
                self._convert_inline(tag.children, p)

            # PRE – Code Blocks
            elif tag_name == "pre":
                # Get the list of pre-formatted strings
                contents = tag.contents

                # We do not style any text inside a code block, so we just write the XML
                # The only content should be one ``code`` block and a long line of text broken up by ``\r\n``
                if self.report_type == "pptx":
                    # Place new textbox to the mid-right to keep it out of the way
                    if contents:
                        top = Inches(1.65)
                        left = Inches(8)
                        width = Inches(4.5)
                        height = Inches(3)
                        # Create new textbox, textframe, paragraph, and run
                        textbox = self.finding_slide.shapes.add_textbox(left, top, width, height)
                        text_frame = textbox.text_frame
                        for content in contents:
                            for code in content:
                                parts = code.split("\r\n")
                                for code_line in parts:
                                    p = text_frame.add_paragraph()
                                    # Align left to anticipate a monospaced font for code
                                    p.alignment = PP_ALIGN.LEFT
                                    run = p.add_run()
                                    # Insert code block and apply formatting
                                    run.text = code_line
                                    font = run.font
                                    font.size = Pt(11)
                                    font.name = "Courier New"
                else:
                    if contents:
                        for content in contents:
                            for code in content:
                                parts = code.split("\r\n")
                                for code_line in parts:
                                    # Create paragraph and apply the ``CodeBlock`` style
                                    p = self.sacrificial_doc.add_paragraph(code_line)
                                    p.style = "CodeBlock"
                                    p.alignment = WD_ALIGN_PARAGRAPH.LEFT

            # OL & UL - Ordered/Numbered & Unordered Lists
            elif tag_name in ("ol", "ul"):
                # Ordered/numbered lists need numbers and linked paragraphs
                prev_p = None
                num = bool(tag_name == "ol")
                # In HTML, sub-items in a list are nested HTML lists
                # We need to check every list item for formatted and additional lists
                # While tracking which level of the list we are working with
                prev_p = self._parse_html_lists(tag, prev_p, num, 0)

            # BLOCKQUOTE - Blockquote Sections
            elif tag_name == "blockquote":
                # Get the tag's contents to check for additional formatting
                contents = tag.contents

                # PowerPoint lacks a blockquote style, so we just add a basic paragraph
                if self.report_type == "pptx":
                    p = self.finding_body_shape.text_frame.add_paragraph()
                else:
                    p = self.sacrificial_doc.add_paragraph()
                    p.style = "Blockquote"

                # Pass the contents and new paragraph on to drill down into nested formatting
                self._convert_inline(contents, p)
            else:
                if not isinstance(tag, NavigableString):
                    logger.warning(
                        "Encountered an unknown tag inside of the finding HTML: %s",
                        tag_name,
                    )

    def _convert_inline(self, contents: Iterable[Tag], par, styles=None):
        """
        Process BeautifulSoup4 ``Tag`` objects containing nested HTML tags.

        **Parameters**

        ``contents`` : Tag.contents
            Contents of a BS4 ``Tag``
        ``par`` : Paragraph
            Word docx ``Paragraph`` object
        ``styles`` : dict
            Override default styles with a provided dict
        """

        def merge_styles(run_styles, parent_styles):
            """Merge the two style dictionaries, keeping the ``run_styles`` values."""
            for key in parent_styles:
                # Same values in both dicts
                if run_styles[key] == parent_styles[key]:
                    pass
                else:
                    # Run has no value but parent has one
                    if not run_styles[key] and parent_styles[key]:
                        run_styles[key] = parent_styles[key]
            return run_styles

        def check_tags(tag, styles_dict):
            """Check the tag name and update the provided styles dictionary as needed."""
            tag_name = tag.name

            # A ``code`` tag here is inline code inside a ``p`` tag
            if tag_name == "code":
                styles_dict["inline_code"] = True

            # These tags apply italic, bold, and underline styles but appear rarely
            elif tag_name == "em":
                styles_dict["italic"] = True
            elif tag_name in ("strong", "b"):
                styles_dict["bold"] = True
            elif tag_name == "u":
                styles_dict["underline"] = True

            # The ``sub`` and ``sup`` tags designates subscript and superscript
            elif tag_name == "sub":
                styles_dict["subscript"] = True
            elif tag_name == "sup":
                styles_dict["superscript"] = True

            # The ``del`` tag applies a strikethrough style
            elif tag_name == "del":
                styles_dict["strikethrough"] = True

            # A ``span`` tag will usually contain one or more classes for formatting
            # Empty spans usually only appear in place of a non-breaking spaces
            elif tag_name == "span":
                # Check existence of supported classes for run styles
                if "class" in tag.attrs:
                    tag_attrs = tag.attrs["class"]
                    if "italic" in tag_attrs:
                        styles_dict["italic"] = True
                    if "bold" in tag_attrs:
                        styles_dict["bold"] = True
                    if "underline" in tag_attrs:
                        styles_dict["underline"] = True
                    if "highlight" in tag_attrs:
                        styles_dict["highlight"] = True

                # Check existence of supported character styles
                if "style" in tag.attrs:
                    tag_style = self._get_styles(tag)
                    if "font-size" in tag_style:
                        styles_dict["font_size"] = tag_style["font-size"]
                    if "font-family" in tag_style:
                        styles_dict["font_family"] = tag_style["font-family"]
                    if "color" in tag_style:
                        styles_dict["font_color"] = tag_style["color"]
                    if "background-color" in tag_style:
                        styles_dict["highlight"] = tag_style["background-color"]

            # An ``a`` tag is a hyperlink
            elif tag_name == "a":
                styles_dict["hyperlink"] = True
                styles_dict["hyperlink_url"] = tag["href"]

            # Any other tags are unexpected and ignored
            else:
                if tag_name not in TAG_ALLOWLIST:
                    logger.warning("Encountered an unexpected nested HTML tag: %s", tag_name)

            return styles_dict

        # Begin with HTML that could have indefinitely nested tags
        for part in contents:
            # Track the styles for this first parent tag or use the defaults if none
            if styles:
                parent_styles = styles.copy()
                run_styles = styles.copy()
            else:
                parent_styles = ReportConstants.DEFAULT_STYLE_VALUES.copy()
                run_styles = ReportConstants.DEFAULT_STYLE_VALUES.copy()

            # Get each part's ``name`` to check if it's a ``Tag`` object
            # A plain string will return ``None`` - no HTML tag
            part_name = part.name

            if part_name == "span" and "data-gw-evidence" in part.attrs:
                # Replace evidence placeholder
                evidence = self.rtcwc.evidences.get(part.attrs["data-gw-evidence"])
                if not evidence:
                    continue
                self._process_evidence(evidence, par)
            elif part_name == "span" and "data-gw-caption" in part.attrs:
                # Replace caption placeholder
                ref_name = part.attrs["data-gw-caption"]
                if self.report_type == "pptx":
                    if ref_name:
                        run = par.add_run()
                        run.text = f"See {ref_name}"
                        font = run.font
                        font.italic = True
                else:
                    par.style = "Caption"
                    if ref_name:
                        self._make_figure(par, ref_name)
                    else:
                        self._make_figure(par)
            elif part_name == "span" and "data-gw-ref" in part.attrs:
                # Replace reference holder
                ref_name = part.attrs["data-gw-ref"]
                if self.report_type == "pptx":
                    run = par.add_run()
                    run.text = f"See {ref_name}"
                    font = run.font
                    font.italic = True
                else:
                    self._make_cross_ref(par, ref_name)
            elif part_name:
                # Get the top-level's styles first as it applies to all future runs

                # Update styles based on the tag and properties
                parent_styles = check_tags(part, parent_styles)

                # Split part into list of text and ``Tag`` objects
                part_contents = part.contents

                # With parent styles recorded, process each child same as above
                for tag in part_contents:
                    tag_name = tag.name

                    if tag_name and tag_name in TAG_ALLOWLIST:
                        # Construct text to be written as part of this run
                        content_text = ""

                        # Split part into list of text and ``Tag`` objects
                        tag_contents = tag.contents

                        if len(tag_contents) >= 1:
                            # Update styles based on the tag and properties
                            run_styles = check_tags(tag, run_styles)

                            # Check for a hyperlink formatted with additional styles
                            if tag_contents[0].name:
                                if tag_contents[0].name == "a":
                                    run_styles["hyperlink"] = True
                                    run_styles["hyperlink_url"] = tag_contents[0]["href"]
                                    content_text = tag_contents[0].text

                            # Combine the styles to carry them over to the next loop
                            merged_styles = merge_styles(run_styles, parent_styles)

                            # Recursively process the nested tags
                            self._convert_inline(tag_contents, par, styles=merged_styles)
                    elif tag_name:
                        logger.warning(
                            "Ignoring a nested HTML tag not in the allowlist: %s",
                            tag_name,
                        )
                        continue
                    else:
                        content_text = tag

                    # Combine the parent's and run's styles, favoring the run's
                    merged_styles = merge_styles(run_styles, parent_styles)

                    # Write the text for this run
                    par = self._replace_and_write(content_text, par, merged_styles)

                    # Reset temporary run styles
                    run_styles = ReportConstants.DEFAULT_STYLE_VALUES.copy()

            # There are no tags to process, so write the string
            else:
                if isinstance(part, NavigableString):
                    par = self._replace_and_write(part, par, parent_styles)
                else:
                    par = self._replace_and_write(part.text, par)
        return par

    def _parse_html_lists(self, tag: Tag, prev_p, num: bool, level: int):
        """
        Recursively parse deeply nested lists. This checks for ``<ol>`` or ``<ul>`` tags
        and keeps parsing until all nested lists are found and processed.

        Returns the last paragraph object created.

        **Parameters**

        ``tag``
            BeautifulSoup 4 Tag object to parse
        ``prev_p``
            Previous paragraph to link the next list item
        ``num``
            Boolean to determine if the line item will be numbered
        ``level``
            Indentation level for the list item (Defaults to 0)
        """
        p = prev_p
        # Loop over the contents to find nested lists
        # Nested lists are ``ol`` or ``ul`` tags inside ``li`` tags
        for part in tag.children:
            # Handle ``li`` tags which might contain more lists
            if part.name == "li":
                li_contents = part.contents
                # A length of ``1`` means there are no nested tags
                if len(li_contents) == 1:
                    p = self._create_list_paragraph(prev_p, level, num)
                    if li_contents[0].name:
                        self._convert_inline(li_contents, p)
                    else:
                        self._replace_and_write(part.text, p)
                # Bigger lists mean more tags, so process nested tags
                else:
                    # Handle nested lists
                    p = self._parse_nested_html_lists(part, prev_p, num, level)
                # Track the paragraph used for this list item to link subsequent list paragraphs
                prev_p = p
            # If ``ol`` tag encountered, increment ``level`` and switch to numbered list
            elif part.name == "ol":
                level += 1
                p = self._parse_html_lists(part, prev_p, True, level)
            # If ``ul`` tag encountered, increment ``level`` and switch to bulleted list
            elif part.name == "ul":
                level += 1
                p = self._parse_html_lists(part, prev_p, False, level)
            # No change in list type, so proceed with writing the line
            elif part.name:
                p = self._create_list_paragraph(prev_p, level, num)
                self._convert_inline(part, p)
            # Handle tags that are not handled above
            else:
                if not isinstance(part, NavigableString):
                    logger.warning("Encountered an unknown tag for a list: %s", part.name)
                else:
                    if part.strip() != "":
                        p = self._create_list_paragraph(prev_p, level, num)
                        self._replace_and_write(part.strip(), p)
        # Return last paragraph created
        return p

    def _parse_nested_html_lists(self, tag: Tag, prev_p, num: bool, level: int = 0):
        """
        Recursively parse deeply nested lists. This checks for ``<ol>`` or ``<ul>`` tags
        and keeps parsing until all nested lists are found and processed.

        Returns the last paragraph object created.

        **Parameters**

        ``part``
            BeautifulSoup 4 Tag object to parse
        ``prev_p``
            Previous paragraph to link the next list item
        ``num``
            Boolean to determine if the line item will be numbered
        ``level``
            Indentation level for the list item (Defaults to 0)
        """
        # Check if this element has any ``ul`` or ``ol`` tags anywhere in its contents
        if tag.ol or tag.ul:
            temp = []
            nested_list = None
            li_contents = tag.contents
            # Loop over the contents of the list item to find the nested lists
            for part in tag:
                # Newlines will appear between elements, ignore them
                if part != "\n":
                    # If the tag name is ``ul`` or ``ol``, tack it as a nested list
                    if part.name in ("ol", "ul"):
                        num = bool(part.name == "ol")
                        nested_list = part
                    # Put everything else in a temporary list to be processed later
                    else:
                        temp.append(part)

            # Make the list paragraph here to pick up any changes to the list style (e.g, ``num``)
            p = self._create_list_paragraph(prev_p, level, num)

            # If ``temp`` isn't empty, process it like any other line
            if temp:
                # A length of ``1`` means no nested tags
                if len(temp) == 1:
                    # If the first list item is a ``Tag`` process for styling
                    if temp[0].name:
                        self._convert_inline(temp, p)
                    # Otherwise, just write the XML
                    else:
                        self._replace_and_write(temp[0], p)
                else:
                    self._convert_inline(temp, p)

            # If we have nested list(s), recursively process them by re-entering ``_parse_html_lists``
            if nested_list:
                # Increment the indentation level
                if not li_contents[0] == "\n":
                    level += 1
                p = self._parse_html_lists(nested_list, p, num, level)
        # No nested list items, proceed as normal
        # This is where we catch ``li`` tags with nested tags like hyperlinks
        else:
            p = self._create_list_paragraph(prev_p, level, num)
            self._convert_inline(tag.contents, p)
        return p

    def _create_list_paragraph(self, prev_p, level: int, num: bool = False, alignment=WD_ALIGN_PARAGRAPH.LEFT):
        """
        Create a new paragraph in the document for a list.

        **Parameters**

        ``prev_p``
            Previous paragraph to link the next list item
        ``level``
            Indentation level for the list item
        ``num``
            Boolean to determine if the line item will be numbered (Default: False)
        """
        if self.report_type == "pptx":
            # Move to new paragraph/line and indent bullets based on level
            p = self.finding_body_shape.text_frame.add_paragraph()
            p.level = level
        else:
            make_list = True
            styles = self.sacrificial_doc.styles
            try:
                if num:
                    list_style = styles["Number List"]
                else:
                    list_style = styles["Bullet List"]
            except Exception:
                if "List Paragraph" in styles:
                    list_style = styles["List Paragraph"]
                else:
                    list_style = styles["Normal"]
                    make_list = False
            p = self.sacrificial_doc.add_paragraph(style=list_style)
            if make_list:
                p = self._list_number(p, prev=prev_p, level=level, num=num)
            p.alignment = alignment
        return p

    def _list_number(self, par, prev=None, level: int | None = None, num: bool = True):
        """
        Makes the specified paragraph a list item with a specific level and optional restart.

        An attempt will be made to retrieve an abstract numbering style that corresponds
        to the style of the paragraph. If that is not possible, the default numbering or
        bullet style will be used based on the ``num`` parameter.

        Source:
            https://github.com/python-openxml/python-docx/issues/25#issuecomment-400787031

        **Parameters**

        ``par`` : docx.paragraph.Paragraph
            The docx paragraph to turn into a list item.
        ``prev`` : docx.paragraph.Paragraph or None
            The previous paragraph in the list. If specified, the numbering and styles will
            be taken as a continuation of this paragraph. If omitted, a new numbering scheme
            will be started.
        ``level`` : int or None
            The level of the paragraph within the outline. If ``prev`` is set, defaults
            to the same level as in ``prev``. Otherwise, defaults to zero.
        ``num`` : bool
            If ``prev`` is :py:obj:`None` and the style of the paragraph does not correspond
            to an existing numbering style, this will determine whether the list will
            be numbered or bulleted. The result is not guaranteed, but is fairly safe for
            most Word templates.
        """
        xpath_options = {
            True: {"single": "count(w:lvl)=1 and ", "level": 0},
            False: {"single": "", "level": level},
        }

        def style_xpath(prefer_single=True):
            """The style comes from the outer-scope variable ``par.style.name``."""
            style = par.style.style_id
            return (
                "w:abstractNum[" '{single}w:lvl[@w:ilvl="{level}"]/w:pStyle[@w:val="{style}"]' "]/@w:abstractNumId"
            ).format(style=style, **xpath_options[prefer_single])

        def type_xpath(prefer_single=True):
            """The type is from the outer-scope variable ``num``."""
            t = "decimal" if num else "bullet"
            return (
                "w:abstractNum[" '{single}w:lvl[@w:ilvl="{level}"]/w:numFmt[@w:val="{type}"]' "]/@w:abstractNumId"
            ).format(type=t, **xpath_options[prefer_single])

        def get_abstract_id():
            """
            Select as follows:

            1. Match single-level by style (get min ID)
            2. Match exact style and level (get min ID)
            3. Match single-level decimal/bullet types (get min ID)
            4. Match decimal/bullet in requested level (get min ID)
            """
            for fn in (style_xpath, type_xpath):
                for prefer_single in (True, False):
                    xpath = fn(prefer_single)
                    ids = numbering.xpath(xpath)
                    if ids:
                        return min(int(x) for x in ids)
            return 0

        if prev is None or prev._p.pPr is None or prev._p.pPr.numPr is None or prev._p.pPr.numPr.numId is None:
            if level is None:
                level = 0
            numbering = self.sacrificial_doc.part.numbering_part.numbering_definitions._numbering
            # Compute the abstract ID first by style, then by ``num``
            abstract = get_abstract_id()
            # Set the concrete numbering based on the abstract numbering ID
            num = numbering.add_num(abstract)
            # Make sure to override the abstract continuation property
            num.add_lvlOverride(ilvl=level).add_startOverride(1)
            # Extract the newly-allocated concrete numbering ID
            num = num.numId
        else:
            if level is None:
                level = prev._p.pPr.numPr.ilvl.val
            # Get the previous concrete numbering ID
            num = prev._p.pPr.numPr.numId.val
        par._p.get_or_add_pPr().get_or_add_numPr().get_or_add_numId().val = num
        par._p.get_or_add_pPr().get_or_add_numPr().get_or_add_ilvl().val = level

        return par

    def _replace_and_write(self, text, par, styles=ReportConstants.DEFAULT_STYLE_VALUES.copy()):
        """
        Write the provided text to Office XML.

        **Parameters**

        ``text`` : string
            Text to check for keywords
        ``par`` : Paragraph
            Paragraph for the processed text
        ``styles`` : dict
            Copy of ``ReportConstants.DEFAULT_STYLE_VALUES`` with styles for the text
        """

        # Remove any newlines to avoid creating unwanted blank lines
        text = text.replace("\r\n", "")

        # Handle hyperlinks based on Office report type
        # Easy with ``python-pptx`` API, but custom work required for ``python-docx``
        if styles["hyperlink"] and styles["hyperlink_url"]:
            if self.report_type == "pptx":
                run = par.add_run()
                run.text = text
                run.hyperlink.address = styles["hyperlink_url"]
            else:
                # For Word, this code is modified from this issue:
                #   https://github.com/python-openxml/python-docx/issues/384
                # Get an ID from the ``document.xml.rels`` file
                part = par.part
                r_id = part.relate_to(
                    styles["hyperlink_url"],
                    docx.opc.constants.RELATIONSHIP_TYPE.HYPERLINK,
                    is_external=True,
                )
                # Create the ``w:hyperlink`` tag and add needed values
                hyperlink = docx.oxml.shared.OxmlElement("w:hyperlink")
                hyperlink.set(
                    docx.oxml.shared.qn("r:id"),
                    r_id,
                )
                # Create the ``w:r`` and ``w:rPr`` elements
                new_run = docx.oxml.shared.OxmlElement("w:r")
                rPr = docx.oxml.shared.OxmlElement("w:rPr")
                new_run.append(rPr)
                new_run.text = text
                hyperlink.append(new_run)
                # Create a new Run object and add the hyperlink into it
                run = par.add_run()
                run._r.append(hyperlink)
                # A workaround for the lack of a hyperlink style
                if "Hyperlink" in self.sacrificial_doc.styles:
                    run.style = "Hyperlink"
                else:
                    run.font.color.theme_color = MSO_THEME_COLOR_INDEX.HYPERLINK
                    run.font.underline = True
        else:
            run = par.add_run()
            run.text = text

        # Apply font-based styles that work for both APIs
        font = run.font
        font.bold = styles["bold"]
        font.italic = styles["italic"]
        font.underline = styles["underline"]
        if styles["font_color"]:
            font.color.rgb = styles["font_color"]
        font.name = styles["font_family"]
        font.size = styles["font_size"]
        if styles["inline_code"]:
            if self.report_type == "pptx":
                font.name = "Courier New"
            else:
                run.style = "CodeInline"
                font.no_proof = True

        # These styles require extra work due to limitations of the ``python-pptx`` API
        if styles["highlight"]:
            if self.report_type == "pptx":
                rPr = run._r.get_or_add_rPr()
                highlight = OxmlElement("a:highlight")
                srgbClr = OxmlElement("a:srgbClr")
                srgbClr.set("val", "FFFF00")
                highlight.append(srgbClr)
                rPr.append(highlight)
            else:
                font.highlight_color = WD_COLOR_INDEX.YELLOW
        if styles["strikethrough"]:
            if self.report_type == "pptx":
                font._element.set("strike", "sngStrike")
            else:
                font.strike = styles["strikethrough"]
        if styles["subscript"]:
            if self.report_type == "pptx":
                font._element.set("baseline", "-25000")
            else:
                font.subscript = styles["subscript"]
        if styles["superscript"]:
            if self.report_type == "pptx":
                font._element.set("baseline", "30000")
            else:
                font.superscript = styles["superscript"]
        return par

    def _get_styles(self, tag):
        """
        Get styles from an BS4 ``Tag`` object's ``styles`` attribute and convert
        the string to a dictionary.

        **Parameters**

        ``tag`` : Tag
            BS4 ``Tag`` with a ``styles`` attribute
        """
        tag_styles = {}
        style_str = tag.attrs["style"]
        # Filter any blanks from the split
        style_list = list(filter(None, style_str.split(";")))
        for style in style_list:
            temp = style.split(":")
            key = temp[0].strip()
            value = temp[1].strip()
            try:
                if key == "font-size":
                    # Remove the "pt" from the font size and convert it to a ``Pt`` object
                    value = value.replace("pt", "")
                    value = float(value)
                    value = Pt(value)
                if key == "font-family":
                    # Some fonts include a list of values – get just the first one
                    font_list = value.split(",")
                    priority_font = font_list[0].replace("'", "").replace('"', "").strip()
                    value = priority_font
                if key == "color":
                    # Convert the color hex value to an ``RGBColor`` object
                    value = value.replace("#", "")
                    n = 2
                    hex_color = [hex(int(value[i : i + n], 16)) for i in range(0, len(value), n)]
                    if self.report_type == "pptx":
                        value = PptxRGBColor(*map(lambda v: int(v, 16), hex_color))
                    else:
                        value = RGBColor(*map(lambda v: int(v, 16), hex_color))
                tag_styles[key] = value
            except Exception:
                logger.exception("Failed to convert one of the inline styles for a text run")
        return tag_styles

    def _make_figure(self, par, ref=None):
        return self.rtcwc.rtc.make_figure(par, ref)

    def _process_evidence(
        self,
        evidence,
        par,
    ):
        return self.rtcwc.rtc.make_evidence(
            evidence,
            par,
            sacrificial_doc=self.sacrificial_doc,
            finding_slide=self.finding_slide,
        )

    def _make_cross_ref(self, par, ref):
        """
        Append a text run configured as a cross-reference to the provided paragraph.

        **Parameters**

        ``par`` : docx.paragraph.Paragraph
            Paragraph to alter
        ``ref`` : string
            The ``w:name`` value of the target bookmark
        """
        # Start the field character run for the label and number
        run = par.add_run()
        r = run._r
        fldChar = OxmlElement("w:fldChar")
        fldChar.set(qn("w:fldCharType"), "begin")
        r.append(fldChar)

        # Add field code instructions with ``instrText`` that points to the target bookmark
        run = par.add_run()
        r = run._r
        instrText = OxmlElement("w:instrText")
        instrText.text = f" REF _Ref{ref} \\h "
        r.append(instrText)

        # An optional ``separate`` value to enforce a space between label and number
        run = par.add_run()
        r = run._r
        fldChar = OxmlElement("w:fldChar")
        fldChar.set(qn("w:fldCharType"), "separate")
        r.append(fldChar)

        # Add runs for the figure label and number
        run = par.add_run(self.label_figure)
        # This ``#`` is a placeholder Word will replace with the figure's number
        run = par.add_run("#")

        # Close the  field character run
        run = par.add_run()
        r = run._r
        fldChar = OxmlElement("w:fldChar")
        fldChar.set(qn("w:fldCharType"), "end")
        r.append(fldChar)

        return par


def _valid_xml_char_ordinal(c):
    """
    Clean string to make all characters XML compatible for Word documents.

    Source:
        https://stackoverflow.com/questions/8733233/filtering-out-certain-bytes-in-python

    **Parameters**

    ``c`` : string
        String of characters to validate
    """
    codepoint = ord(c)
    # Conditions ordered by presumed frequency
    return (
        0x20 <= codepoint <= 0xD7FF
        or codepoint in (0x9, 0xA, 0xD)
        or 0xE000 <= codepoint <= 0xFFFD
        or 0x10000 <= codepoint <= 0x10FFFF
    )


def _jinja_evidence(evidence_name):
    """
    `{{evidence(name)}}` function in jinja.
    """
    return Markup("<span data-gw-evidence=\"" + html.escape(evidence_name) + "\"></span>")


def _jinja_caption(caption_name):
    return Markup("<span data-gw-caption=\"" + html.escape(caption_name) + "\"></span>")


def _jinja_ref(ref_name):
    return Markup("<span data-gw-ref=\"" + html.escape(ref_name) + "\"></span>")
