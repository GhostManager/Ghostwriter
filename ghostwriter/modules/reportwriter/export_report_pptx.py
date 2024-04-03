
import io
import logging
from bs4 import BeautifulSoup
from datetime import date

import pptx
from pptx import Presentation
from pptx.parts.presentation import PresentationPart
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.exc import PackageNotFoundError as PptxPackageNotFoundError
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
from pptx.util import Inches
from django.conf import settings
from django.utils.dateformat import format as dateformat

from ghostwriter.commandcenter.models import CompanyInformation
from ghostwriter.modules.reportwriter.export_report_base import ExportReportBase
from ghostwriter.modules.reportwriter.html_to_pptx import HtmlToPptxWithEvidence

logger = logging.getLogger(__name__)


class ExportReportPptx(ExportReportBase):
    ppt_presentation: PresentationPart

    def __init__(self, report, template_loc=None):
        super().__init__(report)

        self.company_config = CompanyInformation.get_solo()

        try:
            self.ppt_presentation = Presentation(template_loc)
        except ValueError:
            logger.exception(
                "Failed to load the provided template document because it is not a PowerPoint file: %s",
                template_loc,
            )
            raise
        except PptxPackageNotFoundError:
            logger.exception(
                "Failed to load the provided template document because file could not be found: %s",
                template_loc,
            )
            raise
        except Exception:
            logger.exception(
                "Failed to load the provided template document for unknown reason: %s",
                template_loc,
            )
            raise

    def _process_rich_text_pptx(self, text, slide, shape, template_vars, evidences):
        """
        Converts HTML from the TinyMCE rich text editor and inserts it into the passed in slide and shape
        """
        text = self.preprocess_rich_text(text, template_vars)
        try:
            HtmlToPptxWithEvidence.run(
                text,
                slide=slide,
                shape=shape,
                evidences=evidences,
            )
        except:
            # Log input text to help diagnose errors
            logger.warning("Input text: %r", text)
            raise

    def run(self) -> io.BytesIO:
        """Generate a complete PowerPoint slide deck for the current report."""

        base_context = self.jinja_richtext_base_context()
        base_evidences = {e["friendly_name"]: e for e in self.data["evidence"]}

        # Loop through the findings to create slides
        findings_stats = {}

        # Calculate finding stats
        for finding in self.data["findings"]:
            findings_stats[finding["severity"]] = 0

        for finding in self.data["findings"]:
            findings_stats[finding["severity"]] += 1

        # Slide styles (From Master Style counting top to bottom from 0..n)
        SLD_LAYOUT_TITLE = 0
        SLD_LAYOUT_TITLE_AND_CONTENT = 1
        SLD_LAYOUT_FINAL = 12

        # Add a title slide
        slide_layout = self.ppt_presentation.slide_layouts[SLD_LAYOUT_TITLE]
        slide = self.ppt_presentation.slides.add_slide(slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        title_shape.text = f'{self.data["client"]["name"]} {self.data["project"]["type"]}'
        text_frame = get_textframe(body_shape)
        # Use ``text_frame.text`` for first line/paragraph or ``text_frame.paragraphs[0]``
        text_frame.text = "Technical Outbrief"
        p = text_frame.add_paragraph()
        p.text = dateformat(date.today(), settings.DATE_FORMAT)

        # Add Agenda slide
        slide_layout = self.ppt_presentation.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
        slide = self.ppt_presentation.slides.add_slide(slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.text = "Agenda"
        body_shape = shapes.placeholders[1]
        text_frame = get_textframe(body_shape)
        text_frame.clear()
        delete_paragraph(text_frame.paragraphs[0])

        write_bullet(text_frame, "Introduction", 0)
        write_bullet(text_frame, "Assessment Details", 0)
        write_bullet(text_frame, "Methodology", 0)
        write_bullet(text_frame, "Assessment Timeline", 0)
        write_bullet(text_frame, "Attack Path Overview", 0)
        write_bullet(text_frame, "Positive Control Observations", 0)
        write_bullet(text_frame, "Findings and Recommendations Overview", 0)
        write_bullet(text_frame, "Next Steps", 0)

        # Add Introduction slide
        slide_layout = self.ppt_presentation.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
        slide = self.ppt_presentation.slides.add_slide(slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.text = "Introduction"
        body_shape = shapes.placeholders[1]
        text_frame = get_textframe(body_shape)
        text_frame.clear()

        if self.data["team"]:
            # Frame needs at least one paragraph to be valid, so don't delete the paragraph
            # if there are no team members
            delete_paragraph(text_frame.paragraphs[0])
            for member in self.data["team"]:
                write_bullet(text_frame, f"{member['name']} – {member['role']}", 0)
                write_bullet(text_frame, member["email"], 1)

        # Add Assessment Details slide
        slide_layout = self.ppt_presentation.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
        slide = self.ppt_presentation.slides.add_slide(slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.text = "Assessment Details"
        body_shape = shapes.placeholders[1]
        text_frame = get_textframe(body_shape)
        text_frame.clear()
        delete_paragraph(text_frame.paragraphs[0])

        write_bullet(
            text_frame, f"{self.data['project']['type']} assessment of {self.data['client']['name']}", 0
        )
        write_bullet(
            text_frame,
            f"Testing performed from {self.data['project']['start_date']} to {self.data['project']['end_date']}",
            1,
        )

        finding_body_shape = shapes.placeholders[1]
        self._process_rich_text_pptx(
            self.data["project"]["note"],
            slide=slide,
            shape=finding_body_shape,
            template_vars=base_context,
            evidences=base_evidences,
        )
        # The  method adds a new paragraph, so we need to get the last one to increase the indent level
        text_frame = get_textframe(finding_body_shape)
        p = text_frame.paragraphs[-1]
        p.level = 1

        if self.data["objectives"]:
            primary_objs = []
            secondary_objs = []
            tertiary_objs = []
            for objective in self.data["objectives"]:
                if objective["priority"] == "Primary":
                    primary_objs.append(objective)
                elif objective["priority"] == "Secondary":
                    secondary_objs.append(objective)
                elif objective["priority"] == "Tertiary":
                    tertiary_objs.append(objective)

            if primary_objs:
                write_bullet(text_frame, "Primary Objectives", 0)
                write_objective_list(text_frame, primary_objs)

            if secondary_objs:
                write_bullet(text_frame, "Secondary Objectives", 0)
                write_objective_list(text_frame, secondary_objs)

            if tertiary_objs:
                write_bullet(text_frame, "Tertiary Objectives", 0)
                write_objective_list(text_frame, tertiary_objs)

        # Add Methodology slide
        slide_layout = self.ppt_presentation.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
        slide = self.ppt_presentation.slides.add_slide(slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.text = "Methodology"

        # Add Timeline slide
        slide_layout = self.ppt_presentation.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
        slide = self.ppt_presentation.slides.add_slide(slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.text = "Assessment Timeline"

        # Delete the default text placeholder
        textbox = shapes[1]
        sp = textbox.element
        sp.getparent().remove(sp)
        # Add a table
        rows = 4
        columns = 2
        left = Inches(1.5)
        top = Inches(2)
        width = Inches(8)
        height = Inches(0.8)
        table = shapes.add_table(rows, columns, left, top, width, height).table
        # Set column width
        table.columns[0].width = Inches(2.0)
        table.columns[1].width = Inches(8.5)
        # Write table headers
        cell = table.cell(0, 0)
        cell.text = "Date"
        cell.fill.solid()
        cell.fill.fore_color.rgb = pptx.dml.color.RGBColor(0x2D, 0x28, 0x69)
        cell = table.cell(0, 1)
        cell.text = "Action Item"
        cell.fill.solid()
        cell.fill.fore_color.rgb = pptx.dml.color.RGBColor(0x2D, 0x28, 0x69)

        # Write date rows
        row_iter = 1
        table.cell(row_iter, 0).text = self.data["project"]["start_date"]
        table.cell(row_iter, 1).text = "Assessment execution began"
        row_iter += 1
        table.cell(row_iter, 0).text = self.data["project"]["end_date"]
        table.cell(row_iter, 1).text = "Assessment execution completed"
        row_iter += 1
        table.cell(row_iter, 0).text = self.data["project"]["end_date"]
        table.cell(row_iter, 1).text = "Draft report delivery"

        # Set all cells alignment to center and vertical center
        for cell in table.iter_cells():
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Add Attack Path Overview slide
        slide_layout = self.ppt_presentation.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
        slide = self.ppt_presentation.slides.add_slide(slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.text = "Attack Path Overview"

        # Add Observations slide
        slide_layout = self.ppt_presentation.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
        slide = self.ppt_presentation.slides.add_slide(slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        title_shape.text = "Positive Observations"
        text_frame = get_textframe(body_shape)

        # If there are observations then write a table
        if len(self.data["observations"]) > 0:
            # Delete the default text placeholder
            textbox = shapes[1]
            sp = textbox.element
            sp.getparent().remove(sp)
            # Add a table
            rows = len(self.data["observations"]) + 1
            columns = 1
            left = Inches(1.5)
            top = Inches(2)
            width = Inches(8)
            height = Inches(0.8)
            table = shapes.add_table(rows, columns, left, top, width, height).table
            # Set column width
            table.columns[0].width = Inches(10.5)
            # Write table headers
            cell = table.cell(0, 0)
            cell.text = "Observation"
            cell.fill.solid()
            cell.fill.fore_color.rgb = pptx.dml.color.RGBColor(0x2D, 0x28, 0x69)
            # Write findings rows
            row_iter = 1
            for observation in self.data["observations"]:
                table.cell(row_iter, 0).text = observation["title"]
                row_iter += 1
            # Set all cells alignment to center and vertical center
            for cell in table.iter_cells():
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        else:
            write_bullet(text_frame, "No observations", 0)

        # Create slide for each observation
        for observation in self.data["observations"]:
            slide_layout = self.ppt_presentation.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
            observation_slide = self.ppt_presentation.slides.add_slide(slide_layout)
            shapes = observation_slide.shapes
            title_shape = shapes.title

            # Prepare text frame
            observation_body_shape = shapes.placeholders[1]
            if observation_body_shape.has_text_frame:
                text_frame = get_textframe(observation_body_shape)
                text_frame.clear()
                delete_paragraph(text_frame.paragraphs[0])
            else:
                text_frame = None

            # Set slide title to title + [severity]
            title_shape.text = f'{observation["title"]}'

            # Add description to the slide body (other sections will appear in the notes)
            if observation.get("description", "").strip():
                observation_context = self.jinja_richtext_base_context()
                self._process_rich_text_pptx(
                    observation["description"],
                    slide=observation_slide,
                    shape=observation_body_shape,
                    template_vars=observation_context,
                    evidences=base_evidences,
                )
            else:
                par = observation_body_shape.add_paragraph()
                run = par.add_run()
                run.text = "No description provided"

            for ev in observation.get("evidence", []):
                HtmlToPptxWithEvidence.make_evidence(observation_slide, ev)

            # Ensure there is at least one paragraph, as required by the spec
            if text_frame is not None and not text_frame.paragraphs:
                text_frame.add_paragraph()

        # Add Findings Overview Slide
        slide_layout = self.ppt_presentation.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
        slide = self.ppt_presentation.slides.add_slide(slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        title_shape.text = "Findings Overview"
        text_frame = get_textframe(body_shape)

        # If there are findings then write a table of findings and severity ratings
        if len(self.data["findings"]) > 0:
            # Delete the default text placeholder
            textbox = shapes[1]
            sp = textbox.element
            sp.getparent().remove(sp)
            # Add a table
            rows = len(self.data["findings"]) + 1
            columns = 2
            left = Inches(1.5)
            top = Inches(2)
            width = Inches(8)
            height = Inches(0.8)
            table = shapes.add_table(rows, columns, left, top, width, height).table
            # Set column width
            table.columns[0].width = Inches(8.5)
            table.columns[1].width = Inches(2.0)
            # Write table headers
            cell = table.cell(0, 0)
            cell.text = "Finding"
            cell.fill.solid()
            cell.fill.fore_color.rgb = pptx.dml.color.RGBColor(0x2D, 0x28, 0x69)
            cell = table.cell(0, 1)
            cell.text = "Severity"
            cell.fill.solid()
            cell.fill.fore_color.rgb = pptx.dml.color.RGBColor(0x2D, 0x28, 0x69)
            # Write findings rows
            row_iter = 1
            for finding in self.data["findings"]:
                table.cell(row_iter, 0).text = finding["title"]
                risk_cell = table.cell(row_iter, 1)
                # Set risk rating
                risk_cell.text = finding["severity"]
                # Set cell color fill type to solid
                risk_cell.fill.solid()
                # Color the risk cell based on corresponding severity color
                cell_color = pptx.dml.color.RGBColor(*map(lambda v: int(v, 16), finding["severity_color_hex"]))
                risk_cell.fill.fore_color.rgb = cell_color
                row_iter += 1
            # Set all cells alignment to center and vertical center
            for cell in table.iter_cells():
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        else:
            write_bullet(text_frame, "No findings", 0)

        # Create slide for each finding
        for finding in self.data["findings"]:
            slide_layout = self.ppt_presentation.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
            finding_slide = self.ppt_presentation.slides.add_slide(slide_layout)
            shapes = finding_slide.shapes
            title_shape = shapes.title

            # Prepare text frame
            finding_body_shape = shapes.placeholders[1]
            if finding_body_shape.has_text_frame:
                text_frame = get_textframe(finding_body_shape)
                text_frame.clear()
                delete_paragraph(text_frame.paragraphs[0])
            else:
                text_frame = None

            # Set slide title to title + [severity]
            title_shape.text = f'{finding["title"]} [{finding["severity"]}]'

            # Add description to the slide body (other sections will appear in the notes)
            if finding.get("description", "").strip():
                finding_context = self.jinja_richtext_finding_context(base_context, finding)
                finding_evidences = base_evidences | {e["friendly_name"]: e for e in finding["evidence"]}
                self._process_rich_text_pptx(
                    finding["description"],
                    slide=finding_slide,
                    shape=finding_body_shape,
                    template_vars=finding_context,
                    evidences=finding_evidences,
                )
            else:
                par = finding_body_shape.add_paragraph()
                run = par.add_run()
                run.text = "No description provided"

            for ev in finding.get("evidence", []):
                HtmlToPptxWithEvidence.make_evidence(finding_slide, ev)

            # Ensure there is at least one paragraph, as required by the spec
            if text_frame is not None and not text_frame.paragraphs:
                text_frame.add_paragraph()

            # Add all finding data to the notes section for easier reference during edits
            entities = prepare_for_pptx(finding["affected_entities"])
            impact = prepare_for_pptx(finding["impact"])
            host_detection = prepare_for_pptx(finding["host_detection_techniques"])
            net_detection = prepare_for_pptx(finding["network_detection_techniques"])
            recommendation = prepare_for_pptx(finding["recommendation"])
            replication = prepare_for_pptx(finding["replication_steps"])
            references = prepare_for_pptx(finding["references"])
            notes_slide = finding_slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            p = text_frame.add_paragraph()
            p.text = f"""
                {finding["severity"].capitalize()}: finding["title"]

                AFFECTED ENTITIES
                {entities}

                IMPACT
                {impact}

                MITIGATION
                {recommendation}

                REPLICATION
                {replication}

                HOST DETECTION
                {host_detection}

                NETWORK DETECTION
                ,
                {net_detection}

                REFERENCES
                {references}
            """.replace(
                "                ", ""
            )

        # Add Recommendations slide
        slide_layout = self.ppt_presentation.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
        slide = self.ppt_presentation.slides.add_slide(slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.text = "Recommendations"

        # Add Next Steps slide
        slide_layout = self.ppt_presentation.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
        slide = self.ppt_presentation.slides.add_slide(slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.text = "Next Steps"

        # Add final slide
        slide_layout = self.ppt_presentation.slide_layouts[SLD_LAYOUT_FINAL]
        slide = self.ppt_presentation.slides.add_slide(slide_layout)
        shapes = slide.shapes
        body_shape = shapes.placeholders[1]
        text_frame = get_textframe(body_shape)
        text_frame.clear()
        p = text_frame.paragraphs[0]
        p.line_spacing = 0.7
        p.text = self.company_config.company_name
        p = text_frame.add_paragraph()
        p.text = self.company_config.company_twitter
        p.line_spacing = 0.7
        p = text_frame.add_paragraph()
        p.text = self.company_config.company_email
        p.line_spacing = 0.7

        # Add footer elements (if there is one) to all slides based on the footer placeholder in the template
        for idx, slide in enumerate(self.ppt_presentation.slides):
            date_placeholder_idx = -1
            footer_placeholder_idx = -1
            slide_number_placeholder_idx = -1
            slide_layout = slide.slide_layout

            for idx, place in enumerate(slide_layout.placeholders):
                if "Footer" in place.name:
                    footer_placeholder_idx = idx
                if "Slide Number" in place.name:
                    slide_number_placeholder_idx = idx
                if "Date" in place.name:
                    date_placeholder_idx = idx

            # Skip the title slide at index 0
            if idx > 0:
                if footer_placeholder_idx > 0:
                    footer_layout_placeholder, footer_placeholder = clone_placeholder(
                        slide, slide_layout, footer_placeholder_idx
                    )
                    footer_placeholder.text = footer_layout_placeholder.text
                if slide_number_placeholder_idx > 0:
                    _, slide_number_placeholder = clone_placeholder(
                        slide, slide_layout, slide_number_placeholder_idx
                    )
                    add_slide_number(slide_number_placeholder)
                if date_placeholder_idx > 0:
                    _, date_placeholder = clone_placeholder(
                        slide, slide_layout, date_placeholder_idx
                    )
                    date_placeholder.text = dateformat(date.today(), settings.DATE_FORMAT)

        # Finalize document and return it for an HTTP response
        out = io.BytesIO()
        self.ppt_presentation.save(out)
        return out


def add_slide_number(txtbox):
    """
    Add a slide number to the provided textbox. Ideally, the textbox should be the slide layout's slide
    number placeholder to match the template.

    Ref: https://stackoverflow.com/a/55816723
    """
    # Get a textbox's paragraph element
    par = txtbox.text_frame.paragraphs[0]._p

    # The slide number is actually a field, so we add a `fld` element to the paragraph
    # The number enclosed in the `a:t` element is the slide number and should auto-update on load/shuffle
    fld_xml = (
        '<a:fld %s id="{1F4E2DE4-8ADA-4D4E-9951-90A1D26586E7}" type="slidenum">\n'
        '  <a:rPr lang="en-US" smtClean="0"/>\n'
        "  <a:t>2</a:t>\n"
        "</a:fld>\n" % nsdecls("a")
    )
    fld = parse_xml(fld_xml)
    par.append(fld)


def clone_placeholder(slide, slide_layout, placeholder_idx):
    """
    Clone a placeholder from the slide master and return the layout and the new shape.
    """
    layout_placeholder = slide_layout.placeholders[placeholder_idx]
    slide.shapes.clone_placeholder(layout_placeholder)

    # The cloned placeholder is now the last shape in the slide
    return layout_placeholder, slide.shapes[-1]


def get_textframe(shape):
    """
    Get the shape's text frame and enable automatic resizing. The resize only
    triggers after opening the file in the PowerPoint application and making a change or saving.
    """
    text_frame = shape.text_frame
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    return text_frame


def write_bullet(text_frame, text, level):
    """Write a bullet to the provided text frame at the specified level."""
    p = text_frame.add_paragraph()
    p.text = text
    p.level = level


def write_objective_list(text_frame, objectives):
    """Write a list of objectives to the provided text frame."""
    for obj in objectives:
        status = obj["status"]
        if obj["complete"]:
            status = "Achieved"
        write_bullet(text_frame, f"{obj['objective']} – {status}", 1)


def prepare_for_pptx(value):
    """Strip HTML and clear 0x0D characters to prepare text for notes slides."""
    try:
        if value:
            return BeautifulSoup(value, "lxml").text.replace("\x0D", "")
        return "N/A"
    except Exception:
        logger.exception("Failed parsing this value for PPTX: %s", value)
        return ""


def delete_paragraph(par):
    """
    Delete the specified paragraph.

    **Parameter**

    ``par``
        Paragraph to delete from the document
    """
    p = par._p
    parent_element = p.getparent()
    if parent_element is not None:
        parent_element.remove(p)
    else:
        logger.warning("Could not delete paragraph in because it had no parent element")
