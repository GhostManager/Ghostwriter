
import io
import logging
from datetime import date

from django.conf import settings
from django.utils.dateformat import format as dateformat
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches
import pptx

from ghostwriter.modules.reportwriter.base.pptx import SLD_LAYOUT_TITLE, SLD_LAYOUT_TITLE_AND_CONTENT, ExportBasePptx, delete_paragraph, get_textframe, write_bullet, write_objective_list
from ghostwriter.modules.reportwriter.project.base import ExportProjectBase

logger = logging.getLogger(__name__)


class ProjectSlidesMixin:
    """
    Adds a function for generating Project-related slides - shared between the project and report exports
    """

    def get_subtitle_shapes(self, shapes):
        """
        Find subtitle placeholder shapes on a slide.

        **Parameters**
            ``shapes``
                The shapes collection from the slide

        **Returns**
            List of subtitle placeholder shapes, ordered by index
        """
        subtitles = []
        for shape in shapes.placeholders:
            try:
                if shape.placeholder_format.type == PP_PLACEHOLDER.SUBTITLE:
                    subtitles.append(shape)
            except AttributeError:
                # Some shapes may not have placeholder_format
                continue
        # Sort by placeholder index to ensure consistent ordering
        subtitles.sort(key=lambda s: s.placeholder_format.idx)
        return subtitles

    def get_placeholder_or_textbox(
        self,
        shapes,
        placeholder_idx,
        left=None,
        top=None,
        width=None,
        height=None,
    ):
        """
        Safely get a placeholder by index, or create a text box fallback if it doesn't exist.

        **Parameters**

        ``shapes``
            The shapes collection from the slide
        ``placeholder_idx``
            The index of the placeholder to retrieve
        ``left``
            Left position for fallback textbox (default: Inches(1))
        ``top``
            Top position for fallback textbox (default: Inches(1.5))
        ``width``
            Width for fallback textbox (default: Inches(8))
        ``height``
            Height for fallback textbox (default: Inches(5))

        **Returns**
            The placeholder shape or a newly created text box
        """
        try:
            return shapes.placeholders[placeholder_idx]
        except KeyError:
            logger.warning(
                "Placeholder %d not found on slide. Creating fallback text box. "
                "This may indicate a template compatibility issue.",
                placeholder_idx,
            )
            # Create a text box as fallback
            left = left if left is not None else Inches(1)
            top = top if top is not None else Inches(1.5)
            width = width if width is not None else Inches(8)
            height = height if height is not None else Inches(5)
            return shapes.add_textbox(left, top, width, height)

    def get_title_or_textbox(self, shapes, title_text):
        """
        Safely get a title shape, or use an existing shape as fallback.

        **Parameters**

        ``shapes``
            The shapes collection from the slide
        ``title_text``
            The text to set on the title

        **Returns**
            The title shape or an existing shape/textbox to use for the title
        """
        # Try to get the title shape using the standard property (works if the slide layout has a title placeholder)
        title_shape = shapes.title
        if title_shape is None:
            # Try to get the first placeholder (typically the title on most layouts)
            try:
                if len(shapes.placeholders) > 0:
                    title_shape = shapes.placeholders[0]
                    logger.warning(
                        "Title placeholder not found via `shapes.title`, using first placeholder. "
                        "This may indicate a template compatibility issue."
                    )
            except (KeyError, IndexError):
                pass

            # If still no title shape, try to use the first shape (not just placeholder) on the slide
            if title_shape is None and len(shapes) > 0:
                title_shape = shapes[0]
                logger.warning(
                    "No title or placeholder found, using first shape on slide. "
                    "This may indicate a template compatibility issue."
                )

            # Only create a new textbox if there are no shapes at all
            if title_shape is None:
                logger.warning(
                    "No shapes found on slide. Creating fallback text box. "
                    "This may indicate a template compatibility issue."
                )
                title_shape = shapes.add_textbox(
                    left=Inches(0.5),
                    top=Inches(0.5),
                    width=Inches(9),
                    height=Inches(1)
                )

        title_shape.text = title_text
        return title_shape

    def create_project_slides(self, base_context):
        # Add a title slide
        slide_layout = self.ppt_presentation.slide_layouts[SLD_LAYOUT_TITLE]
        slide = self.ppt_presentation.slides.add_slide(slide_layout)
        shapes = slide.shapes
        _ = self.get_title_or_textbox(shapes, f'{self.data["client"]["name"]} {self.data["project"]["type"]}')

        # Try to detect and use subtitle placeholders
        subtitle_shapes = self.get_subtitle_shapes(shapes)

        if len(subtitle_shapes) >= 1:
            # Use first subtitle for "Technical Outbrief"
            subtitle_shapes[0].text = "Technical Outbrief"

            if len(subtitle_shapes) >= 2:
                # Use second subtitle for the date
                subtitle_shapes[1].text = dateformat(date.today(), settings.DATE_FORMAT)
            else:
                # Only one subtitle - add date as a second paragraph
                text_frame = get_textframe(subtitle_shapes[0])
                p = text_frame.add_paragraph()
                p.text = dateformat(date.today(), settings.DATE_FORMAT)
        else:
            # No subtitle placeholders found - fall back to using placeholder[1]
            logger.info("No subtitle placeholders detected on title slide, using fallback approach")
            body_shape = self.get_placeholder_or_textbox(
                shapes,
                1,
                left=Inches(1),
                top=Inches(3),
                width=Inches(8),
                height=Inches(2),
            )
            text_frame = get_textframe(body_shape)
            text_frame.text = "Technical Outbrief"
            p = text_frame.add_paragraph()
            p.text = dateformat(date.today(), settings.DATE_FORMAT)

        # Add Agenda slide
        slide_layout = self.ppt_presentation.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
        slide = self.ppt_presentation.slides.add_slide(slide_layout)
        shapes = slide.shapes
        _ = self.get_title_or_textbox(shapes, "Agenda")
        body_shape = self.get_placeholder_or_textbox(shapes, 1)
        text_frame = get_textframe(body_shape)
        text_frame.clear()
        delete_paragraph(text_frame.paragraphs[0])

        write_bullet(text_frame, "Introduction", 0)
        write_bullet(text_frame, "Assessment Details", 0)
        write_bullet(text_frame, "Methodology", 0)
        write_bullet(text_frame, "Assessment Timeline", 0)
        write_bullet(text_frame, "Attack Path Overview", 0)
        write_bullet(text_frame, "Positive Control Observations", 0)
        write_bullet(text_frame, "Findings and Recommendations Overview", 0)
        write_bullet(text_frame, "Next Steps", 0)

        # Add Introduction slide
        slide_layout = self.ppt_presentation.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
        slide = self.ppt_presentation.slides.add_slide(slide_layout)
        shapes = slide.shapes
        _ = self.get_title_or_textbox(shapes, "Introduction")
        body_shape = self.get_placeholder_or_textbox(shapes, 1)
        text_frame = get_textframe(body_shape)
        text_frame.clear()

        if self.data["team"]:
            # Frame needs at least one paragraph to be valid, so don't delete the paragraph
            # if there are no team members
            delete_paragraph(text_frame.paragraphs[0])
            for member in self.data["team"]:
                write_bullet(text_frame, f"{member['name']} – {member['role']}", 0)
                write_bullet(text_frame, member["email"], 1)

        # Add Assessment Details slide
        slide_layout = self.ppt_presentation.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
        slide = self.ppt_presentation.slides.add_slide(slide_layout)
        shapes = slide.shapes
        _ = self.get_title_or_textbox(shapes, "Assessment Details")
        body_shape = self.get_placeholder_or_textbox(shapes, 1)
        text_frame = get_textframe(body_shape)
        text_frame.clear()
        delete_paragraph(text_frame.paragraphs[0])

        write_bullet(
            text_frame, f"{self.data['project']['type']} assessment of {self.data['client']['name']}", 0
        )
        write_bullet(
            text_frame,
            f"Testing performed from {self.data['project']['start_date']} to {self.data['project']['end_date']}",
            1,
        )

        finding_body_shape = body_shape
        self.render_rich_text_pptx(
            base_context["project"]["description_rt"],
            slide=slide,
            shape=finding_body_shape,
        )

        # The  method adds a new paragraph, so we need to get the last one to increase the indent level
        text_frame = get_textframe(finding_body_shape)
        p = text_frame.paragraphs[-1]
        p.level = 1

        if self.data["objectives"]:
            primary_objs = []
            secondary_objs = []
            tertiary_objs = []
            for objective in self.data["objectives"]:
                if objective["priority"] == "Primary":
                    primary_objs.append(objective)
                elif objective["priority"] == "Secondary":
                    secondary_objs.append(objective)
                elif objective["priority"] == "Tertiary":
                    tertiary_objs.append(objective)

            if primary_objs:
                write_bullet(text_frame, "Primary Objectives", 0)
                write_objective_list(text_frame, primary_objs)

            if secondary_objs:
                write_bullet(text_frame, "Secondary Objectives", 0)
                write_objective_list(text_frame, secondary_objs)

            if tertiary_objs:
                write_bullet(text_frame, "Tertiary Objectives", 0)
                write_objective_list(text_frame, tertiary_objs)

        # Add Methodology slide
        slide_layout = self.ppt_presentation.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
        slide = self.ppt_presentation.slides.add_slide(slide_layout)
        shapes = slide.shapes
        _ = self.get_title_or_textbox(shapes, "Methodology")

        # Add Timeline slide
        slide_layout = self.ppt_presentation.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
        slide = self.ppt_presentation.slides.add_slide(slide_layout)
        shapes = slide.shapes
        _ = self.get_title_or_textbox(shapes, "Assessment Timeline")
        body_shape = self.get_placeholder_or_textbox(shapes, 1)

        # Delete the default text placeholder
        sp = body_shape.element
        sp.getparent().remove(sp)
        # Add a table
        rows = 4
        columns = 2
        left = Inches(1.5)
        top = Inches(2)
        width = Inches(8)
        height = Inches(0.8)
        table = shapes.add_table(rows, columns, left, top, width, height).table
        # Set column width
        table.columns[0].width = Inches(2.0)
        table.columns[1].width = Inches(8.5)
        # Write table headers
        cell = table.cell(0, 0)
        cell.text = "Date"
        cell.fill.solid()
        cell.fill.fore_color.rgb = pptx.dml.color.RGBColor(0x2D, 0x28, 0x69)
        cell = table.cell(0, 1)
        cell.text = "Action Item"
        cell.fill.solid()
        cell.fill.fore_color.rgb = pptx.dml.color.RGBColor(0x2D, 0x28, 0x69)

        # Write date rows
        row_iter = 1
        table.cell(row_iter, 0).text = self.data["project"]["start_date"]
        table.cell(row_iter, 1).text = "Assessment execution began"
        row_iter += 1
        table.cell(row_iter, 0).text = self.data["project"]["end_date"]
        table.cell(row_iter, 1).text = "Assessment execution completed"
        row_iter += 1
        table.cell(row_iter, 0).text = self.data["project"]["end_date"]
        table.cell(row_iter, 1).text = "Draft report delivery"

        # Set all cells alignment to center and vertical center
        for cell in table.iter_cells():
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Add Attack Path Overview slide
        slide_layout = self.ppt_presentation.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
        slide = self.ppt_presentation.slides.add_slide(slide_layout)
        shapes = slide.shapes
        _ = self.get_title_or_textbox(shapes, "Attack Path Overview")


class ExportProjectPptx(ExportBasePptx, ExportProjectBase, ProjectSlidesMixin):
    def run(self) -> io.BytesIO:
        base_context = self.map_rich_texts()
        self.create_project_slides(base_context)
        self.process_footers()
        return super().run()
