"""
This module contains the tools required for generating Microsoft Office documents for
reporting. The ``Reportwriter`` class accepts data and produces a docx, xlsx, pptx,
and json using the provided data.
"""

# Standard Libraries
import copy
import html
import io
import json
import logging
import os
import re
from copy import deepcopy
from datetime import date

# Django Imports
from django.conf import settings
from django.utils.dateformat import format as dateformat

# 3rd Party Libraries
import docx
import jinja2
import jinja2.sandbox
import pptx
from bs4 import BeautifulSoup
from docx.enum.style import WD_STYLE_TYPE
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.opc.exceptions import PackageNotFoundError as DocxPackageNotFoundError
from docx.shared import Inches, Pt
from docxtpl import DocxTemplate, RichText
from jinja2.exceptions import TemplateRuntimeError, TemplateSyntaxError, UndefinedError
from markupsafe import Markup
from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.exc import PackageNotFoundError as PptxPackageNotFoundError
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
from rest_framework.renderers import JSONRenderer
from xlsxwriter.workbook import Workbook

# Ghostwriter Libraries
from ghostwriter.commandcenter.models import (
    CompanyInformation,
    ExtraFieldSpec,
    ReportConfiguration,
)
from ghostwriter.modules.custom_serializers import ReportDataSerializer
from ghostwriter.modules.exceptions import InvalidFilterValue
from ghostwriter.modules.linting_utils import LINTER_CONTEXT
from ghostwriter.modules.reportwriter.extensions import IMAGE_EXTENSIONS, TEXT_EXTENSIONS
from ghostwriter.modules.reportwriter.html_to_docx import HtmlToDocxWithEvidence
from ghostwriter.modules.reportwriter.html_to_plain_text import html_to_plain_text
from ghostwriter.modules.reportwriter.html_to_pptx import HtmlToPptxWithEvidence
from ghostwriter.oplog.models import OplogEntry
from ghostwriter.reporting.models import Evidence, Finding, Observation, Report
from ghostwriter.rolodex.models import Client, Project
from ghostwriter.shepherd.models import Domain, StaticServer
from ghostwriter.modules.reportwriter import filters

# Using __name__ resolves to ghostwriter.modules.reporting
logger = logging.getLogger(__name__)


def prepare_jinja2_env(debug=False):
    """Prepare a Jinja2 environment with all custom filters."""
    if debug:
        undefined = jinja2.DebugUndefined
    else:
        undefined = jinja2.make_logging_undefined(logger=logger, base=jinja2.Undefined)

    env = jinja2.sandbox.SandboxedEnvironment(undefined=undefined, extensions=["jinja2.ext.debug"], autoescape=True)
    env.filters["filter_severity"] = filters.filter_severity
    env.filters["filter_type"] = filters.filter_type
    env.filters["strip_html"] = filters.strip_html
    env.filters["compromised"] = filters.compromised
    env.filters["add_days"] = filters.add_days
    env.filters["format_datetime"] = filters.format_datetime
    env.filters["get_item"] = filters.get_item
    env.filters["regex_search"] = filters.regex_search
    env.filters["filter_tags"] = filters.filter_tags

    return env


class ReportConstants:
    """Constant values used for report generation."""

    DEFAULT_STYLE_VALUES = {
        "bold": False,
        "underline": False,
        "italic": False,
        "inline_code": False,
        "strikethrough": False,
        "font_family": None,
        "font_size": None,
        "font_color": None,
        "highlight": None,
        "superscript": False,
        "subscript": False,
        "hyperlink": False,
        "hyperlink_url": None,
    }


class Reportwriter:
    """Generate report documents in Microsoft Office formats and JSON."""

    # Allowlist for HTML tags expected to come from the WYSIWYG
    tag_allowlist = [
        "code",
        "span",
        "p",
        "ul",
        "ol",
        "li",
        "a",
        "em",
        "strong",
        "u",
        "b",
        "pre",
        "sub",
        "sup",
        "del",
        "blockquote",
    ]

    def __init__(self, report_queryset, template_loc=None):
        self.template_loc = template_loc
        self.report_queryset = report_queryset

        # Generate the report JSON
        self.report_json = json.loads(self.generate_json())

        # Get the global report configuration
        global_report_config = ReportConfiguration.get_solo()
        self.company_config = CompanyInformation.get_solo()

        # Track report type for different Office XML
        self.report_type = None

        # Picture border settings for Word
        self.enable_borders = global_report_config.enable_borders
        self.border_color = global_report_config.border_color
        self.border_weight = global_report_config.border_weight

        # Caption options
        prefix_figure = global_report_config.prefix_figure.strip()
        self.prefix_figure = f" {prefix_figure} "
        label_figure = global_report_config.label_figure.strip()
        self.label_figure = f"{label_figure} "
        self.title_case_captions = global_report_config.title_case_captions
        self.title_case_exceptions = global_report_config.title_case_exceptions.split(",")

        # Set up Jinja2 rendering environment + custom filters
        self.jinja_env = prepare_jinja2_env(debug=False)

        self.extra_fields_spec_cache = {}

        logger.info(
            "Generating a report for %s using the template at %s",
            self.report_queryset,
            self.template_loc,
        )

    def _valid_xml_char_ordinal(self, c):
        """
        Clean string to make all characters XML compatible for Word documents.

        Source:
            https://stackoverflow.com/questions/8733233/filtering-out-certain-bytes-in-python

        **Parameters**

        ``c`` : string
            String of characters to validate
        """
        codepoint = ord(c)
        # Conditions ordered by presumed frequency
        return (
            0x20 <= codepoint <= 0xD7FF
            or codepoint in (0x9, 0xA, 0xD)
            or 0xE000 <= codepoint <= 0xFFFD
            or 0x10000 <= codepoint <= 0x10FFFF
        )

    def generate_json(self):
        """
        Export a report as a JSON dictionary for archiving and to generate other report types.
        """

        # Serialize the :model:`rolodex.Project`
        serializer = ReportDataSerializer(
            self.report_queryset,
            exclude=[
                "id",
            ],
        )
        # Render the serialized data as JSON
        report_json = JSONRenderer().render(serializer.data)
        report_json = json.loads(report_json)
        # An extra step to make the JSON "pretty"
        output = json.dumps(report_json, indent=4)

        return output

    def _preprocess_rich_text(self, text, template_vars):
        """
        Does jinja and `{{.item}}` substitutions on rich text, in preparation for feeding into the
        `BaseHtmlToOOXML` subclass.
        """

        if not text:
            return ""

        # Replace old `{{.item}}`` syntax with jinja templates or elements to replace
        def replace_old_tag(match):
            contents = match.group(1).strip()
            # These will be swapped out when parsing the HTML
            if contents.startswith("ref "):
                return _jinja_ref(contents[4:].strip())
            elif contents == "caption":
                return _jinja_caption("")
            elif contents.startswith("caption "):
                return _jinja_caption(contents[8:].strip())
            else:
                return "{{ _old_dot_vars[" + repr(contents.strip()) + "]}}"

        text_old_dot_subbed = re.sub(r"\{\{\.(.*?)\}\}", replace_old_tag, text)

        text_pagebrea_subbed = text_old_dot_subbed.replace(
            "<p><!-- pagebreak --></p>", '<br data-gw-pagebreak="true" />'
        )

        # Run template
        template = self.jinja_env.from_string(text_pagebrea_subbed)
        text_rendered = template.render(template_vars)

        # Filter out XML-incompatible characters
        text_char_filtered = "".join(c for c in text_rendered if self._valid_xml_char_ordinal(c))
        return text_char_filtered

    def _process_rich_text_docx(self, text, template_vars, evidences, p_style=None):
        """
        Converts HTML from the TinyMCE rich text editor to a Word subdoc.
        """
        text = self._preprocess_rich_text(text, template_vars)
        doc = self.word_doc.new_subdoc()
        try:
            HtmlToDocxWithEvidence.run(
                text,
                doc=doc,
                p_style=p_style,
                evidences=evidences,
                figure_label=self.label_figure,
                figure_prefix=self.prefix_figure,
                title_case_captions=self.title_case_captions,
                title_case_exceptions=self.title_case_exceptions,
                border_color_width=(self.border_color, self.border_weight) if self.enable_borders else None,
            )
        except:
            # Log input text to help diagnose errors
            logger.warning("Input text: %r", text)
            raise
        return doc

    def _process_rich_text_pptx(self, text, slide, shape, template_vars, evidences):
        """
        Converts HTML from the TinyMCE rich text editor and inserts it into the passed in slide and shape
        """
        text = self._preprocess_rich_text(text, template_vars)
        try:
            HtmlToPptxWithEvidence.run(
                text,
                slide=slide,
                shape=shape,
                evidences=evidences,
            )
        except:
            # Log input text to help diagnose errors
            logger.warning("Input text: %r", text)
            raise

    def _process_rich_text_xlsx(self, html, template_vars, evidences) -> str:
        """
        Converts HTML from the TinyMCE rich text editor and returns a plain string
        """
        text = self._preprocess_rich_text(html, template_vars)
        return html_to_plain_text(text, evidences)

    def generate_word_docx(self):
        """Generate a complete Word document for the current report."""

        # Create Word document writer using the specified template file
        try:
            self.word_doc = DocxTemplate(self.template_loc)
        except DocxPackageNotFoundError:
            logger.exception(
                "Failed to load the provided template document because file could not be found: %s",
                self.template_loc,
            )
            raise DocxPackageNotFoundError from docx.opc.exceptions.PackageNotFoundError
        except Exception:
            logger.exception("Failed to load the provided template document: %s", self.template_loc)

        # Check for styles
        styles = self.word_doc.styles
        if "CodeBlock" not in styles:
            codeblock_style = styles.add_style("CodeBlock", WD_STYLE_TYPE.PARAGRAPH)
            codeblock_style.base_style = styles["Normal"]
            codeblock_style.hidden = False
            codeblock_style.quick_style = True
            codeblock_style.priority = 2
            # Set font and size
            codeblock_font = codeblock_style.font
            codeblock_font.name = "Courier New"
            codeblock_font.size = Pt(11)
            # Set alignment
            codeblock_par = codeblock_style.paragraph_format
            codeblock_par.alignment = WD_ALIGN_PARAGRAPH.LEFT
            codeblock_par.line_spacing = 1
            codeblock_par.left_indent = Inches(0.2)
            codeblock_par.right_indent = Inches(0.2)

        if "CodeInline" not in styles:
            codeinline_style = styles.add_style("CodeInline", WD_STYLE_TYPE.CHARACTER)
            codeinline_style.hidden = False
            codeinline_style.quick_style = True
            codeinline_style.priority = 3
            # Set font and size
            codeinline_font = codeinline_style.font
            codeinline_font.name = "Courier New"
            codeinline_font.size = Pt(11)

        if "Caption" not in styles:
            caption_style = styles.add_style("Caption", WD_STYLE_TYPE.PARAGRAPH)
            caption_style.hidden = False
            caption_style.quick_style = True
            caption_style.priority = 4
            # Set font and size
            caption_font = caption_style.font
            caption_font.name = "Calibri"
            caption_font.size = Pt(9)
            caption_font.italic = True

        if "Blockquote" not in styles:
            block_style = styles.add_style("Blockquote", WD_STYLE_TYPE.PARAGRAPH)
            block_style.base_style = styles["Normal"]
            block_style.hidden = False
            block_style.quick_style = True
            block_style.priority = 5
            # Set font and size
            block_font = block_style.font
            block_font.name = "Calibri"
            block_font.size = Pt(12)
            block_font.italic = True
            # Set alignment
            block_par = block_style.paragraph_format
            block_par.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
            block_par.left_indent = Inches(0.2)
            block_par.right_indent = Inches(0.2)
            # Keep first and last lines together after repagination
            block_par.widow_control = True

        # Process template context, converting HTML elements to XML as needed
        context = deepcopy(self.report_json)
        context = self._process_richtext(context)

        # Render the Word document + auto-escape any unsafe XML/HTML
        self.word_doc.render(context, self.jinja_env, autoescape=True)

        # Return the final rendered document
        return self.word_doc

    def _jinja_richtext_base_context(self) -> dict:
        """
        Generates a Jinja context for use in rich text fields
        """
        base_context = {
            # `{{.foo}}` converts to `{{obsolete.foo}}`
            "_old_dot_vars": {
                "client": self.report_json["client"]["short_name"] or self.report_json["client"]["name"],
                "project_start": self.report_json["project"]["start_date"],
                "project_end": self.report_json["project"]["end_date"],
                "project_type": self.report_json["project"]["type"].lower(),
            },
            "mk_evidence": _jinja_evidence,
            "mk_caption": _jinja_caption,
            "mk_ref": _jinja_ref,
        }
        base_context.update(self.report_json)
        for evidence in self.report_json["evidence"]:
            if evidence.get("friendly_name"):
                base_context["_old_dot_vars"][evidence["friendly_name"]] = _jinja_evidence(evidence["friendly_name"])
        return base_context

    def _jinja_richtext_finding_context(self, base_context: dict, finding) -> dict:
        """
        Generates a Jinja context for use in finding-related rich text fields
        """
        finding_context = base_context.copy()
        finding_context.update(
            {
                "finding": finding,
                "_old_dot_vars": base_context["_old_dot_vars"].copy(),
            }
        )
        for evidence in finding["evidence"]:
            if evidence.get("friendly_name"):
                finding_context["_old_dot_vars"][evidence["friendly_name"]] = _jinja_evidence(evidence["friendly_name"])
        return finding_context

    def _process_richtext(self, context: dict) -> dict:
        """
        Update the document context with ``RichText`` and ``Subdocument`` objects for
        each finding and any other values editable with a WYSIWYG editor.

        **Parameters**

        ``context``
            Pre-defined template context
        """

        p_style = self.report_queryset.docx_template.p_style

        base_context = self._jinja_richtext_base_context()
        base_evidences = {e["friendly_name"]: e for e in context["evidence"]}

        def base_render(text):
            return self._process_rich_text_docx(text, base_context, base_evidences, p_style)

        self._process_extra_fields(context["extra_fields"], Report, base_render)

        # Findings
        for finding in context["findings"]:
            logger.info("Processing %s", finding["title"])

            finding_context = self._jinja_richtext_finding_context(base_context, finding)
            finding_evidences = base_evidences | {e["friendly_name"]: e for e in finding["evidence"]}

            def finding_render(text):
                return self._process_rich_text_docx(text, finding_context, finding_evidences, p_style)

            self._process_extra_fields(finding["extra_fields"], Finding, finding_render)

            # Create ``RichText()`` object for a colored severity category
            finding["severity_rt"] = RichText(finding["severity"], color=finding["severity_color"])
            finding["cvss_score_rt"] = RichText(finding["cvss_score"], color=finding["severity_color"])
            finding["cvss_vector_rt"] = RichText(finding["cvss_vector"], color=finding["severity_color"])
            # Create subdocuments for each finding section
            finding["affected_entities_rt"] = finding_render(finding["affected_entities"])
            finding["description_rt"] = finding_render(finding["description"])
            finding["impact_rt"] = finding_render(finding["impact"])

            # Include a copy of ``mitigation`` as ``recommendation`` to match legacy context
            mitigation_section = finding_render(finding["mitigation"])
            finding["mitigation_rt"] = mitigation_section
            finding["recommendation_rt"] = mitigation_section

            finding["replication_steps_rt"] = finding_render(finding["replication_steps"])
            finding["host_detection_techniques_rt"] = finding_render(finding["host_detection_techniques"])
            finding["network_detection_techniques_rt"] = finding_render(finding["network_detection_techniques"])
            finding["references_rt"] = finding_render(finding["references"])

        # Client
        context["client"]["note_rt"] = base_render(context["client"]["note"])
        context["client"]["address_rt"] = base_render(context["client"]["address"])
        self._process_extra_fields(context["client"]["extra_fields"], Client, base_render)

        # Project
        context["project"]["note_rt"] = base_render(context["project"]["note"])
        self._process_extra_fields(context["project"]["extra_fields"], Project, base_render)

        # Assignments
        for assignment in context["team"]:
            if isinstance(assignment, dict):
                if assignment["note"]:
                    assignment["note_rt"] = base_render(assignment["note"])

        # Contacts
        for contact in context["client"]["contacts"]:
            if isinstance(contact, dict):
                if contact["note"]:
                    contact["note_rt"] = base_render(contact["note"])

        # Objectives
        for objective in context["objectives"]:
            if isinstance(objective, dict):
                if objective["description"]:
                    objective["description_rt"] = base_render(objective["description"])

        # Scope Lists
        for scope_list in context["scope"]:
            if isinstance(scope_list, dict):
                if scope_list["description"]:
                    scope_list["description_rt"] = base_render(scope_list["description"])

        # Targets
        for target in context["targets"]:
            if isinstance(target, dict):
                if target["note"]:
                    target["note_rt"] = base_render(target["note"])

        # Deconfliction Events
        for event in context["deconflictions"]:
            if isinstance(event, dict):
                if event["description"]:
                    event["description_rt"] = base_render(event["description"])

        # White Cards
        for card in context["whitecards"]:
            if isinstance(card, dict):
                if card["description"]:
                    card["description_rt"] = base_render(card["description"])

        # Infrastructure
        for asset_type in context["infrastructure"]:
            for asset in context["infrastructure"][asset_type]:
                if isinstance(asset, dict):
                    if asset["note"]:
                        asset["note_rt"] = base_render(asset["note"])
        for asset in context["infrastructure"]["domains"]:
            self._process_extra_fields(asset["extra_fields"], Domain, base_render)
        for asset in context["infrastructure"]["servers"]:
            self._process_extra_fields(asset["extra_fields"], StaticServer, base_render)

        # Logs
        for log in context["logs"]:
            for entry in log["entries"]:
                self._process_extra_fields(entry["extra_fields"], OplogEntry, base_render)

        # Observations
        for observation in context["observations"]:
            if observation["description"]:
                observation["description_rt"] = base_render(observation["description"])
            self._process_extra_fields(observation["extra_fields"], Observation, base_render)

        # Report Evidence
        # for evidence in context["evidence"]:
        #    self._process_extra_fields(evidence["extra_fields"], Report, base_render)

        return context

    def _extra_field_specs_for(self, model):
        """
        Gets (and caches) the set of extra fields for a model class.
        """
        label = model._meta.label
        if label in self.extra_fields_spec_cache:
            return self.extra_fields_spec_cache[label]
        else:
            specs = ExtraFieldSpec.objects.filter(target_model=label)
            self.extra_fields_spec_cache[label] = specs
            return specs

    def _process_extra_fields(self, extra_fields, model, render_rich_text):
        specs = self._extra_field_specs_for(model)
        for field in specs:
            if field.internal_name not in extra_fields:
                extra_fields[field.internal_name] = field.default_value()
            if field.type == "rich_text":
                extra_fields[field.internal_name] = render_rich_text(str(extra_fields[field.internal_name]))

    def generate_excel_xlsx(self, memory_object):
        """
        Generate a complete Excel spreadsheet for the current report.

        **Parameters**

        ``memory_object``
            In-memory file-like object to write the Excel spreadsheet to
        """

        # Create an in-memory Excel workbook with a named worksheet
        xlsx_doc = Workbook(memory_object, {"in_memory": True, "strings_to_formulas": False, "strings_to_urls": False})
        worksheet = xlsx_doc.add_worksheet("Findings")

        # Create a format for headers
        bold_format = xlsx_doc.add_format({"bold": True})
        bold_format.set_text_wrap()
        bold_format.set_align("vcenter")
        bold_format.set_align("center")

        # Create a format for affected entities
        asset_format = xlsx_doc.add_format()
        asset_format.set_text_wrap()
        asset_format.set_align("vcenter")
        asset_format.set_align("center")

        # Formatting for severity cells
        severity_format = xlsx_doc.add_format({"bold": True})
        severity_format.set_align("vcenter")
        severity_format.set_align("center")
        severity_format.set_font_color("black")

        # Create a format for everything else
        wrap_format = xlsx_doc.add_format()
        wrap_format.set_text_wrap()
        wrap_format.set_align("vcenter")

        # Create header row for findings
        col = 0
        headers = [
            "Finding",
            "Severity",
            "CVSS Score",
            "CVSS Vector",
            "Affected Entities",
            "Description",
            "Impact",
            "Recommendation",
            "Replication Steps",
            "Host Detection Techniques",
            "Network Detection Techniques",
            "References",
            "Supporting Evidence",
            "Tags",
        ]

        findings_extra_field_specs = self._extra_field_specs_for(Finding)
        base_evidences = {e["friendly_name"]: e for e in self.report_json["evidence"]}

        # Create 30 width columns and then shrink severity to 10
        for header in headers:
            worksheet.write_string(0, col, header, bold_format)
            col += 1
        for field in findings_extra_field_specs:
            worksheet.write_string(0, col, field.display_name, bold_format)
            col += 1
        worksheet.set_column(0, 13, 30)
        worksheet.set_column(1, 1, 10)
        worksheet.set_column(2, 2, 10)
        worksheet.set_column(3, 3, 40)

        # Loop through the findings to create the rest of the worksheet
        col = 0
        row = 1
        base_context = self._jinja_richtext_base_context()
        for finding in self.report_json["findings"]:
            finding_context = self._jinja_richtext_finding_context(base_context, finding)
            finding_evidences = base_evidences | {e["friendly_name"]: e for e in finding["evidence"]}

            # Finding Name
            worksheet.write_string(
                row,
                col,
                self._process_rich_text_xlsx(finding["title"], finding_context, finding_evidences),
                bold_format,
            )
            col += 1

            # Update severity format bg color with the finding's severity color
            severity_format.set_bg_color(finding["severity_color"])

            # Severity and CVSS information
            worksheet.write_string(
                row,
                col,
                self._process_rich_text_xlsx(finding["severity"], finding_context, finding_evidences),
                severity_format,
            )
            col += 1
            if isinstance(finding["cvss_score"], float):
                worksheet.write_number(row, col, finding["cvss_score"], severity_format)
            else:
                worksheet.write_string(
                    row,
                    col,
                    self._process_rich_text_xlsx(finding["cvss_score"], finding_context, finding_evidences),
                    severity_format,
                )
            col += 1
            worksheet.write_string(
                row,
                col,
                self._process_rich_text_xlsx(finding["cvss_vector"], finding_context, finding_evidences),
                severity_format,
            )
            col += 1

            # Affected Entities
            if finding["affected_entities"]:
                worksheet.write_string(
                    row,
                    col,
                    self._process_rich_text_xlsx(finding["affected_entities"], finding_context, finding_evidences),
                    asset_format,
                )
            else:
                worksheet.write_string(row, col, "N/A", asset_format)
            col += 1

            # Description
            worksheet.write_string(
                row,
                col,
                self._process_rich_text_xlsx(finding["description"], finding_context, finding_evidences),
                wrap_format,
            )
            col += 1

            # Impact
            worksheet.write_string(
                row,
                col,
                self._process_rich_text_xlsx(finding["impact"], finding_context, finding_evidences),
                wrap_format,
            )
            col += 1

            # Recommendation
            worksheet.write_string(
                row,
                col,
                self._process_rich_text_xlsx(finding["recommendation"], finding_context, finding_evidences),
                wrap_format,
            )
            col += 1

            # Replication
            worksheet.write_string(
                row,
                col,
                self._process_rich_text_xlsx(finding["replication_steps"], finding_context, finding_evidences),
                wrap_format,
            )
            col += 1

            # Detection
            worksheet.write_string(
                row,
                col,
                self._process_rich_text_xlsx(finding["host_detection_techniques"], finding_context, finding_evidences),
                wrap_format,
            )
            col += 1
            worksheet.write_string(
                row,
                col,
                self._process_rich_text_xlsx(
                    finding["network_detection_techniques"], finding_context, finding_evidences
                ),
                wrap_format,
            )
            col += 1

            # References
            worksheet.write_string(
                row,
                col,
                self._process_rich_text_xlsx(finding["references"], finding_context, finding_evidences),
                wrap_format,
            )
            col += 1

            # Collect the evidence, if any, from the finding's folder and insert inline with description
            try:
                evidence_queryset = Evidence.objects.filter(finding=finding["id"])
            except Evidence.DoesNotExist:
                evidence_queryset = []
            except Exception:
                logger.exception("Query for evidence failed for finding %s", finding["id"])
                evidence_queryset = []
            evidence = [f.filename for f in evidence_queryset if f in TEXT_EXTENSIONS or f in IMAGE_EXTENSIONS]
            finding_evidence_names = "\r\n".join(map(str, evidence))
            worksheet.write_string(
                row,
                col,
                self._process_rich_text_xlsx(finding_evidence_names, finding_context, finding_evidences),
                wrap_format,
            )
            col += 1

            # Tags
            worksheet.write_string(
                row,
                col,
                self._process_rich_text_xlsx(", ".join(finding["tags"]), finding_context, finding_evidences),
                wrap_format,
            )
            col += 1

            # Extra fields
            for field_spec in findings_extra_field_specs:
                field_value = field_spec.value_of(finding["extra_fields"])
                if field_spec.type == "rich_text":
                    field_value = self._process_rich_text_xlsx(field_value, finding_context, finding_evidences)
                else:
                    field_value = str(field_value)
                worksheet.write_string(row, col, field_value, wrap_format)
                col += 1

            # Increment row counter and reset columns before moving on to next finding
            row += 1
            col = 0

        # Add a filter to the worksheet
        worksheet.autofilter("A1:M{}".format(len(self.report_json["findings"]) + 1))

        # Finalize document
        xlsx_doc.close()
        return xlsx_doc

    def generate_powerpoint_pptx(self):
        """Generate a complete PowerPoint slide deck for the current report."""

        self.report_type = "pptx"

        base_context = self._jinja_richtext_base_context()
        base_evidences = {e["friendly_name"]: e for e in self.report_json["evidence"]}

        # Create document writer using the specified template
        try:
            self.ppt_presentation = Presentation(self.template_loc)
        except ValueError:
            logger.exception(
                "Failed to load the provided template document because it is not a PowerPoint file: %s",
                self.template_loc,
            )
            raise ValueError
        except PptxPackageNotFoundError:
            logger.exception(
                "Failed to load the provided template document because file could not be found: %s",
                self.template_loc,
            )
            raise PptxPackageNotFoundError from pptx.exc.PackageNotFoundError
        except Exception:
            logger.exception(
                "Failed to load the provided template document for unknown reason: %s",
                self.template_loc,
            )

        # Loop through the findings to create slides
        findings_stats = {}

        def get_textframe(shape):
            """
            Get the shape's text frame and enable automatic resizing. The resize only
            triggers after opening the file in the PowerPoint application and making a change or saving.
            """
            text_frame = shape.text_frame
            text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            return text_frame

        def write_bullet(text_frame, text, level):
            """Write a bullet to the provided text frame at the specified level."""
            p = text_frame.add_paragraph()
            p.text = text
            p.level = level

        def write_objective_list(text_frame, objectives):
            """Write a list of objectives to the provided text frame."""
            for obj in objectives:
                status = obj["status"]
                if obj["complete"]:
                    status = "Achieved"
                write_bullet(text_frame, f"{obj['objective']} – {status}", 1)

        def prepare_for_pptx(value):
            """Strip HTML and clear 0x0D characters to prepare text for notes slides."""
            try:
                if value:
                    return BeautifulSoup(value, "lxml").text.replace("\x0D", "")
                return "N/A"
            except Exception:
                logger.exception("Failed parsing this value for PPTX: %s", value)
                return ""

        def add_slide_number(txtbox):
            """
            Add a slide number to the provided textbox. Ideally, the textbox should be the slide layout's slide
            number placeholder to match the template.

            Ref: https://stackoverflow.com/a/55816723
            """
            # Get a textbox's paragraph element
            par = txtbox.text_frame.paragraphs[0]._p

            # The slide number is actually a field, so we add a `fld` element to the paragraph
            # The number enclosed in the `a:t` element is the slide number and should auto-update on load/shuffle
            fld_xml = (
                '<a:fld %s id="{1F4E2DE4-8ADA-4D4E-9951-90A1D26586E7}" type="slidenum">\n'
                '  <a:rPr lang="en-US" smtClean="0"/>\n'
                "  <a:t>2</a:t>\n"
                "</a:fld>\n" % nsdecls("a")
            )
            fld = parse_xml(fld_xml)
            par.append(fld)

        def clone_placeholder(slide, slide_layout, placeholder_idx):
            """
            Clone a placeholder from the slide master and return the layout and the new shape.
            """
            layout_placeholder = slide_layout.placeholders[placeholder_idx]
            slide.shapes.clone_placeholder(layout_placeholder)

            # The cloned placeholder is now the last shape in the slide
            return layout_placeholder, slide.shapes[-1]

        # Calculate finding stats
        for finding in self.report_json["findings"]:
            findings_stats[finding["severity"]] = 0

        for finding in self.report_json["findings"]:
            findings_stats[finding["severity"]] += 1

        # Slide styles (From Master Style counting top to bottom from 0..n)
        SLD_LAYOUT_TITLE = 0
        SLD_LAYOUT_TITLE_AND_CONTENT = 1
        SLD_LAYOUT_FINAL = 12

        # Add a title slide
        slide_layout = self.ppt_presentation.slide_layouts[SLD_LAYOUT_TITLE]
        slide = self.ppt_presentation.slides.add_slide(slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        title_shape.text = f'{self.report_json["client"]["name"]} {self.report_json["project"]["type"]}'
        text_frame = get_textframe(body_shape)
        # Use ``text_frame.text`` for first line/paragraph or ``text_frame.paragraphs[0]``
        text_frame.text = "Technical Outbrief"
        p = text_frame.add_paragraph()
        p.text = dateformat(date.today(), settings.DATE_FORMAT)

        # Add Agenda slide
        slide_layout = self.ppt_presentation.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
        slide = self.ppt_presentation.slides.add_slide(slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.text = "Agenda"
        body_shape = shapes.placeholders[1]
        text_frame = get_textframe(body_shape)
        text_frame.clear()
        self._delete_paragraph(text_frame.paragraphs[0])

        write_bullet(text_frame, "Introduction", 0)
        write_bullet(text_frame, "Assessment Details", 0)
        write_bullet(text_frame, "Methodology", 0)
        write_bullet(text_frame, "Assessment Timeline", 0)
        write_bullet(text_frame, "Attack Path Overview", 0)
        write_bullet(text_frame, "Positive Control Observations", 0)
        write_bullet(text_frame, "Findings and Recommendations Overview", 0)
        write_bullet(text_frame, "Next Steps", 0)

        # Add Introduction slide
        slide_layout = self.ppt_presentation.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
        slide = self.ppt_presentation.slides.add_slide(slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.text = "Introduction"
        body_shape = shapes.placeholders[1]
        text_frame = get_textframe(body_shape)
        text_frame.clear()

        if self.report_json["team"]:
            # Frame needs at least one paragraph to be valid, so don't delete the paragraph
            # if there are no team members
            self._delete_paragraph(text_frame.paragraphs[0])
            for member in self.report_json["team"]:
                write_bullet(text_frame, f"{member['name']} – {member['role']}", 0)
                write_bullet(text_frame, member["email"], 1)

        # Add Assessment Details slide
        slide_layout = self.ppt_presentation.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
        slide = self.ppt_presentation.slides.add_slide(slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.text = "Assessment Details"
        body_shape = shapes.placeholders[1]
        text_frame = get_textframe(body_shape)
        text_frame.clear()
        self._delete_paragraph(text_frame.paragraphs[0])

        write_bullet(
            text_frame, f"{self.report_json['project']['type']} assessment of {self.report_json['client']['name']}", 0
        )
        write_bullet(
            text_frame,
            f"Testing performed from {self.report_json['project']['start_date']} to {self.report_json['project']['end_date']}",
            1,
        )

        finding_body_shape = shapes.placeholders[1]
        self._process_rich_text_pptx(
            self.report_json["project"]["note"],
            slide=slide,
            shape=finding_body_shape,
            template_vars=base_context,
            evidences=base_evidences,
        )
        # The  method adds a new paragraph, so we need to get the last one to increase the indent level
        text_frame = get_textframe(finding_body_shape)
        p = text_frame.paragraphs[-1]
        p.level = 1

        if self.report_json["objectives"]:
            primary_objs = []
            secondary_objs = []
            tertiary_objs = []
            for objective in self.report_json["objectives"]:
                if objective["priority"] == "Primary":
                    primary_objs.append(objective)
                elif objective["priority"] == "Secondary":
                    secondary_objs.append(objective)
                elif objective["priority"] == "Tertiary":
                    tertiary_objs.append(objective)

            if primary_objs:
                write_bullet(text_frame, "Primary Objectives", 0)
                write_objective_list(text_frame, primary_objs)

            if secondary_objs:
                write_bullet(text_frame, "Secondary Objectives", 0)
                write_objective_list(text_frame, secondary_objs)

            if tertiary_objs:
                write_bullet(text_frame, "Tertiary Objectives", 0)
                write_objective_list(text_frame, tertiary_objs)

        # Add Methodology slide
        slide_layout = self.ppt_presentation.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
        slide = self.ppt_presentation.slides.add_slide(slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.text = "Methodology"

        # Add Timeline slide
        slide_layout = self.ppt_presentation.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
        slide = self.ppt_presentation.slides.add_slide(slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.text = "Assessment Timeline"

        # Delete the default text placeholder
        textbox = shapes[1]
        sp = textbox.element
        sp.getparent().remove(sp)
        # Add a table
        rows = 4
        columns = 2
        left = Inches(1.5)
        top = Inches(2)
        width = Inches(8)
        height = Inches(0.8)
        table = shapes.add_table(rows, columns, left, top, width, height).table
        # Set column width
        table.columns[0].width = Inches(2.0)
        table.columns[1].width = Inches(8.5)
        # Write table headers
        cell = table.cell(0, 0)
        cell.text = "Date"
        cell.fill.solid()
        cell.fill.fore_color.rgb = pptx.dml.color.RGBColor(0x2D, 0x28, 0x69)
        cell = table.cell(0, 1)
        cell.text = "Action Item"
        cell.fill.solid()
        cell.fill.fore_color.rgb = pptx.dml.color.RGBColor(0x2D, 0x28, 0x69)

        # Write date rows
        row_iter = 1
        table.cell(row_iter, 0).text = self.report_json["project"]["start_date"]
        table.cell(row_iter, 1).text = "Assessment execution began"
        row_iter += 1
        table.cell(row_iter, 0).text = self.report_json["project"]["end_date"]
        table.cell(row_iter, 1).text = "Assessment execution completed"
        row_iter += 1
        table.cell(row_iter, 0).text = self.report_json["project"]["end_date"]
        table.cell(row_iter, 1).text = "Draft report delivery"

        # Set all cells alignment to center and vertical center
        for cell in table.iter_cells():
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Add Attack Path Overview slide
        slide_layout = self.ppt_presentation.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
        slide = self.ppt_presentation.slides.add_slide(slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.text = "Attack Path Overview"

        # Add Observations slide
        slide_layout = self.ppt_presentation.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
        slide = self.ppt_presentation.slides.add_slide(slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        title_shape.text = "Positive Observations"
        text_frame = get_textframe(body_shape)

        # If there are observations then write a table
        if len(self.report_json["observations"]) > 0:
            # Delete the default text placeholder
            textbox = shapes[1]
            sp = textbox.element
            sp.getparent().remove(sp)
            # Add a table
            rows = len(self.report_json["observations"]) + 1
            columns = 1
            left = Inches(1.5)
            top = Inches(2)
            width = Inches(8)
            height = Inches(0.8)
            table = shapes.add_table(rows, columns, left, top, width, height).table
            # Set column width
            table.columns[0].width = Inches(10.5)
            # Write table headers
            cell = table.cell(0, 0)
            cell.text = "Observation"
            cell.fill.solid()
            cell.fill.fore_color.rgb = pptx.dml.color.RGBColor(0x2D, 0x28, 0x69)
            # Write findings rows
            row_iter = 1
            for observation in self.report_json["observations"]:
                table.cell(row_iter, 0).text = observation["title"]
                row_iter += 1
            # Set all cells alignment to center and vertical center
            for cell in table.iter_cells():
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        else:
            write_bullet(text_frame, "No observations", 0)

        # Create slide for each observation
        for observation in self.report_json["observations"]:
            slide_layout = self.ppt_presentation.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
            observation_slide = self.ppt_presentation.slides.add_slide(slide_layout)
            shapes = observation_slide.shapes
            title_shape = shapes.title

            # Prepare text frame
            observation_body_shape = shapes.placeholders[1]
            if observation_body_shape.has_text_frame:
                text_frame = get_textframe(observation_body_shape)
                text_frame.clear()
                self._delete_paragraph(text_frame.paragraphs[0])
            else:
                text_frame = None

            # Set slide title to title + [severity]
            title_shape.text = f'{observation["title"]}'

            # Add description to the slide body (other sections will appear in the notes)
            if observation.get("description", "").strip():
                observation_context = self._jinja_richtext_base_context()
                self._process_rich_text_pptx(
                    observation["description"],
                    slide=observation_slide,
                    shape=observation_body_shape,
                    template_vars=observation_context,
                    evidences=base_evidences,
                )
            else:
                par = observation_body_shape.add_paragraph()
                run = par.add_run()
                run.text = "No description provided"

            for ev in observation.get("evidence", []):
                HtmlToPptxWithEvidence.make_evidence(observation_slide, ev)

            # Ensure there is at least one paragraph, as required by the spec
            if text_frame is not None and not text_frame.paragraphs:
                text_frame.add_paragraph()

        # Add Findings Overview Slide
        slide_layout = self.ppt_presentation.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
        slide = self.ppt_presentation.slides.add_slide(slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        title_shape.text = "Findings Overview"
        text_frame = get_textframe(body_shape)

        # If there are findings then write a table of findings and severity ratings
        if len(self.report_json["findings"]) > 0:
            # Delete the default text placeholder
            textbox = shapes[1]
            sp = textbox.element
            sp.getparent().remove(sp)
            # Add a table
            rows = len(self.report_json["findings"]) + 1
            columns = 2
            left = Inches(1.5)
            top = Inches(2)
            width = Inches(8)
            height = Inches(0.8)
            table = shapes.add_table(rows, columns, left, top, width, height).table
            # Set column width
            table.columns[0].width = Inches(8.5)
            table.columns[1].width = Inches(2.0)
            # Write table headers
            cell = table.cell(0, 0)
            cell.text = "Finding"
            cell.fill.solid()
            cell.fill.fore_color.rgb = pptx.dml.color.RGBColor(0x2D, 0x28, 0x69)
            cell = table.cell(0, 1)
            cell.text = "Severity"
            cell.fill.solid()
            cell.fill.fore_color.rgb = pptx.dml.color.RGBColor(0x2D, 0x28, 0x69)
            # Write findings rows
            row_iter = 1
            for finding in self.report_json["findings"]:
                table.cell(row_iter, 0).text = finding["title"]
                risk_cell = table.cell(row_iter, 1)
                # Set risk rating
                risk_cell.text = finding["severity"]
                # Set cell color fill type to solid
                risk_cell.fill.solid()
                # Color the risk cell based on corresponding severity color
                cell_color = pptx.dml.color.RGBColor(*map(lambda v: int(v, 16), finding["severity_color_hex"]))
                risk_cell.fill.fore_color.rgb = cell_color
                row_iter += 1
            # Set all cells alignment to center and vertical center
            for cell in table.iter_cells():
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        else:
            write_bullet(text_frame, "No findings", 0)

        # Create slide for each finding
        for finding in self.report_json["findings"]:
            slide_layout = self.ppt_presentation.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
            finding_slide = self.ppt_presentation.slides.add_slide(slide_layout)
            shapes = finding_slide.shapes
            title_shape = shapes.title

            # Prepare text frame
            finding_body_shape = shapes.placeholders[1]
            if finding_body_shape.has_text_frame:
                text_frame = get_textframe(finding_body_shape)
                text_frame.clear()
                self._delete_paragraph(text_frame.paragraphs[0])
            else:
                text_frame = None

            # Set slide title to title + [severity]
            title_shape.text = f'{finding["title"]} [{finding["severity"]}]'

            # Add description to the slide body (other sections will appear in the notes)
            if finding.get("description", "").strip():
                finding_context = self._jinja_richtext_finding_context(base_context, finding)
                finding_evidences = base_evidences | {e["friendly_name"]: e for e in finding["evidence"]}
                self._process_rich_text_pptx(
                    finding["description"],
                    slide=finding_slide,
                    shape=finding_body_shape,
                    template_vars=finding_context,
                    evidences=finding_evidences,
                )
            else:
                par = finding_body_shape.add_paragraph()
                run = par.add_run()
                run.text = "No description provided"

            for ev in finding.get("evidence", []):
                HtmlToPptxWithEvidence.make_evidence(finding_slide, ev)

            # Ensure there is at least one paragraph, as required by the spec
            if text_frame is not None and not text_frame.paragraphs:
                text_frame.add_paragraph()

            # Add all finding data to the notes section for easier reference during edits
            entities = prepare_for_pptx(finding["affected_entities"])
            impact = prepare_for_pptx(finding["impact"])
            host_detection = prepare_for_pptx(finding["host_detection_techniques"])
            net_detection = prepare_for_pptx(finding["network_detection_techniques"])
            recommendation = prepare_for_pptx(finding["recommendation"])
            replication = prepare_for_pptx(finding["replication_steps"])
            references = prepare_for_pptx(finding["references"])
            notes_slide = finding_slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            p = text_frame.add_paragraph()
            p.text = f"""
                {finding["severity"].capitalize()}: finding["title"]

                AFFECTED ENTITIES
                {entities}

                IMPACT
                {impact}

                MITIGATION
                {recommendation}

                REPLICATION
                {replication}

                HOST DETECTION
                {host_detection}

                NETWORK DETECTION
                ,
                {net_detection}

                REFERENCES
                {references}
            """.replace(
                "                ", ""
            )

        # Add Recommendations slide
        slide_layout = self.ppt_presentation.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
        slide = self.ppt_presentation.slides.add_slide(slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.text = "Recommendations"

        # Add Next Steps slide
        slide_layout = self.ppt_presentation.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
        slide = self.ppt_presentation.slides.add_slide(slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.text = "Next Steps"

        # Add final slide
        slide_layout = self.ppt_presentation.slide_layouts[SLD_LAYOUT_FINAL]
        slide = self.ppt_presentation.slides.add_slide(slide_layout)
        shapes = slide.shapes
        body_shape = shapes.placeholders[1]
        text_frame = get_textframe(body_shape)
        text_frame.clear()
        p = text_frame.paragraphs[0]
        p.line_spacing = 0.7
        p.text = self.company_config.company_name
        p = text_frame.add_paragraph()
        p.text = self.company_config.company_twitter
        p.line_spacing = 0.7
        p = text_frame.add_paragraph()
        p.text = self.company_config.company_email
        p.line_spacing = 0.7

        # Add footer elements (if there is one) to all slides based on the footer placeholder in the template
        for idx, slide in enumerate(self.ppt_presentation.slides):
            date_placeholder_idx = -1
            footer_placeholder_idx = -1
            slide_number_placeholder_idx = -1
            slide_layout = slide.slide_layout

            for idx, place in enumerate(slide_layout.placeholders):
                if "Footer" in place.name:
                    footer_placeholder_idx = idx
                if "Slide Number" in place.name:
                    slide_number_placeholder_idx = idx
                if "Date" in place.name:
                    date_placeholder_idx = idx

            # Skip the title slide at index 0
            if idx > 0:
                if footer_placeholder_idx > 0:
                    footer_layout_placeholder, footer_placeholder = clone_placeholder(
                        slide, slide_layout, footer_placeholder_idx
                    )
                    footer_placeholder.text = footer_layout_placeholder.text
                if slide_number_placeholder_idx > 0:
                    slide_number_layout_placeholder, slide_number_placeholder = clone_placeholder(
                        slide, slide_layout, slide_number_placeholder_idx
                    )
                    add_slide_number(slide_number_placeholder)
                if date_placeholder_idx > 0:
                    date_layout_placeholder, date_placeholder = clone_placeholder(
                        slide, slide_layout, date_placeholder_idx
                    )
                    date_placeholder.text = dateformat(date.today(), settings.DATE_FORMAT)

        # Finalize document and return it for an HTTP response
        return self.ppt_presentation

    def generate_all_reports(self, docx_template, pptx_template):
        """Generate all available report types and return memory streams for each file."""

        # Generate the docx report - save it in a memory stream
        word_stream = io.BytesIO()
        self.template_loc = docx_template
        word_doc = self.generate_word_docx()
        word_doc.save(word_stream)
        # Generate the xlsx report - save it in a memory stream
        excel_stream = io.BytesIO()
        self.generate_excel_xlsx(excel_stream)
        # Generate the pptx report - save it in a memory stream
        ppt_stream = io.BytesIO()
        self.template_loc = pptx_template
        ppt_doc = self.generate_powerpoint_pptx()
        ppt_doc.save(ppt_stream)
        # Return each memory object
        return json.dumps(self.report_json, indent=4), word_stream, excel_stream, ppt_stream

    def _delete_paragraph(self, par):
        """
        Delete the specified paragraph.

        **Parameter**

        ``par``
            Paragraph to delete from the document
        """
        p = par._p
        parent_element = p.getparent()
        if parent_element is not None:
            parent_element.remove(p)
        else:
            logger.warning("Could not delete paragraph in because it had no parent element")


class TemplateLinter:
    """Lint template files to catch undefined variables and syntax errors."""

    def __init__(self, template):
        self.template = template
        self.template_loc = template.document.path
        self.jinja_template_env = prepare_jinja2_env(debug=True)

    def lint_docx(self):
        """
        Lint the provided Word docx file from :model:`reporting.ReportTemplate`.
        """
        results = {"result": "success", "warnings": [], "errors": []}
        if self.template_loc:
            if os.path.exists(self.template_loc):
                logger.info("Found template file at %s", self.template_loc)
                try:
                    # Step 1: Load the document as a template
                    template_document = DocxTemplate(self.template_loc)
                    logger.info("Template loaded for linting")

                    undefined_vars = template_document.get_undeclared_template_variables(self.jinja_template_env)
                    if undefined_vars:
                        for variable in undefined_vars:
                            if variable not in LINTER_CONTEXT:
                                results["warnings"].append(f"Potential undefined variable: {variable}")

                    # Step 2: Check document's styles
                    document_styles = template_document.styles
                    if "Bullet List" not in document_styles:
                        results["warnings"].append(
                            "Template is missing a recommended style (see documentation): Bullet List"
                        )
                    if "Number List" not in document_styles:
                        results["warnings"].append(
                            "Template is missing a recommended style (see documentation): Number List"
                        )
                    if "CodeBlock" not in document_styles:
                        results["warnings"].append(
                            "Template is missing a recommended style (see documentation): CodeBlock"
                        )
                    if "CodeInline" not in document_styles:
                        results["warnings"].append(
                            "Template is missing a recommended style (see documentation): CodeInline"
                        )
                    if "Caption" not in document_styles:
                        results["warnings"].append(
                            "Template is missing a recommended style (see documentation): Caption"
                        )
                    if "List Paragraph" not in document_styles:
                        results["warnings"].append(
                            "Template is missing a recommended style (see documentation): List Paragraph"
                        )
                    if "Blockquote" not in document_styles:
                        results["warnings"].append(
                            "Template is missing a recommended style (see documentation): Blockquote"
                        )
                    if "Blockquote" not in document_styles:
                        results["warnings"].append(
                            "Template is missing a recommended style (see documentation): Blockquote"
                        )
                    if "Table Grid" not in document_styles:
                        results["errors"].append("Template is missing a required style (see documentation): Table Grid")
                    if self.template.p_style:
                        if self.template.p_style not in document_styles:
                            results["warnings"].append(
                                f"Template is missing your configured default paragraph style: {self.template.p_style}"
                            )

                    if results["warnings"]:
                        results["result"] = "warning"

                    logger.info("Completed Word style checks")

                    # Step 3: Prepare context
                    context = copy.deepcopy(LINTER_CONTEXT)
                    for field in ExtraFieldSpec.objects.filter(target_model=Report._meta.label):
                        context["extra_fields"][field.internal_name] = field.default_value()
                    for field in ExtraFieldSpec.objects.filter(target_model=Project._meta.label):
                        context["project"]["extra_fields"][field.internal_name] = field.default_value()
                    for field in ExtraFieldSpec.objects.filter(target_model=Client._meta.label):
                        context["client"]["extra_fields"][field.internal_name] = field.default_value()
                    for field in ExtraFieldSpec.objects.filter(target_model=Finding._meta.label):
                        for finding in context["findings"]:
                            finding["extra_fields"][field.internal_name] = field.default_value()
                    for field in ExtraFieldSpec.objects.filter(target_model=OplogEntry._meta.label):
                        for log in context["logs"]:
                            for entry in log["entries"]:
                                entry["extra_fields"][field.internal_name] = field.default_value()
                    for field in ExtraFieldSpec.objects.filter(target_model=Domain._meta.label):
                        for domain in context["infrastructure"]["domains"]:
                            domain["extra_fields"][field.internal_name] = field.default_value()
                    for field in ExtraFieldSpec.objects.filter(target_model=StaticServer._meta.label):
                        for server in context["infrastructure"]["servers"]:
                            server["extra_fields"][field.internal_name] = field.default_value()
                    for field in ExtraFieldSpec.objects.filter(target_model=Observation._meta.label):
                        for obs in context["observations"]:
                            obs["extra_fields"][field.internal_name] = field.default_value()

                    # Step 4: Test rendering the document
                    try:
                        template_document.render(context, self.jinja_template_env, autoescape=True)
                        undefined_vars = template_document.undeclared_template_variables
                        if undefined_vars:
                            for variable in undefined_vars:
                                results["warnings"].append(f"Undefined variable: {variable}")
                        if results["warnings"]:
                            results["result"] = "warning"
                        logger.info("Completed document rendering test")
                    except TemplateSyntaxError as error:
                        logger.exception("Template syntax error: %s", error)
                        results = {
                            "result": "failed",
                            "errors": [f"Jinja2 template syntax error: {error.message}"],
                        }
                        if error.message == "expected token 'end of print statement', got 'such'":
                            results["errors"].append(
                                "The above error means you may have a typo in a variable name or expression"
                            )
                    except UndefinedError as error:
                        logger.error("Template syntax error: %s", error)
                        results = {
                            "result": "failed",
                            "errors": [f"Jinja2 template syntax error: {error.message}"],
                        }
                    except InvalidFilterValue as error:
                        logger.error("Invalid value provided to filter: %s", error)
                        results = {
                            "result": "failed",
                            "errors": [f"Invalid filter value: {error.message}"],
                        }
                    except TypeError as error:
                        logger.error("Invalid value provided to filter or expression: %s", error)
                        results = {
                            "result": "failed",
                            "errors": [f"Invalid value provided to filter or expression: {error}"],
                        }
                    except TemplateRuntimeError as error:
                        logger.error("Invalid filter or expression: %s", error)
                        results = {
                            "result": "failed",
                            "errors": [f"Invalid filter or expression: {error}"],
                        }
                except Exception:
                    logger.exception("Template failed rendering")
                    results = {
                        "result": "failed",
                        "errors": ["Template rendering failed unexpectedly"],
                    }
            else:
                logger.error("Template file path did not exist: %s", self.template_loc)
                results = {
                    "result": "failed",
                    "errors": ["Template file does not exist – upload it again"],
                }
        else:
            logger.error("Received a `None` value for template location")

        logger.info("Template linting completed")
        return results

    def lint_pptx(self):
        """
        Lint the provided PowerPoint pptx file from :model:`reporting.ReportTemplate`.
        """
        results = {"result": "success", "warnings": [], "errors": []}
        if self.template_loc:
            if os.path.exists(self.template_loc):
                logger.info("Found template file at %s", self.template_loc)
                try:
                    # Test 1: Check if the document is a PPTX file
                    template_document = Presentation(self.template_loc)

                    # Test 2: Check for existing slides
                    slide_count = len(template_document.slides)
                    logger.info("Slide count was %s", slide_count)
                    if slide_count > 0:
                        results["warnings"].append(
                            "Template can be used, but it has slides when it should be empty (see documentation)"
                        )
                except ValueError:
                    logger.exception(
                        "Failed to load the provided template document because it is not a PowerPoint file: %s",
                        self.template_loc,
                    )
                    results = {
                        "result": "failed",
                        "errors": ["Template file is not a PowerPoint presentation"],
                    }
                except TypeError as error:
                    logger.error("Invalid value provided to filter or expression: %s", error)
                    results = {
                        "result": "failed",
                        "errors": [f"Invalid value provided to filter or expression: {error}"],
                    }
                except TemplateRuntimeError as error:
                    logger.error("Invalid filter or expression: %s", error)
                    results = {
                        "result": "failed",
                        "errors": [f"Invalid filter or expression: {error}"],
                    }
                except Exception:
                    logger.exception("Template failed rendering")
                    results = {
                        "result": "failed",
                        "errors": ["Template rendering failed unexpectedly"],
                    }
            else:
                logger.error("Template file path did not exist: %s", self.template_loc)
                results = {
                    "result": "failed",
                    "errors": ["Template file does not exist – upload it again"],
                }
        else:
            logger.error("Received a `None` value for template location")

        logger.info("Template linting completed")
        return results


def _jinja_evidence(evidence_name):
    """
    `{{evidence(name)}}` function in jinja.
    """
    return Markup('<span data-gw-evidence="' + html.escape(evidence_name) + '"></span>')


def _jinja_caption(caption_name):
    return Markup('<span data-gw-caption="' + html.escape(caption_name) + '"></span>')


def _jinja_ref(ref_name):
    return Markup('<span data-gw-ref="' + html.escape(ref_name) + '"></span>')
