
import io

from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches
import pptx

from ghostwriter.modules.reportwriter.base.pptx import SLD_LAYOUT_FINAL, SLD_LAYOUT_TITLE_AND_CONTENT, ExportBasePptx, delete_paragraph, get_textframe, prepare_for_pptx, write_bullet
from ghostwriter.modules.reportwriter.project.pptx import ProjectSlidesMixin
from ghostwriter.modules.reportwriter.report.base import ExportReportBase
from ghostwriter.modules.reportwriter.richtext.pptx import HtmlToPptxWithEvidence


class ExportReportPptx(ExportBasePptx, ExportReportBase, ProjectSlidesMixin):
    def run(self) -> io.BytesIO:
        """Generate a complete PowerPoint slide deck for the current report."""

        base_context = self.map_rich_texts()

        # Loop through the findings to create slides
        findings_stats = {}

        # Calculate finding stats
        for finding in base_context["findings"]:
            findings_stats[finding["severity"]] = 0

        for finding in base_context["findings"]:
            findings_stats[finding["severity"]] += 1

        self.create_project_slides(base_context)

        # Add Observations slide
        slide_layout = self.ppt_presentation.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
        slide = self.ppt_presentation.slides.add_slide(slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        title_shape.text = "Positive Observations"
        text_frame = get_textframe(body_shape)

        # If there are observations then write a table
        if len(base_context["observations"]) > 0:
            # Delete the default text placeholder
            textbox = shapes[1]
            sp = textbox.element
            sp.getparent().remove(sp)
            # Add a table
            rows = len(base_context["observations"]) + 1
            columns = 1
            left = Inches(1.5)
            top = Inches(2)
            width = Inches(8)
            height = Inches(0.8)
            table = shapes.add_table(rows, columns, left, top, width, height).table
            # Set column width
            table.columns[0].width = Inches(10.5)
            # Write table headers
            cell = table.cell(0, 0)
            cell.text = "Observation"
            cell.fill.solid()
            cell.fill.fore_color.rgb = pptx.dml.color.RGBColor(0x2D, 0x28, 0x69)
            # Write findings rows
            row_iter = 1
            for observation in base_context["observations"]:
                table.cell(row_iter, 0).text = observation["title"]
                row_iter += 1
            # Set all cells alignment to center and vertical center
            for cell in table.iter_cells():
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        else:
            write_bullet(text_frame, "No observations", 0)

        # Create slide for each observation
        for observation in base_context["observations"]:
            slide_layout = self.ppt_presentation.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
            observation_slide = self.ppt_presentation.slides.add_slide(slide_layout)
            shapes = observation_slide.shapes
            title_shape = shapes.title

            # Prepare text frame
            observation_body_shape = shapes.placeholders[1]
            if observation_body_shape.has_text_frame:
                text_frame = get_textframe(observation_body_shape)
                text_frame.clear()
                delete_paragraph(text_frame.paragraphs[0])
            else:
                text_frame = None

            # Set slide title to title + [severity]
            title_shape.text = f'{observation["title"]}'

            # Add description to the slide body (other sections will appear in the notes)
            if observation.get("description", "").strip():
                self.render_rich_text_pptx(
                    observation["description_rt"],
                    slide=observation_slide,
                    shape=observation_body_shape,
                )
            else:
                par = observation_body_shape.text_frame.add_paragraph()
                run = par.add_run()
                run.text = "No description provided"

            for ev in observation.get("evidence", []):
                HtmlToPptxWithEvidence.make_evidence(observation_slide, ev)

            # Ensure there is at least one paragraph, as required by the spec
            if text_frame is not None and not text_frame.paragraphs:
                text_frame.add_paragraph()

        # Add Findings Overview Slide
        slide_layout = self.ppt_presentation.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
        slide = self.ppt_presentation.slides.add_slide(slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        title_shape.text = "Findings Overview"
        text_frame = get_textframe(body_shape)

        # If there are findings then write a table of findings and severity ratings
        if len(base_context["findings"]) > 0:
            # Delete the default text placeholder
            textbox = shapes[1]
            sp = textbox.element
            sp.getparent().remove(sp)
            # Add a table
            rows = len(base_context["findings"]) + 1
            columns = 2
            left = Inches(1.5)
            top = Inches(2)
            width = Inches(8)
            height = Inches(0.8)
            table = shapes.add_table(rows, columns, left, top, width, height).table
            # Set column width
            table.columns[0].width = Inches(8.5)
            table.columns[1].width = Inches(2.0)
            # Write table headers
            cell = table.cell(0, 0)
            cell.text = "Finding"
            cell.fill.solid()
            cell.fill.fore_color.rgb = pptx.dml.color.RGBColor(0x2D, 0x28, 0x69)
            cell = table.cell(0, 1)
            cell.text = "Severity"
            cell.fill.solid()
            cell.fill.fore_color.rgb = pptx.dml.color.RGBColor(0x2D, 0x28, 0x69)
            # Write findings rows
            row_iter = 1
            for finding in base_context["findings"]:
                table.cell(row_iter, 0).text = finding["title"]
                risk_cell = table.cell(row_iter, 1)
                # Set risk rating
                risk_cell.text = finding["severity"]
                # Set cell color fill type to solid
                risk_cell.fill.solid()
                # Color the risk cell based on corresponding severity color
                cell_color = pptx.dml.color.RGBColor(*map(lambda v: int(v, 16), finding["severity_color_hex"]))
                risk_cell.fill.fore_color.rgb = cell_color
                row_iter += 1
            # Set all cells alignment to center and vertical center
            for cell in table.iter_cells():
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        else:
            write_bullet(text_frame, "No findings", 0)

        # Create slide for each finding
        for finding in base_context["findings"]:
            slide_layout = self.ppt_presentation.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
            finding_slide = self.ppt_presentation.slides.add_slide(slide_layout)
            shapes = finding_slide.shapes
            title_shape = shapes.title

            # Prepare text frame
            finding_body_shape = shapes.placeholders[1]
            if finding_body_shape.has_text_frame:
                text_frame = get_textframe(finding_body_shape)
                text_frame.clear()
                delete_paragraph(text_frame.paragraphs[0])
            else:
                text_frame = None

            # Set slide title to title + [severity]
            title_shape.text = f'{finding["title"]} [{finding["severity"]}]'

            # Add description to the slide body (other sections will appear in the notes)
            if finding.get("description", "").strip():
                self.render_rich_text_pptx(
                    finding["description_rt"],
                    slide=finding_slide,
                    shape=finding_body_shape,
                )
            else:
                par = text_frame.add_paragraph()
                run = par.add_run()
                run.text = "No description provided"

            for ev in finding.get("evidence", []):
                HtmlToPptxWithEvidence.make_evidence(finding_slide, ev)

            # Ensure there is at least one paragraph, as required by the spec
            if text_frame is not None and not text_frame.paragraphs:
                text_frame.add_paragraph()

            # Add all finding data to the notes section for easier reference during edits
            entities = prepare_for_pptx(finding["affected_entities"])
            impact = prepare_for_pptx(finding["impact"])
            host_detection = prepare_for_pptx(finding["host_detection_techniques"])
            net_detection = prepare_for_pptx(finding["network_detection_techniques"])
            recommendation = prepare_for_pptx(finding["recommendation"])
            replication = prepare_for_pptx(finding["replication_steps"])
            references = prepare_for_pptx(finding["references"])
            notes_slide = finding_slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            p = text_frame.add_paragraph()
            p.text = f"""
                {finding["severity"].capitalize()}: finding["title"]

                AFFECTED ENTITIES
                {entities}

                IMPACT
                {impact}

                MITIGATION
                {recommendation}

                REPLICATION
                {replication}

                HOST DETECTION
                {host_detection}

                NETWORK DETECTION
                ,
                {net_detection}

                REFERENCES
                {references}
            """.replace(
                "                ", ""
            )

        # Add Recommendations slide
        slide_layout = self.ppt_presentation.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
        slide = self.ppt_presentation.slides.add_slide(slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.text = "Recommendations"

        # Add Next Steps slide
        slide_layout = self.ppt_presentation.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
        slide = self.ppt_presentation.slides.add_slide(slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.text = "Next Steps"

        # Add final slide
        slide_layout = self.ppt_presentation.slide_layouts[SLD_LAYOUT_FINAL]
        slide = self.ppt_presentation.slides.add_slide(slide_layout)
        shapes = slide.shapes
        body_shape = shapes.placeholders[1]
        text_frame = get_textframe(body_shape)
        text_frame.clear()
        p = text_frame.paragraphs[0]
        p.line_spacing = 0.7
        p.text = self.company_config.company_name
        p = text_frame.add_paragraph()
        p.text = self.company_config.company_twitter
        p.line_spacing = 0.7
        p = text_frame.add_paragraph()
        p.text = self.company_config.company_email
        p.line_spacing = 0.7

        self.process_footers()

        return super().run()
