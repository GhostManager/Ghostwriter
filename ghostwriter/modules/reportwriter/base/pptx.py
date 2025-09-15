
import io
import logging
import os
from datetime import date
from typing import List, Tuple

from django.conf import settings
from django.utils.dateformat import format as dateformat
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.parts.presentation import PresentationPart
from pptx.exc import PackageNotFoundError
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
from pptx.enum.text import MSO_AUTO_SIZE

from ghostwriter.commandcenter.models import CompanyInformation
from ghostwriter.modules.reportwriter.base import ReportExportTemplateError
from ghostwriter.modules.reportwriter.base.base import ExportBase
from ghostwriter.modules.reportwriter.base.html_rich_text import LazilyRenderedTemplate
from ghostwriter.modules.reportwriter.richtext.pptx import HtmlToPptxWithEvidence

logger = logging.getLogger(__name__)


class ExportBasePptx(ExportBase):
    """
    Base class for exporting Pptx (PowerPoint) files

    Subclasses should override `run` to add slides to the `ppt_presentation` field, using `process_rich_text_pptx`
    to template and convert rich text fields, then return `super().run()` to save and return the presentation.
    """
    ppt_presentation: PresentationPart
    company_config: CompanyInformation
    linting: bool

    @classmethod
    def mime_type(cls) -> str:
        return "application/vnd.openxmlformats-officedocument.presentationml.presentation"

    @classmethod
    def extension(cls) -> str:
        return "pptx"

    def __init__(
        self,
        object,
        *,
        template_loc: str = None,
        linting: bool = False,
        **kwargs
    ):
        if "jinja_debug" not in kwargs:
            kwargs["jinja_debug"] = linting
        super().__init__(object, **kwargs)
        self.linting = linting

        try:
            self.ppt_presentation = Presentation(template_loc)
        except PackageNotFoundError as err:
            raise ReportExportTemplateError("Template document file could not be found - try re-uploading it") from err
        except Exception:
            logger.exception(
                "Failed to load the provided template document for unknown reason: %s",
                template_loc,
            )
            raise

        self.company_config = CompanyInformation.get_solo()

    def render_rich_text_pptx(self, rich_text: LazilyRenderedTemplate, slide, shape):
        """
        Renders a `LazilyRenderedTemplate`, converting the HTML from the TinyMCE rich text editor and inserting it into the passed in shape and slide.
        Converts HTML from the TinyMCE rich text editor and inserts it into the passed in slide and shape
        """
        ReportExportTemplateError.map_errors(
            lambda: HtmlToPptxWithEvidence.run(
                rich_text.render_html(),
                slide=slide,
                shape=shape,
                evidences=self.evidences_by_id,
            ),
            getattr(rich_text, "location", None)
        )

    def process_footers(self):
        """
        Add footer elements (if there is one) to all slides based on the footer placeholder in the template
        """
        for idx, slide in enumerate(self.ppt_presentation.slides):
            date_placeholder_idx = -1
            footer_placeholder_idx = -1
            slide_number_placeholder_idx = -1
            slide_layout = slide.slide_layout

            for idx, place in enumerate(slide_layout.placeholders):
                if "Footer" in place.name:
                    footer_placeholder_idx = idx
                if "Slide Number" in place.name:
                    slide_number_placeholder_idx = idx
                if "Date" in place.name:
                    date_placeholder_idx = idx

            # Skip the title slide at index 0
            if idx > 0:
                if footer_placeholder_idx > 0:
                    footer_layout_placeholder, footer_placeholder = clone_placeholder(
                        slide, slide_layout, footer_placeholder_idx
                    )
                    footer_placeholder.text = footer_layout_placeholder.text
                if slide_number_placeholder_idx > 0:
                    _, slide_number_placeholder = clone_placeholder(
                        slide, slide_layout, slide_number_placeholder_idx
                    )
                    add_slide_number(slide_number_placeholder)
                if date_placeholder_idx > 0:
                    _, date_placeholder = clone_placeholder(
                        slide, slide_layout, date_placeholder_idx
                    )
                    date_placeholder.text = dateformat(date.today(), settings.DATE_FORMAT)

    def run(self):
        out = io.BytesIO()
        self.ppt_presentation.save(out)
        return out

    @classmethod
    def lint(cls, template_loc: str) -> Tuple[List[str], List[str]]:
        warnings = []
        errors = []
        try:
            if not os.path.exists(template_loc):
                logger.error("Template file path did not exist: %r", template_loc)
                errors.append("Template file does not exist – upload it again")
                return warnings, errors

            # Test 1: Check if the document is a PPTX file
            template_document = Presentation(template_loc)

            # Test 2: Check for existing slides
            slide_count = len(template_document.slides)
            logger.info("Slide count was %s", slide_count)
            if slide_count > 0:
                warnings.append(
                    "Template can be used, but it has slides when it should be empty (see documentation)"
                )
        except ReportExportTemplateError as error:
            logger.exception("Template failed linting: %s", error)
            errors.append(f"Linting failed: {error}")
        except Exception:
            logger.exception("Template failed linting")
            errors.append("Template rendering failed unexpectedly")

        logger.info("Linting finished: %d warnings, %d errors", len(warnings), len(errors))
        return warnings, errors


# Slide styles (From Master Style counting top to bottom from 0..n)
SLD_LAYOUT_TITLE = 0
SLD_LAYOUT_TITLE_AND_CONTENT = 1
SLD_LAYOUT_FINAL = 12


def add_slide_number(txtbox):
    """
    Add a slide number to the provided textbox. Ideally, the textbox should be the slide layout's slide
    number placeholder to match the template.

    Ref: https://stackoverflow.com/a/55816723
    """
    # Get a textbox's paragraph element
    par = txtbox.text_frame.paragraphs[0]._p

    # The slide number is actually a field, so we add a `fld` element to the paragraph
    # The number enclosed in the `a:t` element is the slide number and should auto-update on load/shuffle
    fld_xml = (
        '<a:fld %s id="{1F4E2DE4-8ADA-4D4E-9951-90A1D26586E7}" type="slidenum">\n'
        '  <a:rPr lang="en-US" smtClean="0"/>\n'
        "  <a:t>2</a:t>\n"
        "</a:fld>\n" % nsdecls("a")
    )
    fld = parse_xml(fld_xml)
    par.append(fld)


def clone_placeholder(slide, slide_layout, placeholder_idx):
    """
    Clone a placeholder from the slide master and return the layout and the new shape.
    """
    layout_placeholder = slide_layout.placeholders[placeholder_idx]
    slide.shapes.clone_placeholder(layout_placeholder)

    # The cloned placeholder is now the last shape in the slide
    return layout_placeholder, slide.shapes[-1]


def get_textframe(shape):
    """
    Get the shape's text frame and enable automatic resizing. The resize only
    triggers after opening the file in the PowerPoint application and making a change or saving.
    """
    text_frame = shape.text_frame
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    return text_frame


def write_bullet(text_frame, text, level):
    """Write a bullet to the provided text frame at the specified level."""
    p = text_frame.add_paragraph()
    p.text = text
    p.level = level


def write_objective_list(text_frame, objectives):
    """Write a list of objectives to the provided text frame."""
    for obj in objectives:
        status = obj["status"]
        if obj["complete"]:
            status = "Achieved"
        write_bullet(text_frame, f"{obj['objective']} – {status}", 1)


def prepare_for_pptx(value):
    """Strip HTML and clear 0x0D characters to prepare text for notes slides."""
    try:
        if value:
            return BeautifulSoup(value, "lxml").text.replace("\x0D", "")
        return "N/A"
    except Exception:
        logger.exception("Failed parsing this value for PPTX: %s", value)
        return ""


def delete_paragraph(par):
    """
    Delete the specified paragraph.

    **Parameter**

    ``par``
        Paragraph to delete from the document
    """
    p = par._p
    parent_element = p.getparent()
    if parent_element is not None:
        parent_element.remove(p)
    else:
        logger.warning("Could not delete paragraph in because it had no parent element")
