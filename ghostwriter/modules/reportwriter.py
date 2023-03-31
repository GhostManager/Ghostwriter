"""
This module contains the tools required for generating Microsoft Office documents for
reporting. The ``Reportwriter`` class accepts data and produces a docx, xlsx, pptx,
and json using the provided data.
"""

# Standard Libraries
import io
import json
import logging
import os
import random
import re
from datetime import timedelta
from string import ascii_letters

# Django Imports
from django.conf import settings
from django.utils.dateformat import format as dateformat

# 3rd Party Libraries
import docx
import jinja2
import jinja2.sandbox
import pptx
from bs4 import BeautifulSoup, NavigableString
from dateutil.parser import parse as parse_datetime
from dateutil.parser._parser import ParserError
from docx.enum.dml import MSO_THEME_COLOR_INDEX
from docx.enum.style import WD_STYLE_TYPE
from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_COLOR_INDEX
from docx.image.exceptions import UnrecognizedImageError
from docx.opc.exceptions import PackageNotFoundError as DocxPackageNotFoundError
from docx.oxml.shared import OxmlElement, qn
from docx.shared import Inches, Pt, RGBColor
from docxtpl import DocxTemplate, RichText
from jinja2.exceptions import TemplateRuntimeError, TemplateSyntaxError, UndefinedError
from pptx import Presentation
from pptx.dml.color import RGBColor as PptxRGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.exc import PackageNotFoundError as PptxPackageNotFoundError
from rest_framework.renderers import JSONRenderer
from xlsxwriter.workbook import Workbook

# Ghostwriter Libraries
from ghostwriter.commandcenter.models import CompanyInformation, ReportConfiguration
from ghostwriter.modules.custom_serializers import ReportDataSerializer
from ghostwriter.modules.exceptions import InvalidFilterValue
from ghostwriter.modules.linting_utils import LINTER_CONTEXT
from ghostwriter.reporting.models import Evidence

# Custom code
from ghostwriter.stratum.enums import (
    DifficultyExploitColor,
    FindingStatusColor,
    Severity,
    get_value_from_key,
)
from ghostwriter.stratum.findings_chart import build_bar_chart, build_pie_chart
from ghostwriter.stratum.sd_graph import build_sd_graph, plt

# Using __name__ resolves to ghostwriter.modules.reporting
logger = logging.getLogger(__name__)


# Custom Jinja2 filters for DOCX templates
def filter_severity(findings, allowlist):
    """
    Filter list of findings to return only those with a severity in the allowlist.

    **Parameters**

    ``findings``
        List of dictionary objects (JSON) for findings
    ``allowlist``
        List of strings matching severity categories to allow through filter
    """
    filtered_values = []
    if isinstance(allowlist, list):
        allowlist = [severity.lower() for severity in allowlist]
    else:
        raise InvalidFilterValue(
            f'Allowlist passed into `filter_severity()` filter is not a list ("{allowlist}"); must be like `["Critical", "High"]`'
        )
    try:
        for finding in findings:
            if finding["severity"].lower() in allowlist:
                filtered_values.append(finding)
    except (KeyError, TypeError):
        logger.exception("Error parsing ``findings`` as a list of dictionaries: %s", findings)
        raise InvalidFilterValue(
            "Invalid list of findings passed into `filter_severity()` filter; must be the `{{ findings }}` object"
        )
    return filtered_values


def filter_type(findings, allowlist):
    """
    Filter list of findings to return only those with a type in the allowlist.

    **Parameters**

    ``findings``
        List of dictionary objects (JSON) for findings
    ``allowlist``
        List of strings matching severity categories to allow through filter
    """
    filtered_values = []
    if isinstance(allowlist, list):
        allowlist = [type.lower() for type in allowlist]
    else:
        raise InvalidFilterValue(
            f'Allowlist passed into `filter_type()` filter is not a list ("{allowlist}"); must be like `["Network", "Web"]`'
        )
    try:
        for finding in findings:
            if finding["finding_type"].lower() in allowlist:
                filtered_values.append(finding)
    except (KeyError, TypeError):
        logger.exception("Error parsing ``findings`` as a list of dictionaries: %s", findings)
        raise InvalidFilterValue(
            "Invalid list of findings passed into `filter_type()` filter; must be the `{{ findings }}` object"
        )
    return filtered_values


def strip_html(s):
    """
    Strip HTML tags from the provided HTML while preserving newlines created by
    ``<br />`` and ``<p>`` tags and spaces.

    **Parameters**

    ``s``
        String of HTML text to strip down
    """
    html = BeautifulSoup(s, "lxml")
    output = ""
    for tag in html.descendants:
        if isinstance(tag, str):
            output += tag
        elif tag.name in ("br", "p"):
            output += "\n"
    return output.strip()


def compromised(targets):
    """
    Filter list of targets to return only those marked as compromised.

    **Parameters**

    ``targets``
        List of dictionary objects (JSON) for targets
    """
    filtered_targets = []
    try:
        for target in targets:
            if target["compromised"]:
                filtered_targets.append(target)
    except (KeyError, TypeError):
        logger.exception("Error parsing ``targets`` as a list of dictionaries: %s", targets)
        raise InvalidFilterValue(
            "Invalid list of targets passed into `compromised()` filter; must be the `{{ targets }}` object"
        )
    return filtered_targets


def add_days(date, days):
    """
    Add a number of business days to a date.

    **Parameters**

    ``date``
        Date string to add business days to
    ``days``
        Number of business days to add to the date
    """
    new_date = None
    try:
        days = int(days)
    except ValueError:
        logger.exception("Error parsing ``days`` as an integer: %s", days)
        raise InvalidFilterValue(f'Invalid integer ("{days}") passed into `add_days()` filter')

    try:
        date_obj = parse_datetime(date)
        # Loop until all days added
        if days > 0:
            while days > 0:
                # Add one day to the date
                date_obj += timedelta(days=1)
                # Check if the day is a business day
                weekday = date_obj.weekday()
                if weekday >= 5:
                    # Return to the top (Sunday is 6)
                    continue
                # Decrement the number of days to add
                days -= 1
        else:
            # Same as above but in reverse for negative days
            while days < 0:
                date_obj -= timedelta(days=1)
                weekday = date_obj.weekday()
                if weekday >= 5:
                    continue
                days += 1
        new_date = dateformat(date_obj, settings.DATE_FORMAT)
    except ParserError:
        logger.exception("Error parsing ``date`` as a date: %s", date)
        raise InvalidFilterValue(f'Invalid date string ("{date}") passed into `add_days()` filter')
    return new_date


def format_datetime(date, new_format):
    """
    Change the format of a given date string.

    **Parameters**

    ``date``
        Date string to modify
    ``format_str``
        The format of the provided date
    """
    formatted_date = None
    try:
        date_obj = parse_datetime(date)
        formatted_date = dateformat(date_obj, new_format)
    except ParserError:
        formatted_date = date
        logger.exception("Error parsing ``date`` as a date: %s", date)
        raise InvalidFilterValue(f'Invalid date string ("{date}") passed into `format_datetime()` filter')
    return formatted_date


def sort_findings(findings):
    # Using math to with appropriate weights to make sure mediums severities with low exploit don't pass highs or crits.
    severities = {
        Severity.CRIT.value.lower(): 200,
        Severity.HIGH.value.lower(): 65,
        Severity.MED.value.lower(): 20,
        Severity.LOW.value.lower(): 5,
        Severity.BP.value.lower(): 1,
    }
    diff_of_exploit = {
        Severity.LOW.value.lower(): 3,
        Severity.MED.value.lower(): 2,
        Severity.HIGH.value.lower(): 1,
    }

    for finding in findings:
        weight = severities[finding["severity"].lower()] * diff_of_exploit.get(
            strip_html(finding["host_detection_techniques"]).lower(), 1
        )
        finding["weight"] = weight

    return sorted(findings, key=lambda f: f["weight"], reverse=True)


def prepare_jinja2_env(debug=False):
    """Prepare a Jinja2 environment with all custom filters."""
    if debug:
        undefined = jinja2.DebugUndefined
    else:
        undefined = jinja2.make_logging_undefined(logger=logger, base=jinja2.Undefined)

    env = jinja2.sandbox.SandboxedEnvironment(undefined=undefined, extensions=["jinja2.ext.debug"], autoescape=True)
    env.filters["filter_severity"] = filter_severity
    env.filters["filter_type"] = filter_type
    env.filters["strip_html"] = strip_html
    env.filters["compromised"] = compromised
    env.filters["add_days"] = add_days
    env.filters["format_datetime"] = format_datetime
    env.filters["sort_findings"] = sort_findings

    return env


class ReportConstants:
    """Constant values used for report generation."""

    DEFAULT_STYLE_VALUES = {
        "bold": False,
        "underline": False,
        "italic": False,
        "inline_code": False,
        "strikethrough": False,
        "font_family": None,
        "font_size": None,
        "font_color": None,
        "highlight": None,
        "superscript": False,
        "subscript": False,
        "hyperlink": False,
        "hyperlink_url": None,
    }


class Reportwriter:
    """Generate report documents in Microsoft Office formats and JSON."""

    # Allowlist for HTML tags expected to come from the WYSIWYG
    tag_allowlist = [
        "code",
        "span",
        "p",
        "ul",
        "ol",
        "li",
        "a",
        "em",
        "strong",
        "u",
        "b",
        "pre",
        "sub",
        "sup",
        "del",
        "blockquote",
    ]

    # Allowlist for evidence file extensions / filetypes
    image_extensions = ["png", "jpeg", "jpg"]
    text_extensions = ["txt", "ps1", "py", "md", "log"]

    def __init__(self, report_queryset, template_loc=None):
        self.template_loc = template_loc
        self.report_queryset = report_queryset

        # Get the global report configuration
        global_report_config = ReportConfiguration.get_solo()
        self.company_config = CompanyInformation.get_solo()

        # Track report type for different Office XML
        self.report_type = None

        # Picture border settings for Word
        self.enable_borders = global_report_config.enable_borders
        self.border_color = global_report_config.border_color
        self.border_weight = global_report_config.border_weight

        # Caption options
        prefix_figure = global_report_config.prefix_figure.strip()
        self.prefix_figure = f" {prefix_figure} "
        label_figure = global_report_config.label_figure.strip()
        self.label_figure = f"{label_figure} "
        label_table = global_report_config.label_table.strip()
        self.label_table = f"{label_table} "

        # Set up Jinja2 rendering environment + custom filters
        self.jinja_env = prepare_jinja2_env(debug=False)

        logger.info(
            "Generating a report for %s using the template at %s",
            self.report_queryset,
            self.template_loc,
        )

    def _valid_xml_char_ordinal(self, c):
        """
        Clean string to make all characters XML compatible for Word documents.

        Source:
            https://stackoverflow.com/questions/8733233/filtering-out-certain-bytes-in-python

        **Parameters**

        ``c`` : string
            String of characters to validate
        """
        codepoint = ord(c)
        # Conditions ordered by presumed frequency
        return (
            0x20 <= codepoint <= 0xD7FF
            or codepoint in (0x9, 0xA, 0xD)
            or 0xE000 <= codepoint <= 0xFFFD
            or 0x10000 <= codepoint <= 0x10FFFF
        )

    def generate_json(self):
        """
        Export a report as a JSON dictionary for archiving and to generate other report types.
        """

        # Serialize the :model:`rolodex.Project`
        serializer = ReportDataSerializer(
            self.report_queryset,
            exclude=[
                "id",
            ],
        )
        # Render the serialized data as JSON
        report_json = JSONRenderer().render(serializer.data)
        report_json = json.loads(report_json)
        # An extra step to make the JSON "pretty"
        output = json.dumps(report_json, indent=4)

        return output

    def _make_figure(self, par, ref=None):
        """
        Append a text run configured as an auto-incrementing figure to the provided
        paragraph. The label and number are wrapped in ``w:bookmarkStart`` and
        ``w:bookmarkEnd``.

        Source:
            https://github.com/python-openxml/python-docx/issues/359

        **Parameters**

        ``par`` : docx.paragraph.Paragraph
            Paragraph to alter
        ``ref`` : string
            String to use as the ``w:name`` value for the bookmark
        """

        def generate_ref():
            """Generate a random eight character reference ID."""
            return random.randint(10000000, 99999999)

        if ref:
            ref = f"_Ref{ref}"
        else:
            ref = f"_Ref{generate_ref()}"
        # Start a bookmark run with the figure label
        p = par._p
        bookmark_start = OxmlElement("w:bookmarkStart")
        bookmark_start.set(qn("w:name"), ref)
        bookmark_start.set(qn("w:id"), "0")
        p.append(bookmark_start)

        # Add the figure label
        run = par.add_run(self.label_figure)

        # Append XML for a new field character run
        run = par.add_run()
        r = run._r
        fldChar = OxmlElement("w:fldChar")
        fldChar.set(qn("w:fldCharType"), "begin")
        r.append(fldChar)

        # Add field code instructions with ``instrText``
        run = par.add_run()
        r = run._r
        instrText = OxmlElement("w:instrText")
        # Sequential figure with arabic numbers
        instrText.text = " SEQ Figure \\* ARABIC"
        r.append(instrText)

        # An optional ``separate`` value to enforce a space between label and number
        run = par.add_run()
        r = run._r
        fldChar = OxmlElement("w:fldChar")
        fldChar.set(qn("w:fldCharType"), "separate")
        r.append(fldChar)

        # Include ``#`` as a placeholder for the number when Word updates fields
        run = par.add_run("#")
        r = run._r
        # Close the field character run
        fldChar = OxmlElement("w:fldChar")
        fldChar.set(qn("w:fldCharType"), "end")
        r.append(fldChar)

        # End the bookmark after the number
        p = par._p
        bookmark_end = OxmlElement("w:bookmarkEnd")
        bookmark_end.set(qn("w:id"), "0")
        p.append(bookmark_end)

    def _make_cross_ref(self, par, ref):
        """
        Append a text run configured as a cross-reference to the provided paragraph.

        **Parameters**

        ``par`` : docx.paragraph.Paragraph
            Paragraph to alter
        ``ref`` : string
            The ``w:name`` value of the target bookmark
        """
        # Start the field character run for the label and number
        run = par.add_run()
        r = run._r
        fldChar = OxmlElement("w:fldChar")
        fldChar.set(qn("w:fldCharType"), "begin")
        r.append(fldChar)

        # Add field code instructions with ``instrText`` that points to the target bookmark
        run = par.add_run()
        r = run._r
        instrText = OxmlElement("w:instrText")
        instrText.text = f" REF _Ref{ref} \\h "
        r.append(instrText)

        # An optional ``separate`` value to enforce a space between label and number
        run = par.add_run()
        r = run._r
        fldChar = OxmlElement("w:fldChar")
        fldChar.set(qn("w:fldCharType"), "separate")
        r.append(fldChar)

        # Add runs for the figure label and number
        run = par.add_run(self.label_figure)
        # This ``#`` is a placeholder Word will replace with the figure's number
        run = par.add_run("#")

        # Close the  field character run
        run = par.add_run()
        r = run._r
        fldChar = OxmlElement("w:fldChar")
        fldChar.set(qn("w:fldCharType"), "end")
        r.append(fldChar)

        return par

    def _list_number(self, par, prev=None, level=None, num=True):
        """
        Makes the specified paragraph a list item with a specific level and optional restart.

        An attempt will be made to retrieve an abstract numbering style that corresponds
        to the style of the paragraph. If that is not possible, the default numbering or
        bullet style will be used based on the ``num`` parameter.

        Source:
            https://github.com/python-openxml/python-docx/issues/25#issuecomment-400787031

        **Parameters**

        ``par`` : docx.paragraph.Paragraph
            The docx paragraph to turn into a list item.
        ``prev`` : docx.paragraph.Paragraph or None
            The previous paragraph in the list. If specified, the numbering and styles will
            be taken as a continuation of this paragraph. If omitted, a new numbering scheme
            will be started.
        ``level`` : int or None
            The level of the paragraph within the outline. If ``prev`` is set, defaults
            to the same level as in ``prev``. Otherwise, defaults to zero.
        ``num`` : bool
            If ``prev`` is :py:obj:`None` and the style of the paragraph does not correspond
            to an existing numbering style, this will determine wether or not the list will
            be numbered or bulleted. The result is not guaranteed, but is fairly safe for
            most Word templates.
        """
        xpath_options = {
            True: {"single": "count(w:lvl)=1 and ", "level": 0},
            False: {"single": "", "level": level},
        }

        def style_xpath(prefer_single=True):
            """
            The style comes from the outer-scope variable ``par.style.name``.
            """
            style = par.style.style_id
            return (
                "w:abstractNum[" '{single}w:lvl[@w:ilvl="{level}"]/w:pStyle[@w:val="{style}"]' "]/@w:abstractNumId"
            ).format(style=style, **xpath_options[prefer_single])

        def type_xpath(prefer_single=True):
            """
            The type is from the outer-scope variable ``num``.
            """
            t = "decimal" if num else "bullet"
            return (
                "w:abstractNum[" '{single}w:lvl[@w:ilvl="{level}"]/w:numFmt[@w:val="{type}"]' "]/@w:abstractNumId"
            ).format(type=t, **xpath_options[prefer_single])

        def get_abstract_id():
            """
            Select as follows:

            1. Match single-level by style (get min ID)
            2. Match exact style and level (get min ID)
            3. Match single-level decimal/bullet types (get min ID)
            4. Match decimal/bullet in requested level (get min ID)
            """
            for fn in (style_xpath, type_xpath):
                for prefer_single in (True, False):
                    xpath = fn(prefer_single)
                    ids = numbering.xpath(xpath)
                    if ids:
                        return min(int(x) for x in ids)
            return 0

        if prev is None or prev._p.pPr is None or prev._p.pPr.numPr is None or prev._p.pPr.numPr.numId is None:
            if level is None:
                level = 0
            numbering = self.sacrificial_doc.part.numbering_part.numbering_definitions._numbering
            # Compute the abstract ID first by style, then by ``num``
            abstract = get_abstract_id()
            # Set the concrete numbering based on the abstract numbering ID
            num = numbering.add_num(abstract)
            # Make sure to override the abstract continuation property
            num.add_lvlOverride(ilvl=level).add_startOverride(1)
            # Extract the newly-allocated concrete numbering ID
            num = num.numId
        else:
            if level is None:
                level = prev._p.pPr.numPr.ilvl.val
            # Get the previous concrete numbering ID
            num = prev._p.pPr.numPr.numId.val
        par._p.get_or_add_pPr().get_or_add_numPr().get_or_add_numId().val = num
        par._p.get_or_add_pPr().get_or_add_numPr().get_or_add_ilvl().val = level

        return par

    def _get_styles(self, tag):
        """
        Get styles from an BS4 ``Tag`` object's ``styles`` attribute and convert
        the string to a dictionary.

        **Parameters**

        ``tag`` : Tag
            BS4 ``Tag`` with a ``styles`` attribute
        """
        tag_styles = {}
        style_str = tag.attrs["style"]
        # Filter any blanks from the split
        style_list = list(filter(None, style_str.split(";")))
        for style in style_list:
            temp = style.split(":")
            key = temp[0].strip()
            value = temp[1].strip()
            try:
                if key == "font-size":
                    # Remove the "pt" from the font size and convert it to a ``Pt`` object
                    value = value.replace("pt", "")
                    value = float(value)
                    value = Pt(value)
                if key == "font-family":
                    # Some fonts include a list of values – get just the first one
                    font_list = value.split(",")
                    priority_font = font_list[0].replace("'", "").replace('"', "").strip()
                    value = priority_font
                if key == "color":
                    # Convert the color hex value to an ``RGBColor`` object
                    value = value.replace("#", "")
                    n = 2
                    hex_color = [hex(int(value[i : i + n], 16)) for i in range(0, len(value), n)]
                    if self.report_type == "pptx":
                        value = PptxRGBColor(*map(lambda v: int(v, 16), hex_color))
                    else:
                        value = RGBColor(*map(lambda v: int(v, 16), hex_color))
                tag_styles[key] = value
            except Exception:
                logger.exception("Failed to convert one of the inline styles for a text run")
        return tag_styles

    def _process_evidence(self, evidence, par):
        """
        Process the specified evidence file for the named finding to add it to the Word document.

        **Parameters**

        ``evidence`` : dict
            Evidence to be processed
        ``par`` : Paragraph
            Paragraph meant to hold the evidence
        """
        file_path = settings.MEDIA_ROOT + "/" + evidence["path"]
        extension = file_path.split(".")[-1].lower()

        # First, check if the file still exists on disk
        if os.path.exists(file_path):
            # Next, check if the file is approved and handle as either text or image
            if extension in self.text_extensions:
                with open(file_path, "r") as evidence_contents:
                    # Read in evidence text
                    evidence_text = evidence_contents.read()
                    if self.report_type == "pptx":
                        if par:
                            self._delete_paragraph(par)
                        top = Inches(1.65)
                        left = Inches(8)
                        width = Inches(4.5)
                        height = Inches(3)
                        # Create new textbox, textframe, paragraph, and run
                        textbox = self.finding_slide.shapes.add_textbox(left, top, width, height)
                        text_frame = textbox.text_frame
                        p = text_frame.paragraphs[0]
                        run = p.add_run()
                        # Insert evidence and apply formatting
                        run.text = evidence_text
                        font = run.font
                        font.size = Pt(11)
                        font.name = "Courier New"
                    else:
                        # Drop in text evidence using the ``CodeBlock`` style
                        par.text = evidence_text
                        par.style = "CodeBlock"
                        par.alignment = WD_ALIGN_PARAGRAPH.LEFT
                        # Add a caption paragraph below the evidence
                        p = self.sacrificial_doc.add_paragraph(style="Caption")
                        ref_name = re.sub(
                            "[^A-Za-z0-9]+",
                            "",
                            evidence["friendly_name"],
                        )
                        self._make_figure(p, ref_name)
                        run = p.add_run(self.prefix_figure + evidence["caption"])
            elif extension in self.image_extensions:
                # Drop in the image at the full 6.5" width and add the caption
                if self.report_type == "pptx":
                    if par:
                        self._delete_paragraph(par)
                    # Place new textbox to the mid-right
                    top = Inches(1.65)
                    left = Inches(8)
                    width = Inches(4.5)
                    self.finding_slide.shapes.add_picture(file_path, left, top, width=width)
                else:
                    par.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    run = par.add_run()
                    try:
                        # Add the picture to the document and then add a border
                        run.add_picture(file_path, width=Inches(6.5))
                    except docx.image.exceptions.UnrecognizedImageError:
                        logger.exception(
                            "Evidence file known as %s (%s) was not recognized as a %s file.",
                            evidence["friendly_name"],
                            file_path,
                            extension,
                        )
                        error_msg = (
                            f'The evidence file, `{evidence["friendly_name"]},` was not recognized as a {extension} file. '
                            "Try opening it, exporting as desired type, and re-uploading it."
                        )
                        raise UnrecognizedImageError(error_msg) from docx.image.exceptions.UnrecognizedImageError

                    if self.enable_borders:
                        # Add the border – see Ghostwriter Wiki for documentation
                        inline_class = run._r.xpath("//wp:inline")[-1]
                        inline_class.attrib["distT"] = "0"
                        inline_class.attrib["distB"] = "0"
                        inline_class.attrib["distL"] = "0"
                        inline_class.attrib["distR"] = "0"

                        # Set the shape's "effect extent" attributes to the border weight
                        effect_extent = OxmlElement("wp:effectExtent")
                        effect_extent.set("l", str(self.border_weight))
                        effect_extent.set("t", str(self.border_weight))
                        effect_extent.set("r", str(self.border_weight))
                        effect_extent.set("b", str(self.border_weight))
                        # Insert just below ``<wp:extent>`` or it will not work
                        inline_class.insert(1, effect_extent)

                        # Find inline shape properties – ``pic:spPr``
                        pic_data = run._r.xpath("//pic:spPr")[-1]
                        # Assemble OXML for a solid border
                        ln_xml = OxmlElement("a:ln")
                        ln_xml.set("w", str(self.border_weight))
                        solidfill_xml = OxmlElement("a:solidFill")
                        color_xml = OxmlElement("a:srgbClr")
                        color_xml.set("val", self.border_color)
                        solidfill_xml.append(color_xml)
                        ln_xml.append(solidfill_xml)
                        pic_data.append(ln_xml)

                    # Create the caption for the image
                    p = self.sacrificial_doc.add_paragraph(style="Caption")
                    ref_name = re.sub("[^A-Za-z0-9]+", "", evidence["friendly_name"])
                    self._make_figure(p, ref_name)
                    run = p.add_run(self.prefix_figure + evidence["caption"])
            # Skip unapproved files
            else:
                par = None
        else:
            raise FileNotFoundError(file_path)

    def _delete_paragraph(self, par):
        """
        Delete the specified paragraph.

        **Parameter**

        ``par``
            Paragraph to delete from the document
        """
        p = par._p
        parent_element = p.getparent()
        if parent_element is not None:
            parent_element.remove(p)
        else:
            logger.warning("Could not delete paragraph in because it had no parent element")

    def _write_xml(self, text, par, styles):
        """
        Write the provided text to Office XML.

        **Parameters**

        ``text`` : string
            Text to check for keywords
        ``par`` : Paragraph
            Paragraph for the processed text
        ``styles`` : dict
            Copy of ``ReportConstants.DEFAULT_STYLE_VALUES`` with styles for the text
        """
        # Handle hyperlinks based on Office report type
        # Easy with ``python-pptx`` API, but custom work required for ``python-docx``
        if styles["hyperlink"] and styles["hyperlink_url"]:
            if self.report_type == "pptx":
                run = par.add_run()
                run.text = text
                run.hyperlink.address = styles["hyperlink_url"]
            else:
                # For Word, this code is modified from this issue:
                #   https://github.com/python-openxml/python-docx/issues/384
                # Get an ID from the ``document.xml.rels`` file
                part = par.part
                r_id = part.relate_to(
                    styles["hyperlink_url"],
                    docx.opc.constants.RELATIONSHIP_TYPE.HYPERLINK,
                    is_external=True,
                )
                # Create the ``w:hyperlink`` tag and add needed values
                hyperlink = docx.oxml.shared.OxmlElement("w:hyperlink")
                hyperlink.set(
                    docx.oxml.shared.qn("r:id"),
                    r_id,
                )
                # Create the ``w:r`` and ``w:rPr`` elements
                new_run = docx.oxml.shared.OxmlElement("w:r")
                rPr = docx.oxml.shared.OxmlElement("w:rPr")
                new_run.append(rPr)
                new_run.text = text
                hyperlink.append(new_run)
                # Create a new Run object and add the hyperlink into it
                run = par.add_run()
                run._r.append(hyperlink)
                # A workaround for the lack of a hyperlink style
                if "Hyperlink" in self.sacrificial_doc.styles:
                    run.style = "Hyperlink"
                else:
                    run.font.color.theme_color = MSO_THEME_COLOR_INDEX.HYPERLINK
                    run.font.underline = True
        else:
            run = par.add_run()
            run.text = text

        # Apply font-based styles that work for both APIs
        font = run.font
        font.bold = styles["bold"]
        font.italic = styles["italic"]
        font.underline = styles["underline"]
        if styles["font_color"]:
            font.color.rgb = styles["font_color"]
        font.name = styles["font_family"]
        font.size = styles["font_size"]
        if styles["inline_code"]:
            if self.report_type == "pptx":
                font.name = "Courier New"
            else:
                run.style = "CodeInline"
                font.no_proof = True

        # These styles require extra work due to limitations of the ``python-pptx`` API
        if styles["highlight"]:
            if self.report_type == "pptx":
                rPr = run._r.get_or_add_rPr()
                highlight = OxmlElement("a:highlight")
                srgbClr = OxmlElement("a:srgbClr")
                srgbClr.set("val", "FFFF00")
                highlight.append(srgbClr)
                rPr.append(highlight)
            else:
                font.highlight_color = WD_COLOR_INDEX.YELLOW
        if styles["strikethrough"]:
            if self.report_type == "pptx":
                font._element.set("strike", "sngStrike")
            else:
                font.strike = styles["strikethrough"]
        if styles["subscript"]:
            if self.report_type == "pptx":
                font._element.set("baseline", "-25000")
            else:
                font.subscript = styles["subscript"]
        if styles["superscript"]:
            if self.report_type == "pptx":
                font._element.set("baseline", "30000")
            else:
                font.superscript = styles["superscript"]

    def _add_image(self, par, fig, filename, pad=0.1, image_width=None, image_height=None):
        # Build the filepath to save the figure and add to report
        # Strip special chars except for - and _
        allowed = ascii_letters + "-" + "_"
        new_file_name = "".join(list(filter(allowed.__contains__, filename)))
        directory = f'{settings.MEDIA_ROOT}/evidence/{self.report_json["project"]["id"]}'

        if not os.path.exists(directory):
            # Create a new directory because it does not exist
            os.makedirs(directory)

        filepath = f"{directory}/{new_file_name}.png"
        # Save the figure as a png to the file system under the report directory to be saved into the report
        fig.savefig(filepath, pad_inches=pad, bbox_inches="tight", dpi=fig.get_dpi())

        # Replace figure in report with saved image
        # Use the filename as a label for replacing the text with the image
        run = par.add_run()
        width = Inches(image_width) if image_width else None
        height = Inches(image_height) if image_height else None

        # The image_width and image_height are separate if we want to change the image
        # dimensions but not the figure
        # For example, we only care about setting the figure height to a specific value
        # but don't care about the width of the image
        run.add_picture(filepath, width=width, height=height)
        # Close the current figure window to clear up memory
        plt.close(fig)

    def _replace_and_write(self, text, par, finding, styles=ReportConstants.DEFAULT_STYLE_VALUES.copy()):
        """
        Find and replace template keywords in the provided text.

        **Parameters**

        ``text`` : string
            Text to check for keywords
        ``par`` : Paragraph
            Paragraph for the processed text
        ``styles`` : dict
            Copy of ``ReportConstants.DEFAULT_STYLE_VALUES`` with styles for the text
        """
        # Remove any newlines to avoid creating unwanted blank lines
        text = text.replace("\r\n", "")

        # Perform static text replacements
        # Do this first so strings are not detected as potential expressions–e.g., ``{{.ref ...}}``
        if "{{.client}}" in text:
            if self.report_json["client"]["short_name"]:
                text = text.replace("{{.client}}", self.report_json["client"]["short_name"])
            else:
                text = text.replace("{{.client}}", self.report_json["client"]["name"])
        if "{{.project_start}}" in text:
            text = text.replace("{{.project_start}}", self.report_json["project"]["start_date"])
        if "{{.project_end}}" in text:
            text = text.replace("{{.project_end}}", self.report_json["project"]["end_date"])
        if "{{.project_type}}" in text:
            text = text.replace(
                "{{.project_type}}",
                self.report_json["project"]["project_type"].lower(),
            )

        # Use regex to search for expressions to process
        keyword_regex = r"\{\{\.(.*?)\}\}"
        # Find all strings like ``{{. foo}}``
        match = re.findall(keyword_regex, text)

        # Loop over all regex matches to determine if they are expressions
        cross_refs = []
        if match:
            for var in match:
                # Check for and track cross-references separately for later action
                if var.startswith("ref "):
                    # Track reference with the curly braces restored for later
                    cross_refs.append("{{.%s}}" % var)
                # Process anything that is not a cross-reference now in the match loop
                else:
                    keyword = var.replace("}}", "").replace("{{.", "").strip()

                    # Transform caption placeholders into figures
                    if keyword.startswith("caption"):
                        ref_name = keyword.lstrip("caption ")
                        ref_name = re.sub("[^A-Za-z0-9]+", "", ref_name)
                        text = text.replace("{{.%s}}" % keyword, "")
                        if self.report_type == "pptx":
                            if ref_name:
                                run = par.add_run()
                                run.text = f"See {ref_name}"
                                font = run.font
                                font.italic = True
                        else:
                            par.style = "Caption"
                            if ref_name:
                                self._make_figure(par, ref_name)
                            else:
                                self._make_figure(par)
                            par.add_run(self.prefix_figure + text)
                        # Captions are on their own line so return
                        return par

                    # Bar charts
                    def _build_bar_chart(self, par, keyword, label):
                        par.alignment = WD_ALIGN_PARAGRAPH.LEFT
                        fig = build_bar_chart(self.report_json["totals"][label])
                        # Subtracting - 3.5 from the width to make it fit perfectly on the page with the font sizes
                        # to prevent overlapping on the x-axis and make it more readable
                        self._add_image(par, fig, keyword, image_width=fig.get_figwidth() - 3.5, image_height=fig.get_figheight())

                    if keyword == "chart_bar":
                        _build_bar_chart(self, par, keyword, "chart_data")
                        return par

                    if keyword == "chart_bar_external":
                        _build_bar_chart(self, par, keyword, "chart_data_external")
                        return par

                    if keyword == "chart_bar_internal":
                        _build_bar_chart(self, par, keyword, "chart_data_internal")
                        return par

                    # SD Graphs
                    def _build_sd_graph(self, par, keyword, label):
                        par.alignment = WD_ALIGN_PARAGRAPH.CENTER
                        fig = build_sd_graph(self.report_json["totals"][label])
                        self._add_image(par, fig, keyword, image_height=fig.get_figheight())

                    # Unfortunately had to register new tags until GW has a proper way to build charts
                    # into reports
                    if keyword == "chart_sdscore_appsec":
                        _build_sd_graph(self, par, keyword, "sd_score_appsec")
                        return par

                    if keyword == "chart_sdscore_cloud":
                        _build_sd_graph(self, par, keyword, "sd_score_cloud")
                        return par

                    if keyword == "chart_sdscore_external":
                        _build_sd_graph(self, par, keyword, "sd_score_external")
                        return par

                    if keyword == "chart_sdscore_internal":
                        _build_sd_graph(self, par, keyword, "sd_score_internal")
                        return par

                    if keyword == "chart_sdscore_physical":
                        _build_sd_graph(self, par, keyword, "sd_score_physical")
                        return par

                    if keyword == "chart_sdscore_wireless":
                        _build_sd_graph(self, par, keyword, "sd_score_wireless")
                        return par

                    # Pie charts
                    def _build_pie(self, par, keyword, label):
                        chart_data = self.report_json["totals"][label]
                        total_findings = 0
                        for r in chart_data:
                            total_findings += sum(r[1:])

                        fig = build_pie_chart(chart_data, total_findings)
                        self._add_image(par, fig, keyword, image_height=fig.get_figheight())

                    if keyword == "chart_pie_internal":
                        _build_pie(self, par, keyword, "chart_data_internal")
                        return par

                    if keyword == "chart_pie_external":
                        _build_pie(self, par, keyword, "chart_data_external")
                        return par

                    # Handle evidence files
                    if "evidence" in finding:
                        if (
                            keyword
                            # and keyword in finding["evidence"]
                            and any(ev["friendly_name"] == keyword for ev in finding["evidence"])
                            and not keyword.startswith("ref ")
                        ):
                            logger.debug(
                                "Identified `%s` as an evidence file attached to this finding",
                                keyword,
                            )
                            for ev in finding["evidence"]:
                                if ev["friendly_name"] == keyword:
                                    self._process_evidence(ev, par)
                                    return par
                        else:
                            self._write_xml(text, par, styles)
                    else:
                        self._write_xml(text, par, styles)
        else:
            self._write_xml(text, par, styles)

        # Transform any cross-references into bookmarks
        if cross_refs:
            # Split-up line while keeping cross-references intact
            cross_ref_regex = r"({{.ref.*?}})"
            exploded_text = re.split(cross_ref_regex, text)
            # Loop over the text to replace cross-reference expressions with OXML bookmarks
            for part in exploded_text:
                if part in cross_refs:
                    # Assemble an alphanumeric (no spaces) bookmark name from the tag
                    part = part.replace("}}", "").replace("{{.", "").strip()
                    ref_name = re.sub("[^A-Za-z0-9]+", "", part.lstrip("ref "))
                    if self.report_type == "pptx":
                        run = par.add_run()
                        run.text = f"See {ref_name}"
                        font = run.font
                        font.italic = True
                    else:
                        self._make_cross_ref(
                            par,
                            ref_name,
                        )
                else:
                    self._write_xml(part, par, styles)
            # return par

        return par

    def _process_nested_html_tags(self, contents, par, finding, styles=None):
        """
        Process BeautifulSoup4 ``Tag`` objects containing nested HTML tags.

        **Parameters**

        ``contents`` : Tag.contents
            Contents of a BS4 ``Tag``
        ``par`` : Paragraph
            Word docx ``Paragraph`` object
        ``finding`` : dict
            Current finding (JSON) being processed
        ``styles`` : dict
            Override default styles with a provided dict
        """

        def merge_styles(run_styles, parent_styles):
            """Merge the two style dictionaries, keeping the ``run_styles`` values."""
            for key in parent_styles:
                # Same values in both dicts
                if run_styles[key] == parent_styles[key]:
                    pass
                else:
                    # Run has no value but parent has one
                    if not run_styles[key] and parent_styles[key]:
                        run_styles[key] = parent_styles[key]
            return run_styles

        def check_tags(tag, styles_dict):
            """Check the tag name and update the provided styles dictionary as needed."""
            tag_name = tag.name

            # A ``code`` tag here is inline code inside of a ``p`` tag
            if tag_name == "code":
                styles_dict["inline_code"] = True

            # These tags apply italic, bold, and underline styles but appear rarely
            elif tag_name == "em":
                styles_dict["italic"] = True
            elif tag_name in ("strong", "b"):
                styles_dict["bold"] = True
            elif tag_name == "u":
                styles_dict["underline"] = True

            # The ``sub`` and ``sup`` tags designates subscript and superscript
            elif tag_name == "sub":
                styles_dict["subscript"] = True
            elif tag_name == "sup":
                styles_dict["superscript"] = True

            # The ``del`` tag applies a strikethrough style
            elif tag_name == "del":
                styles_dict["strikethrough"] = True

            # A ``span`` tag will usually contain one or more classes for formatting
            # Empty spans usually only appear in place of a non-breaking spaces
            elif tag_name == "span":
                # Check existence of supported classes for run styles
                if "class" in tag.attrs:
                    tag_attrs = tag.attrs["class"]
                    if "italic" in tag_attrs:
                        styles_dict["italic"] = True
                    if "bold" in tag_attrs:
                        styles_dict["bold"] = True
                    if "underline" in tag_attrs:
                        styles_dict["underline"] = True
                    if "highlight" in tag_attrs:
                        styles_dict["highlight"] = True

                # Check existence of supported character styles
                if "style" in tag.attrs:
                    tag_style = self._get_styles(tag)
                    if "font-size" in tag_style:
                        styles_dict["font_size"] = tag_style["font-size"]
                    if "font-family" in tag_style:
                        styles_dict["font_family"] = tag_style["font-family"]
                    if "color" in tag_style:
                        styles_dict["font_color"] = tag_style["color"]
                    if "background-color" in tag_style:
                        styles_dict["highlight"] = tag_style["background-color"]

            # An ``a`` tag is a hyperlink
            elif tag_name == "a":
                styles_dict["hyperlink"] = True
                styles_dict["hyperlink_url"] = tag["href"]

            # Any other tags are unexpected and ignored
            else:
                if tag_name not in self.tag_allowlist:
                    logger.warning("Encountered an unexpected nested HTML tag: %s", tag_name)

            return styles_dict

        # Begin with HTML that could have indefinitely nested tags
        for part in contents:
            # Track the styles for this first parent tag or use the defaults if none
            if styles:
                parent_styles = styles.copy()
                run_styles = styles.copy()
            else:
                parent_styles = ReportConstants.DEFAULT_STYLE_VALUES.copy()
                run_styles = ReportConstants.DEFAULT_STYLE_VALUES.copy()

            # Get each part's ``name`` to check if it's a ``Tag`` object
            # A plain string will return ``None`` - no HTML tag
            part_name = part.name

            # Get the top-level's styles first as it applies to all future runs
            if part_name:
                # Update styles based on the tag and properties
                parent_styles = check_tags(part, parent_styles)

                # Split part into list of text and ``Tag`` objects
                part_contents = part.contents

                # With parent styles recorded, process each child same as above
                for tag in part_contents:
                    tag_name = tag.name

                    if tag_name and tag_name in self.tag_allowlist:
                        # Construct text to be written as part of this run
                        content_text = ""

                        # Split part into list of text and ``Tag`` objects
                        tag_contents = tag.contents

                        if len(tag_contents) >= 1:
                            # Update styles based on the tag and properties
                            run_styles = check_tags(tag, run_styles)

                            # Check for a hyperlink formatted with additional styles
                            if tag_contents[0].name:
                                if tag_contents[0].name == "a":
                                    run_styles["hyperlink"] = True
                                    run_styles["hyperlink_url"] = tag_contents[0]["href"]
                                    content_text = tag_contents[0].text

                            # Combine the styles to carry them over to the next loop
                            merged_styles = merge_styles(run_styles, parent_styles)

                            # Recursively process the nested tags
                            self._process_nested_html_tags(tag_contents, par, finding, styles=merged_styles)
                    elif tag_name:
                        logger.warning(
                            "Ignoring a nested HTML tag not in the allowlist: %s",
                            tag_name,
                        )
                        continue
                    else:
                        content_text = tag

                    # Combine the parent's and run's styles, favoring the run's
                    merged_styles = merge_styles(run_styles, parent_styles)

                    # Write the text for this run
                    par = self._replace_and_write(content_text, par, finding, merged_styles)

                    # Reset temporary run styles
                    run_styles = ReportConstants.DEFAULT_STYLE_VALUES.copy()

            # There are no tags to process, so write the string
            else:
                if isinstance(part, NavigableString):
                    par = self._replace_and_write(part, par, finding, parent_styles)
                else:
                    par = self._replace_and_write(part.text, par, finding)
        return par

    def _create_list_paragraph(self, prev_p, level, num=False, alignment=WD_ALIGN_PARAGRAPH.LEFT):
        """
        Create a new paragraph in the document for a list.

        **Parameters**

        ``prev_p``
            Previous paragraph to link the next list item
        ``level``
            Indentation level for the list item
        ``num``
            Boolean to determine if the line item will be numbered (Default: False)
        """
        if self.report_type == "pptx":
            # Move to new paragraph/line and indent bullets based on level
            p = self.finding_body_shape.text_frame.add_paragraph()
            p.level = level
        else:
            make_list = True
            styles = self.sacrificial_doc.styles
            try:
                if num:
                    list_style = styles["Number List"]
                else:
                    list_style = styles["Bullet List"]
            except Exception:
                if "List Paragraph" in styles:
                    list_style = styles["List Paragraph"]
                else:
                    list_style = styles["Normal"]
                    make_list = False
            p = self.sacrificial_doc.add_paragraph(style=list_style)
            if make_list:
                p = self._list_number(p, prev=prev_p, level=level, num=num)
            p.alignment = alignment
        return p

    def _parse_nested_html_lists(self, tag, prev_p, num, finding, level=0):
        """
        Recursively parse deeply nested lists. This checks for ``<ol>`` or ``<ul>`` tags
        and keeps parsing until all nested lists are found and processed.

        Returns the last paragraph object created.

        **Parameters**

        ``part``
            BeautifulSoup 4 Tag object to parse
        ``prev_p``
            Previous paragraph to link the next list item
        ``num``
            Boolean to determine if the line item will be numbered
        ``finding``
            Report finding currently being processed
        ``level``
            Indentation level for the list item (Defaults to 0)
        """
        # Check if this element has any ``ul`` or ``ol`` tags anywhere in its contents
        if tag.ol or tag.ul:
            temp = []
            nested_list = None
            li_contents = tag.contents
            # Loop over the contents of the list item to find the nested lists
            for part in tag:
                # Newlines will appear between elements, ignore them
                if part != "\n":
                    # If the tag name is ``ul`` or ``ol``, tack it as a nested list
                    if part.name in ("ol", "ul"):
                        num = bool(part.name == "ol")
                        nested_list = part
                    # Put everything else in a temporary list to be processed later
                    else:
                        temp.append(part)

            # Make the list paragraph here to pick up any changes to the list style (e.g, ``num``)
            p = self._create_list_paragraph(prev_p, level, num)

            # If ``temp`` isn't empty, process it like any other line
            if temp:
                # A length of ``1`` means no nested tags
                if len(temp) == 1:
                    # If the first list item is a ``Tag`` process for styling
                    if temp[0].name:
                        self._process_nested_html_tags(temp, p, finding)
                    # Otherwise, just write the XML
                    else:
                        self._replace_and_write(temp[0], p, finding)
                else:
                    self._process_nested_html_tags(temp, p, finding)

            # If we have nested list(s), recursively process them by re-entering ``_parse_html_lists``
            if nested_list:
                # Increment the indentation level
                if not li_contents[0] == "\n":
                    level += 1
                p = self._parse_html_lists(nested_list, p, num, finding, level)
        # No nested list items, proceed as normal
        # This is where we catch ``li`` tags with nested tags like hyperlinks
        else:
            p = self._create_list_paragraph(prev_p, level, num)
            self._process_nested_html_tags(tag.contents, p, finding)
        return p

    def _parse_html_lists(self, tag, prev_p, num, finding, level=0):
        """
        Recursively parse deeply nested lists. This checks for ``<ol>`` or ``<ul>`` tags
        and keeps parsing until all nested lists are found and processed.

        Returns the last paragraph object created.

        **Parameters**

        ``tag``
            BeautifulSoup 4 Tag object to parse
        ``prev_p``
            Previous paragraph to link the next list item
        ``num``
            Boolean to determine if the line item will be numbered
        ``finding``
            Report finding currently being processed
        ``level``
            Indentation level for the list item (Defaults to 0)
        """
        # Get the individuals contents of the tag
        contents = tag.contents
        # Loop over the contents to find nested lists
        # Nested lists are ``ol`` or ``ul`` tags inside ``li`` tags
        for part in contents:
            # Handle ``li`` tags which might contain more lists
            if part.name == "li":
                li_contents = part.contents
                # A length of ``1`` means there are no nested tags
                if len(li_contents) == 1:
                    p = self._create_list_paragraph(prev_p, level, num)
                    if li_contents[0].name:
                        self._process_nested_html_tags(li_contents, p, finding)
                    else:
                        self._replace_and_write(part.text, p, finding)
                # Bigger lists mean more tags, so process nested tags
                else:
                    # Handle nested lists
                    p = self._parse_nested_html_lists(part, prev_p, num, finding, level)
                # Track the paragraph used for this list item to link subsequent list paragraphs
                prev_p = p
            # If ``ol`` tag encountered, increment ``level`` and switch to numbered list
            elif part.name == "ol":
                level += 1
                p = self._parse_html_lists(part, prev_p, True, finding, level)
            # If ``ul`` tag encountered, increment ``level`` and switch to bulleted list
            elif part.name == "ul":
                level += 1
                p = self._parse_html_lists(part, prev_p, False, finding, level)
            # No change in list type, so proceed with writing the line
            elif part.name:
                p = self._create_list_paragraph(prev_p, level, num)
                self._process_nested_html_tags(part, p, finding)
            # Handle tags that are not handled above
            else:
                if not isinstance(part, NavigableString):
                    logger.warning("Encountered an unknown tag for a list: %s", part.name)
                else:
                    if part.strip() != "":
                        p = self._create_list_paragraph(prev_p, level, num)
                        self._replace_and_write(part.strip(), p, finding)
        # Return last paragraph created
        return p

    def _process_text_xml(self, text, finding=None):
        """
        Process the provided text from the specified finding to parse keywords for
        evidence placement and formatting for Office XML.

        **Parameters**

        ``text``
            Text to convert to Office XML
        ``finding``
            Current report finding being processed
        """
        if text:
            # Clean text to make it XML compatible for Office XML
            text = "".join(c for c in text if self._valid_xml_char_ordinal(c))

            # Parse the HTML into a BS4 soup object
            soup = BeautifulSoup(text, "lxml")

            # Each WYSIWYG field begins with ``<html><body>`` so get the contents of ``body``
            body = soup.find("body")
            contents_list = body.contents

            # Loop over all strings and ``bs4.element.Tag`` objects
            for tag in contents_list:
                # Get the HTML tag's name to determine next steps
                tag_name = tag.name

                # Hn – Headings
                if tag_name in ["h1", "h2", "h3", "h4", "h5", "h6"]:
                    if self.report_type == "pptx":
                        # No headings in PPTX, so add a new line and bold it as a pseudo-heading
                        p = self.finding_body_shape.text_frame.add_paragraph()
                        run = p.add_run()
                        run.text = tag.text
                        font = run.font
                        font.bold = True
                    else:
                        heading_num = int(tag_name[1])
                        # Add the heading to the document
                        # This discards any inline formatting, but that should be managed
                        # by editing the style in the template
                        p = self.sacrificial_doc.add_heading(tag.text, heading_num)

                # P – Paragraphs
                elif tag_name == "p":
                    # Get the tag's contents to check for additional formatting
                    contents = tag.contents

                    # Add a paragraph to the document based on doc type
                    if self.report_type == "pptx":
                        p = self.finding_body_shape.text_frame.add_paragraph()
                        ALIGNMENT = PP_ALIGN
                    else:
                        p = self.sacrificial_doc.add_paragraph()
                        ALIGNMENT = WD_ALIGN_PARAGRAPH

                    # Check for alignment classes on the ``p`` tag
                    if "class" in tag.attrs:
                        tag_attrs = tag.attrs["class"]
                        if "left" in tag_attrs:
                            p.alignment = ALIGNMENT.LEFT
                        if "center" in tag_attrs:
                            p.alignment = ALIGNMENT.CENTER
                        if "right" in tag_attrs:
                            p.alignment = ALIGNMENT.RIGHT
                        if "justify" in tag_attrs:
                            p.alignment = ALIGNMENT.JUSTIFY

                    # Pass the contents and new paragraph on to drill down into nested formatting
                    self._process_nested_html_tags(contents, p, finding)

                # PRE – Code Blocks
                elif tag_name == "pre":
                    # Get the list of pre-formatted strings
                    contents = tag.contents

                    # We do not style any text inside a code block, so we just write the XML
                    # The only content should be one ``code`` block and a long line of text broken up by ``\r\n``
                    if self.report_type == "pptx":
                        # Place new textbox to the mid-right to keep it out of the way
                        if contents:
                            top = Inches(1.65)
                            left = Inches(8)
                            width = Inches(4.5)
                            height = Inches(3)
                            # Create new textbox, textframe, paragraph, and run
                            textbox = self.finding_slide.shapes.add_textbox(left, top, width, height)
                            text_frame = textbox.text_frame
                            for content in contents:
                                for code in content:
                                    parts = code.split("\r\n")
                                    for code_line in parts:
                                        p = text_frame.add_paragraph()
                                        # Align left to anticipate a monospaced font for code
                                        p.alignment = PP_ALIGN.LEFT
                                        run = p.add_run()
                                        # Insert code block and apply formatting
                                        run.text = code_line
                                        font = run.font
                                        font.size = Pt(11)
                                        font.name = "Courier New"
                    else:
                        if contents:
                            for content in contents:
                                for code in content:
                                    parts = code.split("\r\n")
                                    for code_line in parts:
                                        # Create paragraph and apply the ``CodeBlock`` style
                                        p = self.sacrificial_doc.add_paragraph(code_line)
                                        p.style = "CodeBlock"
                                        p.alignment = WD_ALIGN_PARAGRAPH.LEFT

                # OL & UL – Ordered/Numbered & Unordered Lists
                elif tag_name in ("ol", "ul"):
                    # Ordered/numbered lists need numbers and linked paragraphs
                    prev_p = None
                    num = bool(tag_name == "ol")
                    # In HTML, sub-items in a list are nested HTML lists
                    # We need to check every list item for formatted and additional lists
                    # While tracking which level of the list we are working with
                    level = 0
                    prev_p = self._parse_html_lists(tag, prev_p, num, finding, level)

                # BLOCKQUOTE – Blockquote Sections
                elif tag_name == "blockquote":
                    # Get the tag's contents to check for additional formatting
                    contents = tag.contents

                    # PowerPoint lacks a blockquote style, so we just add a basic paragraph
                    if self.report_type == "pptx":
                        p = self.finding_body_shape.text_frame.add_paragraph()
                        ALIGNMENT = PP_ALIGN
                    else:
                        p = self.sacrificial_doc.add_paragraph()
                        p.style = "Blockquote"

                    # Pass the contents and new paragraph on to drill down into nested formatting
                    self._process_nested_html_tags(contents, p, finding)
                else:
                    if not isinstance(tag, NavigableString):
                        logger.warning(
                            "Encountered an unknown tag inside of the finding HTML: %s",
                            tag_name,
                        )

    def generate_word_docx(self):
        """
        Generate a complete Word document for the current report.
        """
        # Generate the JSON for the report
        self.report_json = json.loads(self.generate_json())

        # Create Word document writer using the specified template file
        try:
            self.word_doc = DocxTemplate(self.template_loc)
        except DocxPackageNotFoundError:
            logger.exception(
                "Failed to load the provided template document because file could not be found: %s",
                self.template_loc,
            )
            raise DocxPackageNotFoundError from docx.opc.exceptions.PackageNotFoundError
        except Exception:
            logger.exception("Failed to load the provided template document: %s", self.template_loc)

        # Check for styles
        styles = self.word_doc.styles
        if "CodeBlock" not in styles:
            codeblock_style = styles.add_style("CodeBlock", WD_STYLE_TYPE.PARAGRAPH)
            codeblock_style.base_style = styles["Normal"]
            codeblock_style.hidden = False
            codeblock_style.quick_style = True
            codeblock_style.priority = 2
            # Set font and size
            codeblock_font = codeblock_style.font
            codeblock_font.name = "Courier New"
            codeblock_font.size = Pt(11)
            # Set alignment
            codeblock_par = codeblock_style.paragraph_format
            codeblock_par.alignment = WD_ALIGN_PARAGRAPH.LEFT
            codeblock_par.line_spacing = 1
            codeblock_par.left_indent = Inches(0.2)
            codeblock_par.right_indent = Inches(0.2)

        if "CodeInline" not in styles:
            codeinline_style = styles.add_style("CodeInline", WD_STYLE_TYPE.CHARACTER)
            codeinline_style.hidden = False
            codeinline_style.quick_style = True
            codeinline_style.priority = 3
            # Set font and size
            codeinline_font = codeinline_style.font
            codeinline_font.name = "Courier New"
            codeinline_font.size = Pt(11)

        if "Caption" not in styles:
            caption_style = styles.add_style("Caption", WD_STYLE_TYPE.PARAGRAPH)
            caption_style.hidden = False
            caption_style.quick_style = True
            caption_style.priority = 4
            # Set font and size
            caption_font = caption_style.font
            caption_font.name = "Calibri"
            caption_font.size = Pt(9)
            caption_font.italic = True

        if "Blockquote" not in styles:
            block_style = styles.add_style("Blockquote", WD_STYLE_TYPE.PARAGRAPH)
            block_style.base_style = styles["Normal"]
            block_style.hidden = False
            block_style.quick_style = True
            block_style.priority = 5
            # Set font and size
            block_font = block_style.font
            block_font.name = "Calibri"
            block_font.size = Pt(12)
            block_font.italic = True
            # Set alignment
            block_par = block_style.paragraph_format
            block_par.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
            block_par.left_indent = Inches(0.2)
            block_par.right_indent = Inches(0.2)
            # Keep first and last lines together after repagination
            block_par.widow_control = True

        # Process template context, converting HTML elements to XML as needed
        context = self._process_richtext(self.report_json)

        # Render the Word document + auto-escape any unsafe XML/HTML
        self.word_doc.render(context, self.jinja_env, autoescape=True)

        # Return the final rendered document
        return self.word_doc

    def _process_richtext(self, context: dict) -> dict:
        """
        Update the document context with ``RichText`` and ``Subdocument`` objects for
        each finding and any other values editable with a WYSIWYG editor.

        **Parameters**

        ``context``
            Pre-defined template context
        """

        def render_subdocument(section, finding):
            if section:
                self.sacrificial_doc = self.word_doc.new_subdoc()
                self._process_text_xml(section, finding)
                return self.sacrificial_doc
            return None

        # Findings
        for finding in context["findings"]:
            logger.info("Processing %s", finding["title"])
            # Create ``RichText()`` object for a colored severity category
            finding["severity_rt"] = RichText(finding["severity"], color=finding["severity_color"])
            finding["cvss_score_rt"] = RichText(finding["cvss_score"], color=finding["severity_color"])
            finding["cvss_vector_rt"] = RichText(finding["cvss_vector"], color=finding["severity_color"])
            # Create subdocuments for each finding section
            finding["affected_entities_rt"] = render_subdocument(finding["affected_entities"], finding)
            finding["description_rt"] = render_subdocument(finding["description"], finding)
            finding["impact_rt"] = render_subdocument(finding["impact"], finding)

            # Include a copy of ``mitigation`` as ``recommendation`` to match legacy context
            mitigation_section = render_subdocument(finding["mitigation"], finding)
            finding["mitigation_rt"] = mitigation_section
            finding["recommendation_rt"] = mitigation_section

            finding["replication_steps_rt"] = render_subdocument(finding["replication_steps"], finding)
            finding["host_detection_techniques_rt"] = RichText(
                strip_html(finding["host_detection_techniques"]),
                color=get_value_from_key(DifficultyExploitColor, strip_html(finding["host_detection_techniques"]),)
            )
            finding["network_detection_techniques_rt"] = RichText(
                strip_html(finding["network_detection_techniques"]),
                color=get_value_from_key(FindingStatusColor, strip_html(finding["network_detection_techniques"]),)
            )
            finding["references_rt"] = render_subdocument(finding["references"], finding)

        # Client Notes
        context["client"]["note_rt"] = render_subdocument(context["client"]["note"], finding=None)
        context["client"]["address_rt"] = render_subdocument(context["client"]["address"], finding=None)

        # Project Notes
        context["project"]["note_rt"] = render_subdocument(context["project"]["note"], finding=None)

        # Bar Charts
        context["project"]["chart_bar"] = "<p>{{.chart_bar}}</p>"
        context["project"]["chart_bar_rt"] = render_subdocument(
            context["project"]["chart_bar"], finding=None
        )
        context["project"]["chart_bar_external"] = "<p>{{.chart_bar_external}}</p>"
        context["project"]["chart_bar_external_rt"] = render_subdocument(
            context["project"]["chart_bar_external"], finding=None
        )
        context["project"]["chart_bar_internal"] = "<p>{{.chart_bar_internal}}</p>"
        context["project"]["chart_bar_internal_rt"] = render_subdocument(
            context["project"]["chart_bar_internal"], finding=None
        )

        # SD Graphs
        context["project"]["chart_sdscore_appsec"] = "<p>{{.chart_sdscore_appsec}}</p>"
        context["project"]["chart_sdscore_appsec_rt"] = render_subdocument(
            context["project"]["chart_sdscore_appsec"], finding=None
        )
        context["project"]["chart_sdscore_cloud"] = "<p>{{.chart_sdscore_cloud}}</p>"
        context["project"]["chart_sdscore_cloud_rt"] = render_subdocument(
            context["project"]["chart_sdscore_cloud"], finding=None
        )
        context["project"]["chart_sdscore_wireless"] = "<p>{{.chart_sdscore_wireless}}</p>"
        context["project"]["chart_sdscore_wireless_rt"] = render_subdocument(
            context["project"]["chart_sdscore_wireless"], finding=None
        )
        context["project"]["chart_sdscore_external"] = "<p>{{.chart_sdscore_external}}</p>"
        context["project"]["chart_sdscore_external_rt"] = render_subdocument(
            context["project"]["chart_sdscore_external"], finding=None
        )
        context["project"]["chart_sdscore_internal"] = "<p>{{.chart_sdscore_internal}}</p>"
        context["project"]["chart_sdscore_internal_rt"] = render_subdocument(
            context["project"]["chart_sdscore_internal"], finding=None
        )
        context["project"]["chart_sdscore_physical"] = "<p>{{.chart_sdscore_physical}}</p>"
        context["project"]["chart_sdscore_physical_rt"] = render_subdocument(
            context["project"]["chart_sdscore_physical"], finding=None
        )

        # Pie Charts
        context["project"]["chart_pie_external"] = "<p>{{.chart_pie_external}}</p>"
        context["project"]["chart_pie_external_rt"] = render_subdocument(
            context["project"]["chart_pie_external"], finding=None
        )
        context["project"]["chart_pie_internal"] = "<p>{{.chart_pie_internal}}</p>"
        context["project"]["chart_pie_internal_rt"] = render_subdocument(
            context["project"]["chart_pie_internal"], finding=None
        )

        # Assignments
        for assignment in context["team"]:
            if isinstance(assignment, dict):
                if assignment["note"]:
                    assignment["note_rt"] = render_subdocument(assignment["note"], finding=None)

        # Contacts
        for contact in context["client"]["contacts"]:
            if isinstance(contact, dict):
                if contact["note"]:
                    contact["note_rt"] = render_subdocument(contact["note"], finding=None)

        # Objectives
        for objective in context["objectives"]:
            if isinstance(objective, dict):
                if objective["description"]:
                    objective["description_rt"] = render_subdocument(objective["description"], finding=None)

        # Scope Lists
        for scope_list in context["scope"]:
            if isinstance(scope_list, dict):
                if scope_list["description"]:
                    scope_list["description_rt"] = render_subdocument(scope_list["description"], finding=None)

        # Targets
        for target in context["targets"]:
            if isinstance(target, dict):
                if target["note"]:
                    target["note_rt"] = render_subdocument(target["note"], finding=None)

        # Deconfliction Events
        for event in context["deconflictions"]:
            if isinstance(event, dict):
                if event["description"]:
                    event["description_rt"] = render_subdocument(event["description"], finding=None)

        # White Cards
        for card in context["whitecards"]:
            if isinstance(card, dict):
                if card["description"]:
                    card["description_rt"] = render_subdocument(card["description"], finding=None)

        # Infrastructure
        for asset_type in context["infrastructure"]:
            for asset in context["infrastructure"][asset_type]:
                if isinstance(asset, dict):
                    if asset["note"]:
                        asset["note_rt"] = render_subdocument(asset["note"], finding=None)

        return context

    def _process_text_xlsx(self, html, text_format, finding):
        """
        Process the provided text from the specified finding to parse keywords for
        evidence placement and formatting in xlsx documents.

        **Parameters**

        ``html``
            HTML content to parse with BeautifulSoup 4
        ``text_format``
            Format to apply to the Excel worksheet cell – defined in ``generate_excel_xlsx()``
        ``finding``
            Current report finding being processed
        """
        # Regex for searching for bracketed template placeholders, e.g. {{.client}}
        keyword_regex = r"\{\{\.(.*?)\}\}"

        # Strip out all HTML tags because we can't format text runs for XLSX
        if html:
            text = BeautifulSoup(html, "lxml").text
        else:
            text = ""

        # Perform the necessary replacements
        if "{{.client}}" in text:
            if self.report_json["client"]["short_name"]:
                text = text.replace("{{.client}}", self.report_json["client"]["short_name"])
            else:
                text = text.replace("{{.client}}", self.report_json["client"]["name"])
        if "{{.project_start}}" in text:
            text = text.replace("{{.project_start}}", self.report_json["project"]["start_date"])
        if "{{.project_end}}" in text:
            text = text.replace("{{.project_end}}", self.report_json["project"]["end_date"])
        if "{{.project_type}}" in text:
            text = text.replace("{{.project_type}}", self.report_json["project"]["project_type"].lower())

        # No evidence or captions in workbook cells
        text = text.replace("{{.caption}}", "Caption \u2013 ")

        # Find/replace evidence keywords to make everything human readable
        match = re.findall(keyword_regex, text)
        if match:
            for keyword in match:
                if "evidence" in finding:
                    if (
                        keyword
                        # and keyword in finding["evidence"]
                        and any(ev["friendly_name"] == keyword for ev in finding["evidence"])
                        and not keyword.startswith("ref ")
                    ):
                        for ev in finding["evidence"]:
                            if ev["friendly_name"] == keyword:
                                text = text.replace(
                                    "{{." + keyword + "}}",
                                    "\n<See Report for Evidence File: {}>\nCaption \u2013 {}".format(
                                        ev["friendly_name"],
                                        ev["caption"],
                                    ),
                                )
                # Ignore any other non-keyword string that happens to be inside braces
                else:
                    pass

        self.worksheet.write(self.row, self.col, text, text_format)

    def generate_excel_xlsx(self, memory_object):
        """
        Generate a complete Excel spreadsheet for the current report.
        """

        # Generate the JSON for the report
        self.report_json = json.loads(self.generate_json())

        # Create xlsxwriter in memory with a named worksheet
        xlsx_doc = memory_object
        self.worksheet = xlsx_doc.add_worksheet("Findings")

        # Create a format for headers
        bold_format = xlsx_doc.add_format({"bold": True})
        bold_format.set_text_wrap()
        bold_format.set_align("vcenter")

        # Create a format for affected entities
        asset_format = xlsx_doc.add_format()
        asset_format.set_text_wrap()
        asset_format.set_align("vcenter")
        asset_format.set_align("center")

        # Create a format for everything else
        wrap_format = xlsx_doc.add_format()
        wrap_format.set_text_wrap()
        wrap_format.set_align("vcenter")

        # Create header row for findings
        self.col = 0
        headers = [
            "Finding",
            "Severity",
            "Affected Entities",
            "Description",
            "Impact",
            "Recommendation",
            "Replication Steps",
            "Host Detection Techniques",
            "Network Detection Techniques",
            "References",
            "Supporting Evidence",
        ]

        # Create 30 width columns and then shrink severity to 10
        for header in headers:
            self.worksheet.write(0, self.col, header, bold_format)
            self.col += 1
        self.worksheet.set_column(0, 10, 30)
        self.worksheet.set_column(1, 1, 10)

        # Loop through the findings to create the rest of the worksheet
        self.col = 0
        self.row = 1
        for finding in self.report_json["findings"]:
            # Finding Name
            self.worksheet.write(self.row, self.col, finding["title"], wrap_format)
            self.col += 1

            # Severity
            severity_format = xlsx_doc.add_format({"bold": True})
            severity_format.set_align("vcenter")
            severity_format.set_align("center")
            severity_format.set_font_color("black")

            # Color the cell based on corresponding severity color
            severity_format.set_bg_color(finding["severity_color"])
            self.worksheet.write(self.row, 1, finding["severity"], severity_format)
            self.col += 1

            # Affected Entities
            if finding["affected_entities"]:
                self._process_text_xlsx(finding["affected_entities"], asset_format, finding)
            else:
                self.worksheet.write(self.row, self.col, "N/A", asset_format)
            self.col += 1

            # Description
            self._process_text_xlsx(finding["description"], wrap_format, finding)
            self.col += 1

            # Impact
            self._process_text_xlsx(finding["impact"], wrap_format, finding)
            self.col += 1

            # Recommendation
            self._process_text_xlsx(finding["recommendation"], wrap_format, finding)
            self.col += 1

            # Replication
            self._process_text_xlsx(finding["replication_steps"], wrap_format, finding)
            self.col += 1

            # Detection
            self._process_text_xlsx(finding["host_detection_techniques"], wrap_format, finding)
            self.col += 1
            self._process_text_xlsx(finding["network_detection_techniques"], wrap_format, finding)
            self.col += 1

            # References
            self._process_text_xlsx(finding["references"], wrap_format, finding)
            self.col += 1

            # Collect the evidence, if any, from the finding's folder and insert inline with description
            try:
                evidence_queryset = Evidence.objects.filter(finding=finding["id"])
            except Evidence.DoesNotExist:
                evidence_queryset = []
            except Exception:
                logger.exception("Query for evidence failed for finding %s", finding["id"])
                evidence_queryset = []
            evidence = [
                f.document.name for f in evidence_queryset if f in self.image_extensions or self.text_extensions
            ]
            finding_evidence_names = "\r\n".join(map(str, evidence))
            self.worksheet.write(self.row, self.col, finding_evidence_names, wrap_format)

            # Increment row counter and reset columns before moving on to next finding
            self.row += 1
            self.col = 0

        # Add a filter to the worksheet
        self.worksheet.autofilter("A1:J{}".format(len(self.report_json["findings"]) + 1))

        # Finalize document
        xlsx_doc.close()
        return xlsx_doc

    def generate_powerpoint_pptx(self):
        """
        Generate a complete PowerPoint slide deck for the current report.
        """
        self.report_type = "pptx"
        self.report_json = json.loads(self.generate_json())

        # Create document writer using the specified template
        try:
            self.ppt_presentation = Presentation(self.template_loc)
        except ValueError:
            logger.exception(
                "Failed to load the provided template document because it is not a PowerPoint file: %s",
                self.template_loc,
            )
            raise ValueError
        except PptxPackageNotFoundError:
            logger.exception(
                "Failed to load the provided template document because file could not be found: %s",
                self.template_loc,
            )
            raise PptxPackageNotFoundError from pptx.exc.PackageNotFoundError
        except Exception:
            logger.exception(
                "Failed to load the provided template document for unknown reason: %s",
                self.template_loc,
            )

        # Loop through the findings to create slides
        findings_stats = {}

        def get_textframe(shape):
            """
            Get the shape's text frame and enable automatic resizing. The resize only
            triggers after opening the file in the PowerPoint application and making a change or saving.
            """
            text_frame = shape.text_frame
            text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            return text_frame

        # Calculate finding stats
        for finding in self.report_json["findings"]:
            findings_stats[finding["severity"]] = 0

        for finding in self.report_json["findings"]:
            findings_stats[finding["severity"]] += 1

        # Slide styles (From Master Style counting top to bottom from 0..n)
        SLD_LAYOUT_TITLE = 0
        SLD_LAYOUT_TITLE_AND_CONTENT = 1
        SLD_LAYOUT_FINAL = 12

        # Add a title slide
        slide_layout = self.ppt_presentation.slide_layouts[SLD_LAYOUT_TITLE]
        slide = self.ppt_presentation.slides.add_slide(slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        title_shape.text = self.company_config.company_name
        text_frame = get_textframe(body_shape)
        # Use ``text_frame.text`` for first line/paragraph or ``text_frame.paragraphs[0]``
        text_frame.text = f'{self.report_json["project"]["type"]} Debrief'
        p = text_frame.add_paragraph()
        p.text = self.report_json["client"]["name"]

        # Add Agenda slide
        slide_layout = self.ppt_presentation.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
        slide = self.ppt_presentation.slides.add_slide(slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.text = "Agenda"
        body_shape = shapes.placeholders[1]
        text_frame = get_textframe(body_shape)

        # Add Introduction slide
        slide_layout = self.ppt_presentation.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
        slide = self.ppt_presentation.slides.add_slide(slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.text = "Introduction"
        body_shape = shapes.placeholders[1]
        text_frame = get_textframe(body_shape)

        # Add Methodology slide
        slide_layout = self.ppt_presentation.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
        slide = self.ppt_presentation.slides.add_slide(slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.text = "Methodology"
        body_shape = shapes.placeholders[1]
        text_frame = get_textframe(body_shape)

        # Add Attack Path Overview slide
        slide_layout = self.ppt_presentation.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
        slide = self.ppt_presentation.slides.add_slide(slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.text = "Attack Path Overview"
        body_shape = shapes.placeholders[1]
        text_frame = get_textframe(body_shape)

        # Add Findings Overview Slide
        slide_layout = self.ppt_presentation.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
        slide = self.ppt_presentation.slides.add_slide(slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        title_shape.text = "Findings Overview"
        text_frame = get_textframe(body_shape)
        for key, value in findings_stats.items():
            p = text_frame.add_paragraph()
            p.text = "{} Findings".format(key)
            p.level = 0
            p = text_frame.add_paragraph()
            p.text = str(value)
            p.level = 1

        # Add Findings Overview Slide 2
        # If there are findings then write a table of findings and severity ratings
        if len(self.report_json["findings"]) > 0:
            # Delete the default text placeholder
            textbox = shapes[1]
            sp = textbox.element
            sp.getparent().remove(sp)
            # Add a table
            rows = len(self.report_json["findings"]) + 1
            columns = 2
            left = Inches(1.5)
            top = Inches(2)
            width = Inches(8)
            height = Inches(0.8)
            table = shapes.add_table(rows, columns, left, top, width, height).table
            # Set column width
            table.columns[0].width = Inches(8.5)
            table.columns[1].width = Inches(2.0)
            # Write table headers
            cell = table.cell(0, 0)
            cell.text = "Finding"
            cell.fill.solid()
            cell.fill.fore_color.rgb = pptx.dml.color.RGBColor(0x2D, 0x28, 0x69)
            cell = table.cell(0, 1)
            cell.text = "Severity"
            cell.fill.solid()
            cell.fill.fore_color.rgb = pptx.dml.color.RGBColor(0x2D, 0x28, 0x69)
            # Write findings rows
            row_iter = 1
            for finding in self.report_json["findings"]:
                table.cell(row_iter, 0).text = finding["title"]
                risk_cell = table.cell(row_iter, 1)
                # Set risk rating
                risk_cell.text = finding["severity"]
                # Set cell color fill type to solid
                risk_cell.fill.solid()
                # Color the risk cell based on corresponding severity color
                cell_color = pptx.dml.color.RGBColor(*map(lambda v: int(v, 16), finding["severity_color_hex"]))
                risk_cell.fill.fore_color.rgb = cell_color
                row_iter += 1
            # Set all cells alignment to center and vertical center
            for cell in table.iter_cells():
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        else:
            p = text_frame.add_paragraph()
            p.text = "No findings"
            p.level = 0

        # Create slide for each finding
        for finding in self.report_json["findings"]:
            slide_layout = self.ppt_presentation.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
            self.finding_slide = self.ppt_presentation.slides.add_slide(slide_layout)
            shapes = self.finding_slide.shapes
            title_shape = shapes.title

            # Prepare text frame
            self.finding_body_shape = shapes.placeholders[1]
            if self.finding_body_shape.has_text_frame:
                text_frame = get_textframe(self.finding_body_shape)
                text_frame.clear()
                self._delete_paragraph(text_frame.paragraphs[0])

            # Set slide title to title + [severity]
            title_shape.text = f'{finding["title"]} [{finding["severity"]}]'

            # Add description to the slide body (other sections will appear in the notes)
            if finding["description"]:
                self._process_text_xml(finding["description"], finding)
            else:
                self._process_text_xml("<p>No description provided</p>", finding)

            if "evidence" in finding:
                for ev in finding["evidence"]:
                    self._process_evidence(ev, par=None)

            def prepare_for_pptx(value):
                """Strip HTML and clear 0x0D characters to prepare text for notes slides."""
                try:
                    if value:
                        return BeautifulSoup(value, "lxml").text.replace("\x0D", "")
                    return "N/A"
                except Exception:
                    logger.exception("Failed parsing this value for PPTX: %s", value)
                    return ""

            # Add all finding data to the notes section for easier reference during edits
            entities = prepare_for_pptx(finding["affected_entities"])
            impact = prepare_for_pptx(finding["impact"])
            host_detection = prepare_for_pptx(finding["host_detection_techniques"])
            net_detection = prepare_for_pptx(finding["network_detection_techniques"])
            recommendation = prepare_for_pptx(finding["recommendation"])
            replication = prepare_for_pptx(finding["replication_steps"])
            references = prepare_for_pptx(finding["references"])
            notes_slide = self.finding_slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            p = text_frame.add_paragraph()
            p.text = f"""
                {finding["severity"].capitalize()}: finding["title"]

                AFFECTED ENTITIES
                {entities}

                IMPACT
                {impact}

                MITIGATION
                {recommendation}

                REPLICATION
                {replication}

                HOST DETECTION
                {host_detection}

                NETWORK DETECTION
                ,
                {net_detection}

                REFERENCES
                {references}
            """.replace(
                "                ", ""
            )

        # Add Observations slide
        slide_layout = self.ppt_presentation.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
        slide = self.ppt_presentation.slides.add_slide(slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        title_shape.text = "Positive Observations"
        text_frame = get_textframe(body_shape)

        # Add Recommendations slide
        slide_layout = self.ppt_presentation.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
        slide = self.ppt_presentation.slides.add_slide(slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        title_shape.text = "Recommendations"
        text_frame = get_textframe(body_shape)

        # Add Conclusion slide
        slide_layout = self.ppt_presentation.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
        slide = self.ppt_presentation.slides.add_slide(slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        title_shape.text = "Conclusion"
        text_frame = get_textframe(body_shape)

        # Add final slide
        slide_layout = self.ppt_presentation.slide_layouts[SLD_LAYOUT_FINAL]
        slide = self.ppt_presentation.slides.add_slide(slide_layout)
        shapes = slide.shapes
        body_shape = shapes.placeholders[1]
        text_frame = get_textframe(body_shape)
        text_frame.clear()
        p = text_frame.paragraphs[0]
        p.line_spacing = 0.7
        p.text = self.company_config.company_name
        p = text_frame.add_paragraph()
        p.text = self.company_config.company_twitter
        p.line_spacing = 0.7
        p = text_frame.add_paragraph()
        p.text = self.company_config.company_email
        p.line_spacing = 0.7

        # Finalize document and return it for an HTTP response
        return self.ppt_presentation

    def generate_all_reports(self, docx_template, pptx_template):
        """
        Generate all available report types and return memory streams for each file.
        """
        # Generate the JSON report - it just needs to be a string object
        self.report_json = json.loads(self.generate_json())
        # Generate the docx report - save it in a memory stream
        word_stream = io.BytesIO()
        self.template_loc = docx_template
        word_doc = self.generate_word_docx()
        word_doc.save(word_stream)
        # Generate the xlsx report - save it in a memory stream
        excel_stream = io.BytesIO()
        workbook = Workbook(excel_stream, {"in_memory": True})
        self.generate_excel_xlsx(workbook)
        # Generate the pptx report - save it in a memory stream
        ppt_stream = io.BytesIO()
        self.template_loc = pptx_template
        ppt_doc = self.generate_powerpoint_pptx()
        ppt_doc.save(ppt_stream)
        # Return each memory object
        return self.report_json, word_stream, excel_stream, ppt_stream


class TemplateLinter:
    """Lint template files to catch undefined variables and syntax errors."""

    def __init__(self, template_loc):
        self.template_loc = template_loc
        self.jinja_template_env = prepare_jinja2_env(debug=True)

    def lint_docx(self):
        """
        Lint the provided Word docx file from :model:`reporting.ReportTemplate`.
        """
        results = {"result": "success", "warnings": [], "errors": []}
        if self.template_loc:
            if os.path.exists(self.template_loc):
                logger.info("Found template file at %s", self.template_loc)
                try:
                    # Step 1: Load the document as a template
                    template_document = DocxTemplate(self.template_loc)
                    logger.info("Template loaded for linting")

                    # Step 2: Check document's styles
                    document_styles = template_document.styles
                    if "Bullet List" not in document_styles:
                        results["warnings"].append(
                            "Template is missing a recommended style (see documentation): Bullet List"
                        )
                    if "Number List" not in document_styles:
                        results["warnings"].append(
                            "Template is missing a recommended style (see documentation): Number List"
                        )
                    if "CodeBlock" not in document_styles:
                        results["warnings"].append(
                            "Template is missing a recommended style (see documentation): CodeBlock"
                        )
                    if "CodeInline" not in document_styles:
                        results["warnings"].append(
                            "Template is missing a recommended style (see documentation): CodeInline"
                        )
                    if "Caption" not in document_styles:
                        results["warnings"].append(
                            "Template is missing a recommended style (see documentation): Caption"
                        )
                    if "List Paragraph" not in document_styles:
                        results["warnings"].append(
                            "Template is missing a recommended style (see documentation): List Paragraph"
                        )
                    if "Blockquote" not in document_styles:
                        results["warnings"].append(
                            "Template is missing a recommended style (see documentation): Blockquote"
                        )
                    logger.info("Completed Word style checks")

                    # Step 3: Test rendering the document
                    try:
                        template_document.render(LINTER_CONTEXT, self.jinja_template_env, autoescape=True)
                        undefined_vars = template_document.undeclared_template_variables
                        if undefined_vars:
                            for variable in undefined_vars:
                                results["warnings"].append(f"Undefined variable: {variable}")
                        if results["warnings"]:
                            results["result"] = "warning"
                        logger.info("Completed document rendering test")
                    except TemplateSyntaxError as error:
                        logger.exception("Template syntax error: %s", error)
                        results = {
                            "result": "failed",
                            "errors": [f"Jinja2 template syntax error: {error.message}"],
                        }
                    except UndefinedError as error:
                        logger.error("Template syntax error: %s", error)
                        results = {
                            "result": "failed",
                            "errors": [f"Jinja2 template syntax error: {error.message}"],
                        }
                    except InvalidFilterValue as error:
                        logger.error("Invalid value provided to filter: %s", error)
                        results = {
                            "result": "failed",
                            "errors": [f"Invalid filter value: {error.message}"],
                        }
                    except TypeError as error:
                        logger.error("Invalid value provided to filter or expression: %s", error)
                        results = {
                            "result": "failed",
                            "errors": [f"Invalid value provided to filter or expression: {error}"],
                        }
                    except TemplateRuntimeError as error:
                        logger.error("Invalid filter or expression: %s", error)
                        results = {
                            "result": "failed",
                            "errors": [f"Invalid filter or expression: {error}"],
                        }
                except Exception:
                    logger.exception("Template failed rendering")
                    results = {
                        "result": "failed",
                        "errors": ["Template rendering failed unexpectedly"],
                    }
            else:
                logger.error("Template file path did not exist: %s", self.template_loc)
                results = {
                    "result": "failed",
                    "errors": ["Template file does not exist – upload it again"],
                }
        else:
            logger.error("Received a `None` value for template location")

        logger.info("Template linting completed")
        return results

    def lint_pptx(self):
        """
        Lint the provided PowerPoint pptx file from :model:`reporting.ReportTemplate`.
        """
        results = {"result": "success", "warnings": [], "errors": []}
        if self.template_loc:
            if os.path.exists(self.template_loc):
                logger.info("Found template file at %s", self.template_loc)
                try:
                    # Test 1: Check if the document is a PPTX file
                    template_document = Presentation(self.template_loc)

                    # Test 2: Check for existing slides
                    slide_count = len(template_document.slides)
                    logger.info("Slide count was %s", slide_count)
                    if slide_count > 0:
                        results["warnings"].append(
                            "Template can be used, but it has slides when it should be empty (see documentation)"
                        )
                except ValueError:
                    logger.exception(
                        "Failed to load the provided template document because it is not a PowerPoint file: %s",
                        self.template_loc,
                    )
                    results = {
                        "result": "failed",
                        "errors": ["Template file is not a PowerPoint presentation"],
                    }
                except TypeError as error:
                    logger.error("Invalid value provided to filter or expression: %s", error)
                    results = {
                        "result": "failed",
                        "errors": [f"Invalid value provided to filter or expression: {error}"],
                    }
                except TemplateRuntimeError as error:
                    logger.error("Invalid filter or expression: %s", error)
                    results = {
                        "result": "failed",
                        "errors": [f"Invalid filter or expression: {error}"],
                    }
                except Exception:
                    logger.exception("Template failed rendering")
                    results = {
                        "result": "failed",
                        "errors": ["Template rendering failed unexpectedly"],
                    }
            else:
                logger.error("Template file path did not exist: %s", self.template_loc)
                results = {
                    "result": "failed",
                    "errors": ["Template file does not exist – upload it again"],
                }
        else:
            logger.error("Received a `None` value for template location")

        logger.info("Template linting completed")
        return results
