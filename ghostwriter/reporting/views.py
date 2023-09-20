"""This contains all the views used by the Reporting application."""

# Standard Libraries
import io
import json
import logging.config
import os
import re
import zipfile
from asgiref.sync import async_to_sync
from datetime import datetime
from os.path import exists
from socket import gaierror

# Django Imports
from django.conf import settings
from django.contrib import messages
from django.contrib.auth import get_user_model
from django.contrib.auth.decorators import login_required
from django.core.files import File
from django.core.files.base import ContentFile
from django.db.models import Q
from django.http import (
    FileResponse,
    Http404,
    HttpResponse,
    HttpResponseRedirect,
    JsonResponse,
)
from django.shortcuts import get_object_or_404, redirect, render
from django.template.loader import render_to_string
from django.urls import reverse, reverse_lazy
from django.utils import dateformat, timezone
from django.views import generic
from django.views.generic.detail import DetailView, SingleObjectMixin
from django.views.generic.edit import CreateView, DeleteView, UpdateView, View
from django.views.generic.list import ListView

# 3rd Party Libraries
from channels.layers import get_channel_layer
from docx.image.exceptions import UnrecognizedImageError
from docx.opc.exceptions import PackageNotFoundError as DocxPackageNotFoundError
from pptx.exc import PackageNotFoundError as PptxPackageNotFoundError

# Ghostwriter Libraries
from ghostwriter.api.utils import (
    ForbiddenJsonResponse,
    get_archives_list,
    get_reports_list,
    get_templates_list,
    verify_finding_access,
    verify_access,
    verify_user_is_privileged,
    RoleBasedAccessControlMixin,
)
from ghostwriter.commandcenter.models import CompanyInformation, ReportConfiguration
from ghostwriter.modules import reportwriter
from ghostwriter.modules.exceptions import MissingTemplate
from ghostwriter.modules.model_utils import to_dict
from ghostwriter.reporting.filters import ArchiveFilter, FindingFilter, ReportFilter
from ghostwriter.reporting.forms import (
    EvidenceForm,
    FindingForm,
    FindingNoteForm,
    LocalFindingNoteForm,
    ReportFindingLinkUpdateForm,
    ReportForm,
    ReportTemplateForm,
    SelectReportTemplateForm,
)
from ghostwriter.reporting.models import (
    Archive,
    Evidence,
    Finding,
    FindingNote,
    FindingType,
    LocalFindingNote,
    Report,
    ReportFindingLink,
    ReportTemplate,
    Severity,
)
from ghostwriter.reporting.resources import FindingResource
from ghostwriter.rolodex.models import Project, ProjectAssignment

channel_layer = get_channel_layer()

User = get_user_model()

# Using __name__ resolves to ghostwriter.reporting.views
logger = logging.getLogger(__name__)


def get_position(report_pk, severity):
    findings = ReportFindingLink.objects.filter(Q(report__pk=report_pk) & Q(severity=severity)).order_by("-position")
    if findings:
        # Set new position to be one above the last/largest position
        last_position = findings[0].position
        return last_position + 1
    return 1


def generate_report_name(report_instance):
    """
    Generate a filename for a report based on the current time and attributes of an
    individual :model:`reporting.Report`. All illegal characters are removed to keep
    the filename browser-friendly.
    """

    def replace_placeholders(name, instance):
        """Replace placeholders in the report name with the appropriate values."""
        company_info = CompanyInformation.get_solo()
        name = name.replace("{title}", instance.title)
        name = name.replace("{company}", company_info.company_name)
        name = name.replace("{client}", instance.project.client.name)
        name = name.replace("{date}", dateformat.format(timezone.now(), settings.DATE_FORMAT))
        name = name.replace("{assessment_type}", instance.project.project_type.project_type)
        return name

    def replace_date_format(name):
        """Replace date format placeholders in the report name with the appropriate values."""
        # Find all strings wrapped in curly braces
        datetime_regex = r"(?<=\{)(.*?)(?=\})"
        for match in re.findall(datetime_regex, name):
            strfmt = dateformat.format(timezone.now(), match)
            name = name.replace(match, strfmt)
        return name

    def replace_chars(name):
        """Remove illegal characters from the report name."""
        name = name.replace("–", "-")
        return re.sub(r"[<>:;\"'/\\|?*.,{}\[\]]", "", name)

    report_config = ReportConfiguration.get_solo()
    report_name = report_config.report_filename
    report_name = replace_placeholders(report_name, report_instance)
    report_name = replace_date_format(report_name)
    report_name = replace_chars(report_name)
    return report_name.strip()


def zip_directory(path, zip_handler):
    """Compress the target directory as a Zip file for archiving."""
    # Walk the target directory
    abs_src = os.path.abspath(path)
    for root, _, files in os.walk(path):
        # Add each file to the zip file handler
        for file in files:
            absname = os.path.abspath(os.path.join(root, file))
            arcname = absname[len(abs_src) + 1 :]
            zip_handler.write(os.path.join(root, file), "evidence/" + arcname)


##################
# AJAX Functions #
##################


@login_required
def ajax_update_report_findings(request):
    """
    Update the ``position`` and ``severity`` fields of all :model:`reporting.ReportFindingLink`
    attached to an individual :model:`reporting.Report`.
    """
    if request.method == "POST" and request.is_ajax():
        pos = request.POST.get("positions")
        report_id = request.POST.get("report")
        weight = request.POST.get("weight")
        order = json.loads(pos)

        report = get_object_or_404(Report, pk=report_id)
        if verify_access(request.user, report.project):
            logger.info(
                "Received AJAX POST to update report %s's %s severity group findings in this order: %s",
                report_id,
                weight,
                ", ".join(order),
            )
            data = {"result": "success"}

            try:
                severity = Severity.objects.get(weight=weight)
            except Severity.DoesNotExist:
                severity = None
                logger.exception("Failed to get Severity object for weight %s", weight)

            if severity:
                counter = 1
                for finding_id in order:
                    if "placeholder" not in finding_id:
                        finding_instance = ReportFindingLink.objects.get(id=finding_id)
                        if finding_instance:
                            finding_instance.severity = severity
                            finding_instance.position = counter
                            finding_instance.save()
                            counter += 1
                        else:
                            logger.error(
                                "Received a finding ID, %s, that did not match an existing finding",
                                finding_id,
                            )
            else:
                data = {"result": "error", "message": "Specified severity weight, {}, is invalid.".format(weight)}
        else:
            logger.error(
                "AJAX request submitted by user %s without access to report %s",
                request.user,
                report_id,
            )
            data = {"result": "error"}
    else:
        data = {"result": "error"}
    return JsonResponse(data)


class UpdateTemplateLintResults(RoleBasedAccessControlMixin, SingleObjectMixin, View):
    """
    Return an updated version of the template following a request to update linter results
    for an individual :model:`reporting.ReportTemplate`.

    **Template**

    :template:`snippets/template_lint_results.html`
    """

    model = ReportTemplate

    def get(self, *args, **kwargs):
        template = self.get_object()
        html = render_to_string(
            "snippets/template_lint_results.html",
            {"reporttemplate": template},
        )
        return HttpResponse(html)


class AssignFinding(RoleBasedAccessControlMixin, SingleObjectMixin, View):
    """
    Copy an individual :model:`reporting.Finding` to create a new
    :model:`reporting.ReportFindingLink` connected to the user's active
    :model:`reporting.Report`.
    """

    model = Finding

    def post(self, *args, **kwargs):
        finding_instance = self.get_object()
        finding_dict = to_dict(finding_instance, resolve_fk=True)

        # Remove the tags from the finding dict to add them later with the ``taggit`` API
        del finding_dict["tags"]
        del finding_dict["tagged_items"]

        # The user must have the ``active_report`` session variable
        active_report = self.request.session.get("active_report", None)
        if active_report:
            try:
                report = Report.objects.get(pk=active_report["id"])
                if not verify_access(self.request.user, report.project):
                    return ForbiddenJsonResponse()
            except Report.DoesNotExist:
                message = "Please select a report to edit before trying to assign a finding."
                data = {"result": "error", "message": message}
                return JsonResponse(data)

            # Clone the selected object to make a new :model:`reporting.ReportFindingLink`
            report_link = ReportFindingLink(
                report=report,
                assigned_to=self.request.user,
                position=get_position(report.id, finding_instance.severity),
                **finding_dict,
            )
            report_link.save()
            report_link.tags.add(*finding_instance.tags.all())

            message = "{} successfully added to your active report.".format(finding_instance)
            data = {"result": "success", "message": message}
            logger.info(
                "Copied %s %s to %s %s (%s %s) by request of %s",
                finding_instance.__class__.__name__,
                finding_instance.id,
                report.__class__.__name__,
                report.id,
                report_link.__class__.__name__,
                report_link.id,
                self.request.user,
            )
        else:
            message = "Please select a report to edit before trying to assign a finding."
            data = {"result": "error", "message": message}
        return JsonResponse(data)


class LocalFindingNoteDelete(RoleBasedAccessControlMixin, SingleObjectMixin, View):
    """Delete an individual :model:`reporting.LocalFindingNote`."""

    model = LocalFindingNote

    def test_func(self):
        note = self.get_object()
        return note.operator.id == self.request.user.id

    def handle_no_permission(self):
        messages.error(self.request, "You do not have permission to access that.")
        return redirect("home:dashboard")

    def post(self, *args, **kwargs):
        note = self.get_object()
        note.delete()
        data = {"result": "success", "message": "Note successfully deleted!"}
        logger.info(
            "Deleted %s %s by request of %s",
            note.__class__.__name__,
            note.id,
            self.request.user,
        )
        return JsonResponse(data)


class FindingNoteDelete(RoleBasedAccessControlMixin, SingleObjectMixin, View):
    """Delete an individual :model:`reporting.FindingNote`."""

    model = FindingNote

    def test_func(self):
        note = self.get_object()
        return note.operator.id == self.request.user.id

    def handle_no_permission(self):
        messages.error(self.request, "You do not have permission to access that.")
        return redirect("home:dashboard")

    def post(self, *args, **kwargs):
        note = self.get_object()
        note.delete()
        data = {"result": "success", "message": "Note successfully deleted!"}
        logger.info(
            "Deleted %s %s by request of %s",
            note.__class__.__name__,
            note.id,
            self.request.user,
        )
        return JsonResponse(data)


class ReportFindingLinkDelete(RoleBasedAccessControlMixin, SingleObjectMixin, View):
    """Delete an individual :model:`reporting.ReportFindingLink`."""

    model = ReportFindingLink

    def test_func(self):
        return verify_access(self.request.user, self.get_object().report.project)

    def handle_no_permission(self):
        return ForbiddenJsonResponse()

    def post(self, *args, **kwargs):
        finding = self.get_object()
        finding.delete()
        data = {
            "result": "success",
            "message": "Successfully deleted {finding} and cleaned up evidence.".format(finding=finding),
        }
        logger.info(
            "Deleted %s %s by request of %s",
            finding.__class__.__name__,
            finding.id,
            self.request.user,
        )

        return JsonResponse(data)


class ReportActivate(RoleBasedAccessControlMixin, SingleObjectMixin, View):
    """Set an individual :model:`reporting.Report` as active for the current user session."""

    model = Report

    def test_func(self):
        return verify_access(self.request.user, self.get_object().project)

    def handle_no_permission(self):
        return ForbiddenJsonResponse()

    def post(self, *args, **kwargs):
        report = self.get_object()
        try:
            self.request.session["active_report"] = {}
            self.request.session["active_report"]["id"] = report.id
            self.request.session["active_report"]["title"] = report.title
            message = "{report} is now your active report and you will be redirected there in 5 seconds...".format(
                report=report.title
            )
            data = {
                "result": "success",
                "report": report.title,
                "report_url": report.get_absolute_url(),
                "message": message,
            }
        except Exception as exception:  # pragma: no cover
            template = "An exception of type {0} occurred. Arguments:\n{1!r}"
            log_message = template.format(type(exception).__name__, exception.args)
            logger.error(log_message)
            data = {
                "result": "error",
                "message": "Could not set the selected report as your active report!",
            }

        return JsonResponse(data)


class ReportStatusToggle(RoleBasedAccessControlMixin, SingleObjectMixin, View):
    """Toggle the ``complete`` field of an individual :model:`rolodex.Report`."""

    model = Report

    def test_func(self):
        return verify_access(self.request.user, self.get_object().project)

    def handle_no_permission(self):
        return ForbiddenJsonResponse()

    def post(self, *args, **kwargs):
        report = self.get_object()
        try:
            if report.complete:
                report.complete = False
                data = {
                    "result": "success",
                    "message": "Report successfully marked as incomplete.",
                    "status": "Draft",
                    "toggle": 0,
                }
            else:
                report.complete = True
                data = {
                    "result": "success",
                    "message": "Report successfully marked as complete.",
                    "status": "Complete",
                    "toggle": 1,
                }
            report.save()
            logger.info(
                "Toggled status of %s %s by request of %s",
                report.__class__.__name__,
                report.id,
                self.request.user,
            )
        except Exception as exception:  # pragma: no cover
            template = "An exception of type {0} occurred. Arguments:\n{1!r}"
            log_message = template.format(type(exception).__name__, exception.args)
            logger.error(log_message)
            data = {"result": "error", "message": "Could not update report's status!"}

        return JsonResponse(data)


class ReportDeliveryToggle(RoleBasedAccessControlMixin, SingleObjectMixin, View):
    """Toggle the ``delivered`` field of an individual :model:`rolodex.Report`."""

    model = Report

    def test_func(self):
        return verify_access(self.request.user, self.get_object().project)

    def handle_no_permission(self):
        return ForbiddenJsonResponse()

    def post(self, *args, **kwargs):
        report = self.get_object()
        try:
            if report.delivered:
                report.delivered = False
                data = {
                    "result": "success",
                    "message": "Report successfully marked as not delivered.",
                    "status": "Not Delivered",
                    "toggle": 0,
                }
            else:
                report.delivered = True
                data = {
                    "result": "success",
                    "message": "Report successfully marked as delivered.",
                    "status": "Delivered",
                    "toggle": 1,
                }
            report.save()
            logger.info(
                "Toggled delivery status of %s %s by request of %s",
                report.__class__.__name__,
                report.id,
                self.request.user,
            )
        except Exception as exception:  # pragma: no cover
            template = "An exception of type {0} occurred. Arguments:\n{1!r}"
            log_message = template.format(type(exception).__name__, exception.args)
            logger.error(log_message)
            data = {
                "result": "error",
                "message": "Could not update report's delivery status!",
            }

        return JsonResponse(data)


class ReportFindingStatusUpdate(RoleBasedAccessControlMixin, SingleObjectMixin, View):
    """Update the ``complete`` field of an individual :model:`reporting.ReportFindingLink`."""

    model = ReportFindingLink

    def test_func(self):
        return verify_access(self.request.user, self.get_object().report.project)

    def handle_no_permission(self):
        return ForbiddenJsonResponse()

    def post(self, *args, **kwargs):
        # Get ``status`` kwargs from the URL
        status = self.kwargs["status"]
        finding = self.get_object()

        try:
            result = "success"
            if status.lower() == "edit":
                finding.complete = False
                message = "Successfully flagged finding for editing."
                display_status = "Needs Editing"
                classes = "burned"
            elif status.lower() == "complete":
                finding.complete = True
                message = "Successfully marking finding as complete."
                display_status = "Ready"
                classes = "healthy"
            else:
                result = "error"
                message = "Could not update the finding's status to: {}".format(status)
                display_status = "Error"
                classes = "burned"
            finding.save()
            # Prepare the JSON response data
            data = {
                "result": result,
                "status": display_status,
                "classes": classes,
                "message": message,
            }
            logger.info(
                "Set status of %s %s to %s by request of %s",
                finding.__class__.__name__,
                finding.id,
                status,
                self.request.user,
            )
        # Return an error message if the query for the requested status returned DoesNotExist
        except Exception as exception:  # pragma: no cover
            template = "An exception of type {0} occurred. Arguments:\n{1!r}"
            log_message = template.format(type(exception).__name__, exception.args)
            logger.error(log_message)
            data = {"result": "error", "message": "Could not update finding's status!"}

        return JsonResponse(data)


class ReportTemplateSwap(RoleBasedAccessControlMixin, SingleObjectMixin, View):
    """Update the ``template`` value for an individual :model:`reporting.Report`."""

    model = Report

    def test_func(self):
        return verify_access(self.request.user, self.get_object().project)

    def handle_no_permission(self):
        return ForbiddenJsonResponse()

    def post(self, *args, **kwargs):
        report = self.get_object()
        docx_template_id = self.request.POST.get("docx_template", None)
        pptx_template_id = self.request.POST.get("pptx_template", None)
        if docx_template_id and pptx_template_id:
            docx_template_query = None
            pptx_template_query = None
            try:
                docx_template_id = int(docx_template_id)
                pptx_template_id = int(pptx_template_id)

                if docx_template_id < 0:
                    report.docx_template = None
                if pptx_template_id < 0:
                    report.pptx_template = None
                if docx_template_id >= 0:
                    docx_template_query = ReportTemplate.objects.get(pk=docx_template_id)
                    report.docx_template = docx_template_query
                if pptx_template_id >= 0:
                    pptx_template_query = ReportTemplate.objects.get(pk=pptx_template_id)
                    report.pptx_template = pptx_template_query
                data = {
                    "result": "success",
                    "message": "Templates successfully updated.",
                }
                report.save()

                # Check template for linting issues
                try:
                    if docx_template_query:
                        template_status = docx_template_query.get_status()
                        data["docx_lint_result"] = template_status
                        if template_status != "success":
                            if template_status == "warning":
                                data[
                                    "docx_lint_message"
                                ] = "Selected Word template has warnings from linter. Check the template before generating a report."
                            elif template_status == "error":
                                data[
                                    "docx_lint_message"
                                ] = "Selected Word template has linting errors and cannot be used to generate a report."
                            elif template_status == "failed":
                                data[
                                    "docx_lint_message"
                                ] = "Selected Word template failed basic linter checks and can't be used to generate a report."
                            else:
                                data[
                                    "docx_lint_message"
                                ] = "Selected Word template has an unknown linter status. Check and lint the template before generating a report."
                            data["docx_url"] = docx_template_query.get_absolute_url()
                except Exception:  # pragma: no cover
                    logger.exception("Failed to get the template status")
                    data["docx_lint_result"] = "failed"
                    data[
                        "docx_lint_message"
                    ] = "Could not retrieve the Word template's linter status. Check and lint the template before generating a report."
                try:
                    if pptx_template_query:
                        template_status = pptx_template_query.get_status()
                        data["pptx_lint_result"] = template_status
                        if template_status != "success":
                            if template_status == "warning":
                                data[
                                    "pptx_lint_message"
                                ] = "Selected PowerPoint template has warnings from linter. Check the template before generating a report."
                            elif template_status == "error":
                                data[
                                    "pptx_lint_message"
                                ] = "Selected PowerPoint template has linting errors and cannot be used to generate a report."
                            elif template_status == "failed":
                                data[
                                    "pptx_lint_message"
                                ] = "Selected PowerPoint template failed basic linter checks and can't be used to generate a report."
                            else:
                                data[
                                    "pptx_lint_message"
                                ] = "Selected PowerPoint template has an unknown linter status. Check and lint the template before generating a report."
                            data["pptx_url"] = pptx_template_query.get_absolute_url()
                except Exception:  # pragma: no cover
                    logger.exception("Failed to get the template status")
                    data["pptx_lint_result"] = "failed"
                    data[
                        "pptx_lint_message"
                    ] = "Could not retrieve the PowerPoint template's linter status. Check and lint the template before generating a report."
                logger.info(
                    "Swapped template for %s %s by request of %s",
                    report.__class__.__name__,
                    report.id,
                    self.request.user,
                )
            except ValueError:
                data = {
                    "result": "error",
                    "message": "Submitted template ID was not an integer.",
                }
                logger.error(
                    "Received one or two invalid (non-integer) template IDs (%s & %s) from a request submitted by %s",
                    docx_template_id,
                    pptx_template_id,
                    self.request.user,
                )
            except ReportTemplate.DoesNotExist:
                data = {
                    "result": "error",
                    "message": "Submitted template ID does not exist.",
                }
                logger.error(
                    "Received one or two invalid (non-existent) template IDs (%s & %s) from a request submitted by %s",
                    docx_template_id,
                    pptx_template_id,
                    self.request.user,
                )
            except Exception:  # pragma: no cover
                data = {
                    "result": "error",
                    "message": "An exception prevented the template change.",
                }
                logger.exception(
                    "Encountered an error trying to update %s %s with template IDs %s & %s from a request submitted by %s",
                    report.__class__.__name__,
                    report.id,
                    docx_template_id,
                    pptx_template_id,
                    self.request.user,
                )
        else:
            data = {"result": "error", "message": "Submitted request was incomplete."}
            logger.warning(
                "Received bad template IDs (%s & %s) from a request submitted by %s",
                docx_template_id,
                pptx_template_id,
                self.request.user,
            )
        return JsonResponse(data)


class ReportTemplateLint(RoleBasedAccessControlMixin, SingleObjectMixin, View):
    """
    Check an individual :model:`reporting.ReportTemplate` for Jinja2 syntax errors
    and undefined variables.
    """

    model = ReportTemplate

    def post(self, *args, **kwargs):
        template = self.get_object()
        linter = reportwriter.TemplateLinter(template=template)
        if template.doc_type.doc_type == "docx":
            results = linter.lint_docx()
        elif template.doc_type.doc_type == "pptx":
            results = linter.lint_pptx()
        else:
            logger.warning(
                "Template had an unknown filetype not supported by the linter: %s",
                template.doc_type,
            )
            results = {}
        template.lint_result = results
        template.save()

        data = results
        if data["result"] == "success":
            data["message"] = "Template linter returned results with no errors or warnings."
        elif not data["result"]:
            data["message"] = f"Template had an unknown filetype not supported by the linter: {template.doc_type}"
        else:
            data["message"] = "Template linter returned results with issues that require attention."

        return JsonResponse(data)


class ReportClone(RoleBasedAccessControlMixin, SingleObjectMixin, View):
    """Create an identical copy of an individual :model:`reporting.Report`."""

    model = Report

    def test_func(self):
        return verify_access(self.request.user, self.get_object().project)

    def handle_no_permission(self):
        return ForbiddenJsonResponse()

    def get(self, *args, **kwargs):
        report_to_clone = self.get_object()
        report_pk = None
        try:
            findings = ReportFindingLink.objects.select_related("report").filter(report=report_to_clone.pk)
            report_to_clone.title = report_to_clone.title + " Copy"
            report_to_clone.complete = False
            report_to_clone.pk = None
            report_to_clone.save()
            report_pk = report_to_clone.pk
            for finding in findings:
                # Get any evidence files attached to the original finding
                evidences = Evidence.objects.filter(finding=finding.pk)
                # Create a clone of this finding attached to the new report
                finding.report = report_to_clone
                finding.pk = None
                finding.save()
                # Clone evidence files and attach them to the new finding
                for evidence in evidences:
                    if exists(evidence.document.path):
                        evidence_file = File(evidence.document, os.path.basename(evidence.document.name))
                        evidence.finding = finding
                        evidence._current_evidence = None
                        evidence.document = evidence_file
                        evidence.pk = None
                        evidence.save()
                    else:
                        logger.warning(
                            "Evidence file not found: %s",
                            evidence.document.path,
                        )
                        messages.warning(
                            self.request,
                            f"An evidence file was missing and could not be copied: {evidence.friendly_name} ({os.path.basename(evidence.document.name)})",
                            extra_tags="alert-warning",
                        )

            logger.info(
                "Cloned %s %s by request of %s",
                report_to_clone.__class__.__name__,
                report_to_clone.id,
                self.request.user,
            )

            messages.success(
                self.request,
                "Successfully cloned your report: {}".format(report_to_clone.title),
                extra_tags="alert-error",
            )
        except Exception as exception:  # pragma: no cover
            template = "An exception of type {0} occurred. Arguments:\n{1!r}"
            log_message = template.format(type(exception).__name__, exception.args)
            logger.error(log_message)

            messages.error(
                self.request,
                "Encountered an error while trying to clone your report: {}".format(exception.args),
                extra_tags="alert-error",
            )

        return HttpResponseRedirect(reverse("reporting:report_detail", kwargs={"pk": report_pk}))


class AssignBlankFinding(RoleBasedAccessControlMixin, SingleObjectMixin, View):
    """
    Create a blank :model:`reporting.ReportFindingLink` entry linked to an individual
    :model:`reporting.Report`.
    """

    model = Report

    def test_func(self):
        return verify_access(self.request.user, self.get_object().project)

    def handle_no_permission(self):
        return ForbiddenJsonResponse()

    def __init__(self):
        self.severity = Severity.objects.order_by("weight").last()
        self.finding_type = FindingType.objects.all().first()
        super().__init__()

    def get(self, *args, **kwargs):
        obj = self.get_object()
        try:
            report_link = ReportFindingLink(
                title="Blank Template",
                severity=self.severity,
                finding_type=self.finding_type,
                report=obj,
                assigned_to=self.request.user,
                position=get_position(obj.id, self.severity),
                added_as_blank=True,
            )
            report_link.save()

            logger.info(
                "Added a blank finding to %s %s by request of %s",
                obj.__class__.__name__,
                obj.id,
                self.request.user,
            )

            messages.success(
                self.request,
                "Successfully added a blank finding to the report",
                extra_tags="alert-success",
            )
        except Exception as exception:  # pragma: no cover
            template = "An exception of type {0} occurred. Arguments:\n{1!r}"
            log_message = template.format(type(exception).__name__, exception.args)
            logger.error(log_message)

            messages.error(
                self.request,
                "Encountered an error while trying to add a blank finding to your report: {}".format(exception.args),
                extra_tags="alert-error",
            )

        return HttpResponseRedirect(reverse("reporting:report_detail", args=(obj.id,)))


class ConvertFinding(RoleBasedAccessControlMixin, SingleObjectMixin, View):
    """
    Create a copy of an individual :model:`reporting.ReportFindingLink` and prepare
    it to be saved as a new :model:`reporting.Finding`.

    **Template**

    :template:`reporting/finding_form.html`
    """

    model = ReportFindingLink

    def test_func(self):
        if verify_access(self.request.user, self.get_object().report.project):
            if verify_finding_access(self.request.user, "create"):
                return True
        return False

    def handle_no_permission(self):
        messages.error(self.request, "You do not have the necessary permission to create new findings.")
        return redirect(reverse("reporting:report_detail", kwargs={"pk": self.get_object().report.pk}))

    def get(self, *args, **kwargs):
        finding_instance = self.get_object()
        try:
            form = FindingForm(
                initial={
                    "title": finding_instance.title,
                    "description": finding_instance.description,
                    "impact": finding_instance.impact,
                    "mitigation": finding_instance.mitigation,
                    "replication_steps": finding_instance.replication_steps,
                    "host_detection_techniques": finding_instance.host_detection_techniques,
                    "network_detection_techniques": finding_instance.network_detection_techniques,
                    "references": finding_instance.references,
                    "severity": finding_instance.severity,
                    "finding_type": finding_instance.finding_type,
                    "cvss_score": finding_instance.cvss_score,
                    "cvss_vector": finding_instance.cvss_vector,
                    "tags": finding_instance.tags.all(),
                }
            )
        except Exception as exception:  # pragma: no cover
            template = "An exception of type {0} occurred. Arguments:\n{1!r}"
            log_message = template.format(type(exception).__name__, exception.args)
            logger.error(log_message)

            messages.error(
                self.request,
                "Encountered an error while trying to convert your finding: {}".format(exception.args),
                extra_tags="alert-error",
            )
            return HttpResponse(status=500)

        return render(self.request, "reporting/finding_form.html", {"form": form})

    def post(self, *args, **kwargs):
        form = FindingForm(self.request.POST)
        if form.is_valid():
            new_finding = form.save()
            return HttpResponseRedirect(reverse("reporting:finding_detail", kwargs={"pk": new_finding.pk}))
        logger.warning(form.errors.as_data())
        return render(self.request, "reporting/finding_form.html", {"form": form})


##################
# View Functions #
##################


@login_required
def index(request):
    """Display the main homepage."""
    return HttpResponseRedirect(reverse("home:dashboard"))


@login_required
def archive_list(request):
    """
    Display a list of all :model:`reporting.Report` marked as archived.

    **Context**

    ``filter``
        Instance of :filter:`reporting.ArchiveFilter`

    **Template**

    :template:`reporting/archives.html`
    """
    archives = get_archives_list(request.user)
    archive_filter = ArchiveFilter(request.GET, queryset=archives)
    return render(request, "reporting/archives.html", {"filter": archive_filter})


@login_required
def upload_evidence_modal_success(request):
    """
    Display message following the successful creation of an individual
    :model:`reporting.Evidence` using a TinyMCE URLDialog.

    **Template**

    :template:`reporting/evidence_modal_success.html`
    """
    return render(request, "reporting/evidence_modal_success.html")


@login_required
def export_findings_to_csv(request):
    """Export all :model:`reporting.Finding` to a csv file for download."""
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    finding_resource = FindingResource()
    dataset = finding_resource.export()
    response = HttpResponse(dataset.csv, content_type="text/csv")
    response["Content-Disposition"] = f'attachment; filename="{timestamp}_findings.csv"'

    return response


################
# View Classes #
################

# CBVs related to :model:`reporting.Finding`


class FindingListView(RoleBasedAccessControlMixin, ListView):
    """
    Display a list of all :model:`reporting.Finding`.

    **Context**

    ``filter``
        Instance of :filter:`reporting.FindingFilter`

    **Template**

    :template:`reporting/finding_list.html`
    """

    model = Finding
    template_name = "reporting/finding_list.html"

    def __init__(self):
        super().__init__()
        self.autocomplete = []

    def get_queryset(self):
        search_term = ""
        findings = (
            Finding.objects.select_related("severity", "finding_type")
            .all()
            .order_by("severity__weight", "-cvss_score", "finding_type", "title")
        )

        # Build autocomplete list
        for finding in findings:
            self.autocomplete.append(finding.title)

        if "finding" in self.request.GET:
            search_term = self.request.GET.get("finding").strip()
            if search_term is None or search_term == "":
                search_term = ""
        if search_term:
            messages.success(
                self.request,
                "Displaying search results for: {}".format(search_term),
                extra_tags="alert-success",
            )
            return findings.filter(Q(title__icontains=search_term) | Q(description__icontains=search_term)).order_by(
                "severity__weight", "-cvss_score", "finding_type", "title"
            )
        return findings

    def get(self, request, *args, **kwarg):
        findings_filter = FindingFilter(request.GET, queryset=self.get_queryset())
        return render(
            request, "reporting/finding_list.html", {"filter": findings_filter, "autocomplete": self.autocomplete}
        )


class FindingDetailView(RoleBasedAccessControlMixin, DetailView):
    """
    Display an individual :model:`reporting.Finding`.

    **Template**

    :template:`reporting/finding_detail.html`
    """

    model = Finding


class FindingCreate(RoleBasedAccessControlMixin, CreateView):
    """
    Create an individual instance of :model:`reporting.Finding`.

    **Context**

    ``cancel_link``
        Link for the form's Cancel button to return to clients list page

    **Template**

    :template:`reporting/finding_form.html`
    """

    model = Finding
    form_class = FindingForm

    def test_func(self):
        return verify_finding_access(self.request.user, "create")

    def handle_no_permission(self):
        messages.error(self.request, "You do not have the necessary permission to create new findings.")
        return redirect("reporting:findings")

    def get_context_data(self, **kwargs):
        ctx = super().get_context_data(**kwargs)
        ctx["cancel_link"] = reverse("reporting:findings")
        return ctx

    def get_success_url(self):
        messages.success(
            self.request,
            "Successfully added {} to the findings library".format(self.object.title),
            extra_tags="alert-success",
        )
        return reverse("reporting:finding_detail", kwargs={"pk": self.object.pk})


class FindingUpdate(RoleBasedAccessControlMixin, UpdateView):
    """
    Update an individual instance of :model:`reporting.Finding`.

    **Context**

    ``cancel_link``
        Link for the form's Cancel button to return to clients list page

    **Template**

    :template:`reporting/finding_form.html`
    """

    model = Finding
    form_class = FindingForm

    def test_func(self):
        return verify_finding_access(self.request.user, "edit")

    def handle_no_permission(self):
        messages.error(self.request, "You do not have the necessary permission to edit findings.")
        return redirect(reverse("reporting:finding_detail", kwargs={"pk": self.get_object().pk}))

    def get_context_data(self, **kwargs):
        ctx = super().get_context_data(**kwargs)
        ctx["cancel_link"] = reverse("reporting:finding_detail", kwargs={"pk": self.object.pk})
        return ctx

    def get_success_url(self):
        messages.success(
            self.request,
            "Master record for {} was successfully updated".format(self.get_object().title),
            extra_tags="alert-success",
        )
        return reverse("reporting:finding_detail", kwargs={"pk": self.object.pk})


class FindingDelete(RoleBasedAccessControlMixin, DeleteView):
    """
    Delete an individual instance of :model:`reporting.Finding`.

    **Context**

    ``object_type``
        String describing what is to be deleted
    ``object_to_be_deleted``
        To-be-deleted instance of :model:`reporting.Finding`
    ``cancel_link``
        Link for the form's Cancel button to return to finding list page

    **Template**

    :template:`confirm_delete.html`
    """

    model = Finding
    template_name = "confirm_delete.html"

    def test_func(self):
        return verify_finding_access(self.request.user, "delete")

    def handle_no_permission(self):
        messages.error(self.request, "You do not have the necessary permission to delete findings.")
        return redirect(reverse("reporting:finding_detail", kwargs={"pk": self.get_object().pk}))

    def get_success_url(self):
        messages.warning(
            self.request,
            "Master record for {} was successfully deleted".format(self.get_object().title),
            extra_tags="alert-warning",
        )
        return reverse_lazy("reporting:findings")

    def get_context_data(self, **kwargs):
        ctx = super().get_context_data(**kwargs)
        queryset = kwargs["object"]
        ctx["object_type"] = "finding master record"
        ctx["object_to_be_deleted"] = queryset.title
        ctx["cancel_link"] = reverse("reporting:findings")
        return ctx


# CBVs related to :model:`reporting.Report`


class ReportListView(RoleBasedAccessControlMixin, ListView):
    """
    Display a list of all :model:`reporting.Report`.

    **Template**

    :template:`reporting/report_list.html`
    """

    model = Finding
    template_name = "reporting/report_list.html"

    def get_queryset(self):
        return get_reports_list(self.request.user)

    def get(self, request, *args, **kwarg):
        reports_filter = ReportFilter(request.GET, queryset=self.get_queryset())
        return render(request, "reporting/report_list.html", {"filter": reports_filter})


class ArchiveView(RoleBasedAccessControlMixin, SingleObjectMixin, View):
    """
    Generate all report types for an individual :model:`reporting.Report`, collect all
    related :model:`reporting.Evidence` and related files, and compress the files into a
    single Zip file for archiving.
    """

    model = Report

    def test_func(self):
        return verify_access(self.request.user, self.get_object().project)

    def handle_no_permission(self):
        messages.error(self.request, "You do not have permission to access that.")
        return redirect("home:dashboard")

    def get(self, *args, **kwargs):
        report_instance = self.get_object()
        try:
            archive_loc = os.path.join(settings.MEDIA_ROOT, "archives/")
            evidence_loc = os.path.join(settings.MEDIA_ROOT, "evidence", str(report_instance.id))
            report_name = generate_report_name(report_instance)

            # Get the templates for Word and PowerPoint
            report_config = ReportConfiguration.get_solo()
            if report_instance.docx_template:
                docx_template = report_instance.docx_template.document.path
            else:
                docx_template = report_config.default_docx_template
                if not docx_template:
                    raise MissingTemplate
            if report_instance.pptx_template:
                pptx_template = report_instance.pptx_template.document.path
            else:
                pptx_template = report_config.default_docx_template
                if not pptx_template:
                    raise MissingTemplate

            engine = reportwriter.Reportwriter(report_instance, template_loc=None)
            json_doc, word_doc, excel_doc, ppt_doc = engine.generate_all_reports(docx_template, pptx_template)

            # Create a zip file in memory and add the reports to it
            zip_buffer = io.BytesIO()
            with zipfile.ZipFile(zip_buffer, "a") as zf:
                zf.writestr("report.json", json_doc)
                zf.writestr("report.docx", word_doc.getvalue())
                zf.writestr("report.xlsx", excel_doc.getvalue())
                zf.writestr("report.pptx", ppt_doc.getvalue())
                zip_directory(evidence_loc, zf)
            zip_buffer.seek(0)
            with open(os.path.join(archive_loc, report_name + ".zip"), "wb+") as archive_file:
                archive_file = ContentFile(zip_buffer.read(), name=report_name + ".zip")
                new_archive = Archive(
                    project=report_instance.project,
                    report_archive=File(archive_file),
                )
            new_archive.save()
            messages.success(
                self.request,
                "Successfully archived {}!".format(report_instance.title),
                extra_tags="alert-success",
            )
            return HttpResponseRedirect(reverse("reporting:archived_reports"))
        except MissingTemplate:
            logger.error(
                "Archive generation failed for %s %s and user %s because no template was configured.",
                report_instance.__class__.__name__,
                report_instance.id,
                self.request.user,
            )
            messages.error(
                self.request,
                "You do not have a Word or PowerPoint template selected and have not configured a default template.",
                extra_tags="alert-danger",
            )
        except DocxPackageNotFoundError:
            logger.exception(
                "DOCX generation failed for %s %s and user %s because the template file was missing.",
                report_instance.__class__.__name__,
                report_instance.id,
                self.request.user,
            )
            messages.error(
                self.request,
                "Your selected Word template could not be found on the server – try uploading it again.",
                extra_tags="alert-danger",
            )
        except PptxPackageNotFoundError:
            logger.exception(
                "PPTX generation failed for %s %s and user %s because the template file was missing.",
                report_instance.__class__.__name__,
                report_instance.id,
                self.request.user,
            )
            messages.error(
                self.request,
                "Your selected PowerPoint template could not be found on the server – try uploading it again.",
                extra_tags="alert-danger",
            )
        except Exception:
            logger.exception("Error archiving report.")
            messages.error(
                self.request,
                "Failed to generate one or more documents for the archive.",
                extra_tags="alert-danger",
            )
        return HttpResponseRedirect(reverse("reporting:report_detail", kwargs={"pk": report_instance.id}))


class ArchiveDownloadView(RoleBasedAccessControlMixin, SingleObjectMixin, View):
    """Return the target :model:`reporting.Report` archive file for download."""

    model = Archive

    def test_func(self):
        return verify_access(self.request.user, self.get_object().project)

    def handle_no_permission(self):
        messages.error(self.request, "You do not have permission to access that.")
        return redirect("home:dashboard")

    def get(self, *args, **kwargs):
        archive_instance = self.get_object()
        file_path = os.path.join(settings.MEDIA_ROOT, archive_instance.report_archive.path)
        if os.path.exists(file_path):
            with open(file_path, "rb") as archive_file:
                response = HttpResponse(archive_file.read(), content_type="application/x-zip-compressed")
                response["Content-Disposition"] = "attachment; filename=" + os.path.basename(file_path)
                return response
        raise Http404


class ReportDetailView(RoleBasedAccessControlMixin, DetailView):
    """
    Display an individual :model:`reporting.Report`.

    **Template**

    :template:`reporting/report_detail.html`
    """

    model = Report

    def __init__(self):
        super().__init__()
        self.autocomplete = []

    def test_func(self):
        return verify_access(self.request.user, self.get_object().project)

    def handle_no_permission(self):
        messages.error(self.request, "You do not have permission to access that.")
        return redirect("reporting:reports")

    def get_context_data(self, **kwargs):
        ctx = super().get_context_data(**kwargs)
        form = SelectReportTemplateForm(instance=self.object)
        form.fields["docx_template"].queryset = ReportTemplate.objects.filter(
            Q(doc_type__doc_type="docx") & Q(client=self.object.project.client)
            | Q(doc_type__doc_type="docx") & Q(client__isnull=True)
        ).select_related(
            "doc_type",
            "client",
        )
        form.fields["pptx_template"].queryset = ReportTemplate.objects.filter(
            Q(doc_type__doc_type="pptx") & Q(client=self.object.project.client)
            | Q(doc_type__doc_type="pptx") & Q(client__isnull=True)
        ).select_related(
            "doc_type",
            "client",
        )
        ctx["form"] = form

        # Build autocomplete list
        findings = (
            Finding.objects.select_related("severity", "finding_type")
            .all()
            .order_by("severity__weight", "-cvss_score", "finding_type", "title")
        )
        for finding in findings:
            self.autocomplete.append(finding.title)
        ctx["autocomplete"] = self.autocomplete

        return ctx


class ReportCreate(RoleBasedAccessControlMixin, CreateView):
    """
    Create an individual instance of :model:`reporting.Report`.

    **Context**

    ``project``
        Instance of :model:`rolodex.Project` associated with this report
    ``cancel_link``
        Link for the form's Cancel button to return to report list or details page

    **Template**

    :template:`reporting/report_form.html`
    """

    model = Report
    form_class = ReportForm

    def setup(self, request, *args, **kwargs):
        super().setup(request, *args, **kwargs)
        # Check if this request is for a specific project or not
        self.project = ""
        # Determine if ``pk`` is in the kwargs
        if "pk" in self.kwargs:
            pk = self.kwargs.get("pk")
            # Try to get the project from :model:`rolodex.Project`
            if pk:
                try:
                    project = get_object_or_404(Project, pk=self.kwargs.get("pk"))
                    if verify_access(self.request.user, project):
                        self.project = project
                except Project.DoesNotExist:
                    logger.info(
                        "Received report create request for Project ID %s, but that Project does not exist",
                        pk,
                    )

    def get_form_kwargs(self):
        kwargs = super().get_form_kwargs()
        kwargs.update({"project": self.project, "user": self.request.user})
        return kwargs

    def get_context_data(self, **kwargs):
        ctx = super().get_context_data(**kwargs)
        ctx["project"] = self.project
        if self.project:
            ctx["cancel_link"] = reverse("rolodex:project_detail", kwargs={"pk": self.project.pk})
        else:
            ctx["cancel_link"] = reverse("reporting:reports")
        return ctx

    def get_form(self, form_class=None):
        form = super().get_form(form_class)
        if not form.fields["project"].queryset:
            messages.error(
                self.request,
                "There are no active projects for a new report",
                extra_tags="alert-error",
            )
        return form

    def form_valid(self, form):
        form.instance.created_by = self.request.user
        self.request.session["active_report"] = {}
        self.request.session["active_report"]["title"] = form.instance.title
        return super().form_valid(form)

    def get_initial(self):
        if self.project:
            title = "{} {} ({}) Report".format(self.project.client, self.project.project_type, self.project.start_date)
            return {"title": title, "project": self.project.id}
        return super().get_initial()

    def get_success_url(self):
        self.request.session["active_report"]["id"] = self.object.pk
        self.request.session.modified = True
        messages.success(
            self.request,
            "Successfully created new report and set it as your active report",
            extra_tags="alert-success",
        )
        return reverse("reporting:report_detail", kwargs={"pk": self.object.pk})


class ReportUpdate(RoleBasedAccessControlMixin, UpdateView):
    """
    Update an individual instance of :model:`reporting.Report`.

    **Context**

    ``cancel_link``
        Link for the form's Cancel button to return to report's detail page

    **Template**

    :template:`reporting/report_form.html`
    """

    model = Report
    form_class = ReportForm

    def test_func(self):
        return verify_access(self.request.user, self.get_object().project)

    def handle_no_permission(self):
        messages.error(self.request, "You do not have permission to access that.")
        return redirect("reporting:reports")

    def get_form_kwargs(self):
        kwargs = super().get_form_kwargs()
        kwargs.update({"project": self.get_object().project, "user": self.request.user})
        return kwargs

    def get_context_data(self, **kwargs):
        ctx = super().get_context_data(**kwargs)
        ctx["project"] = self.object.project
        ctx["cancel_link"] = reverse("reporting:report_detail", kwargs={"pk": self.object.pk})
        return ctx

    def form_valid(self, form):
        self.request.session["active_report"] = {}
        self.request.session["active_report"]["id"] = form.instance.id
        self.request.session["active_report"]["title"] = form.instance.title
        self.request.session.modified = True
        return super().form_valid(form)

    def get_success_url(self):
        messages.success(self.request, "Successfully updated the report", extra_tags="alert-success")
        return reverse("reporting:report_detail", kwargs={"pk": self.object.pk})


class ReportDelete(RoleBasedAccessControlMixin, DeleteView):
    """
    Delete an individual instance of :model:`reporting.Report`.

    **Context**

    ``object_type``
        String describing what is to be deleted
    ``object_to_be_deleted``
        To-be-deleted instance of :model:`reporting.Report`
    ``cancel_link``
        Link for the form's Cancel button to return to report's detail page

    **Template**

    :template:`confirm_delete.html`
    """

    model = Report
    template_name = "confirm_delete.html"

    def test_func(self):
        return verify_access(self.request.user, self.get_object().project)

    def handle_no_permission(self):
        messages.error(self.request, "You do not have permission to access that.")
        return redirect("reporting:reports")

    def get_success_url(self):
        # Clear user's session if deleted report is their active report
        if self.object.pk == self.request.session["active_report"]["id"]:
            self.request.session["active_report"] = {}
            self.request.session["active_report"]["id"] = ""
            self.request.session["active_report"]["title"] = ""
        self.request.session.modified = True
        messages.warning(
            self.request,
            "Successfully deleted the report and associated evidence files",
            extra_tags="alert-warning",
        )
        return "{}#reports".format(reverse("rolodex:project_detail", kwargs={"pk": self.object.project.id}))

    def get_context_data(self, **kwargs):
        ctx = super().get_context_data(**kwargs)
        queryset = kwargs["object"]
        ctx["cancel_link"] = reverse("rolodex:project_detail", kwargs={"pk": self.object.project.pk})
        ctx["object_type"] = "entire report, evidence and all"
        ctx["object_to_be_deleted"] = queryset.title
        return ctx


class ReportTemplateListView(RoleBasedAccessControlMixin, generic.ListView):
    """
    Display a list of all :model:`reporting.ReportTemplate`.

    **Template**

    :template:`reporting/report_template_list.html`
    """

    model = ReportTemplate
    template_name = "reporting/report_templates_list.html"

    def get_queryset(self):
        user = self.request.user
        queryset = get_templates_list(user)
        return queryset


class ReportTemplateDetailView(RoleBasedAccessControlMixin, DetailView):
    """
    Display an individual :model:`reporting.ReportTemplate`.

    **Template**

    :template:`reporting/report_template_list.html`
    """

    model = ReportTemplate
    template_name = "reporting/report_template_detail.html"

    def test_func(self):
        client = self.get_object().client
        if client:
            return verify_access(self.request.user, client)
        return self.request.user.is_active

    def handle_no_permission(self):
        messages.error(self.request, "You do not have permission to access that.")
        return redirect("reporting:templates")


class ReportTemplateCreate(RoleBasedAccessControlMixin, CreateView):
    """
    Create an individual instance of :model:`reporting.ReportTemplate`.

    **Context**

    ``cancel_link``
        Link for the form's Cancel button to return to template list page

    **Template**

    :template:`report_template_form.html`
    """

    model = ReportTemplate
    form_class = ReportTemplateForm
    template_name = "reporting/report_template_form.html"

    def get_context_data(self, **kwargs):
        ctx = super().get_context_data(**kwargs)
        ctx["cancel_link"] = reverse("reporting:templates")
        return ctx

    def get_initial(self):
        date = datetime.now().strftime("%d %B %Y")
        initial_upload = f'<p><span class="bold">{date}</span></p><p>Initial upload</p>'
        return {"changelog": initial_upload, "p_style": "Normal"}

    def get_success_url(self):
        messages.success(
            self.request,
            "Template successfully uploaded",
            extra_tags="alert-success",
        )
        return reverse("reporting:template_detail", kwargs={"pk": self.object.pk})

    def form_valid(self, form, **kwargs):
        self.object = form.save(commit=False)
        self.object.uploaded_by = self.request.user
        self.object.save()
        form.save_m2m()
        return HttpResponseRedirect(self.get_success_url())

    def get_form_kwargs(self):
        kwargs = super().get_form_kwargs()
        kwargs.update({"user": self.request.user})
        return kwargs


class ReportTemplateUpdate(RoleBasedAccessControlMixin, UpdateView):
    """
    Save an individual instance of :model:`reporting.ReportTemplate`.

    **Context**

    ``cancel_link``
        Link for the form's Cancel button to return to template list page

    **Template**

    :template:`report_template_form.html`
    """

    model = ReportTemplate
    form_class = ReportTemplateForm
    template_name = "reporting/report_template_form.html"

    def test_func(self):
        obj = self.get_object()
        if obj.protected:
            return verify_user_is_privileged(self.request.user)
        return self.request.user.is_active

    def handle_no_permission(self):
        obj = self.get_object()
        messages.error(self.request, "That template is protected – only an admin can edit it.")
        return HttpResponseRedirect(
            reverse(
                "reporting:template_detail",
                args=(obj.pk,),
            )
        )

    def get_context_data(self, **kwargs):
        ctx = super().get_context_data(**kwargs)
        ctx["cancel_link"] = reverse("reporting:templates")
        return ctx

    def get_success_url(self):
        messages.success(
            self.request,
            "Template successfully updated",
            extra_tags="alert-success",
        )
        return reverse("reporting:template_detail", kwargs={"pk": self.object.pk})

    def form_valid(self, form, **kwargs):
        obj = form.save(commit=False)
        obj.uploaded_by = self.request.user
        obj.save()
        form.save_m2m()
        return HttpResponseRedirect(self.get_success_url())

    def get_form_kwargs(self):
        kwargs = super().get_form_kwargs()
        kwargs.update({"user": self.request.user})
        return kwargs


class ReportTemplateDelete(RoleBasedAccessControlMixin, DeleteView):
    """
    Delete an individual instance of :model:`reporting.ReportTemplate`.

    **Context**

    ``object_type``
        String describing what is to be deleted
    ``object_to_be_deleted``
        To-be-deleted instance of :model:`reporting.ReportTemplate`
    ``cancel_link``
        Link for the form's Cancel button to return to template's detail page

    **Template**

    :template:`confirm_delete.html`
    """

    model = ReportTemplate
    template_name = "confirm_delete.html"

    def test_func(self):
        obj = self.get_object()
        if obj.protected:
            return verify_user_is_privileged(self.request.user)
        if obj.client:
            return verify_access(self.request.user, obj.client)
        return self.request.user.is_active

    def handle_no_permission(self):
        messages.error(self.request, "You do not have permission to access that.")
        return HttpResponseRedirect(
            reverse(
                "reporting:templates",
            )
        )

    def get_success_url(self):
        messages.success(
            self.request,
            "Successfully deleted the template and associated file.",
            extra_tags="alert-success",
        )
        return reverse("reporting:templates")

    def get_context_data(self, **kwargs):
        ctx = super().get_context_data(**kwargs)
        queryset = kwargs["object"]
        ctx["cancel_link"] = reverse("reporting:template_detail", kwargs={"pk": queryset.pk})
        ctx["object_type"] = "report template file (and associated file on disk)"
        ctx["object_to_be_deleted"] = queryset.filename
        return ctx


class ReportTemplateDownload(RoleBasedAccessControlMixin, SingleObjectMixin, View):
    """Return the target :model:`reporting.ReportTemplate` template file for download."""

    model = ReportTemplate

    def get(self, *args, **kwargs):
        obj = self.get_object()
        file_path = os.path.join(settings.MEDIA_ROOT, obj.document.path)
        if os.path.exists(file_path):
            return FileResponse(
                open(file_path, "rb"),
                as_attachment=True,
                filename=os.path.basename(file_path),
            )
        raise Http404


class GenerateReportJSON(RoleBasedAccessControlMixin, SingleObjectMixin, View):
    """Generate a JSON report for an individual :model:`reporting.Report`."""

    model = Report

    def test_func(self):
        return verify_access(self.request.user, self.get_object().project)

    def handle_no_permission(self):
        messages.error(self.request, "You do not have permission to access that.")
        return redirect("home:dashboard")

    def get(self, *args, **kwargs):
        obj = self.get_object()

        logger.info(
            "Generating JSON report for %s %s by request of %s",
            obj.__class__.__name__,
            obj.id,
            self.request.user,
        )

        engine = reportwriter.Reportwriter(obj, template_loc=None)
        json_report = engine.generate_json()

        return HttpResponse(json_report, "application/json")


class GenerateReportDOCX(RoleBasedAccessControlMixin, SingleObjectMixin, View):
    """Generate a DOCX report for an individual :model:`reporting.Report`."""

    model = Report

    def test_func(self):
        return verify_access(self.request.user, self.get_object().project)

    def handle_no_permission(self):
        messages.error(self.request, "You do not have permission to access that.")
        return redirect("home:dashboard")

    def get(self, *args, **kwargs):
        obj = self.get_object()

        logger.info(
            "Generating DOCX report for %s %s by request of %s",
            obj.__class__.__name__,
            obj.id,
            self.request.user,
        )

        try:
            report_name = generate_report_name(obj)

            # Get the template for this report
            if obj.docx_template:
                report_template = obj.docx_template
            else:
                report_config = ReportConfiguration.get_solo()
                report_template = report_config.default_docx_template
                if not report_template:
                    raise MissingTemplate
            template_loc = report_template.document.path

            # Check template's linting status
            template_status = report_template.get_status()
            if template_status in ("error", "failed"):
                messages.error(
                    self.request,
                    "The selected report template has linting errors and cannot be used to render a DOCX document",
                    extra_tags="alert-danger",
                )
                return HttpResponseRedirect(reverse("reporting:report_detail", kwargs={"pk": obj.pk}))

            # Template available and passes linting checks, so proceed with generation
            engine = reportwriter.Reportwriter(obj, template_loc)
            docx = engine.generate_word_docx()

            response = HttpResponse(
                content_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            )
            response["Content-Disposition"] = f'attachment; filename="{report_name}.docx"'
            docx.save(response)

            # Send WebSocket message to update user's webpage
            try:
                async_to_sync(channel_layer.group_send)(
                    "report_{}".format(obj.pk),
                    {
                        "type": "status_update",
                        "message": {"status": "success"},
                    },
                )
            except gaierror:
                # WebSocket are unavailable (unit testing)
                pass

            return response
        except ZeroDivisionError:
            logger.error(
                "DOCX generation failed for %s %s and user %s because of an attempt to divide by zero in Jinja2",
                obj.__class__.__name__,
                obj.id,
                self.request.user,
            )
            messages.info(
                self.request,
                "Tip: Before performing math, check if the number is greater than zero",
                extra_tags="alert-danger",
            )
            messages.error(
                self.request,
                "Word document generation failed because the selected template has Jinja2 code that attempts to divide by zero",
                extra_tags="alert-danger",
            )
        except MissingTemplate:
            logger.error(
                "DOCX generation failed for %s %s and user %s because no template was configured",
                obj.__class__.__name__,
                obj.id,
                self.request.user,
            )
            messages.error(
                self.request,
                "You do not have a Word template selected and have not configured a default template.",
                extra_tags="alert-danger",
            )
        except DocxPackageNotFoundError:
            logger.exception(
                "DOCX generation failed for %s %s and user %s because the template file was missing.",
                obj.__class__.__name__,
                obj.id,
                self.request.user,
            )
            messages.error(
                self.request,
                "Your selected Word template could not be found on the server – try uploading it again.",
                extra_tags="alert-danger",
            )
        except FileNotFoundError as error:
            logger.exception(
                "DOCX generation failed for %s %s and user %s because an evidence file was missing",
                obj.__class__.__name__,
                obj.id,
                self.request.user,
            )
            messages.error(
                self.request,
                "Halted document generation because an evidence file is missing: {}".format(error),
                extra_tags="alert-danger",
            )
        except UnrecognizedImageError as error:
            logger.exception(
                "DOCX generation failed for %s %s and user %s because of an unrecognized or corrupt image",
                obj.__class__.__name__,
                obj.id,
                self.request.user,
            )
            messages.error(
                self.request,
                "Encountered an error generating the document: {}".format(error).replace('"', "").replace("'", "`"),
                extra_tags="alert-danger",
            )
        except Exception as error:
            logger.exception(
                "DOCX generation failed unexpectedly for %s %s and user %s",
                obj.__class__.__name__,
                obj.id,
                self.request.user,
            )
            messages.error(
                self.request,
                "Encountered an error generating the document: {}".format(error).replace('"', "").replace("'", "`"),
                extra_tags="alert-danger",
            )

        return HttpResponseRedirect(reverse("reporting:report_detail", kwargs={"pk": obj.pk}))


class GenerateReportXLSX(RoleBasedAccessControlMixin, SingleObjectMixin, View):
    """Generate an XLSX report for an individual :model:`reporting.Report`."""

    model = Report

    def test_func(self):
        return verify_access(self.request.user, self.get_object().project)

    def handle_no_permission(self):
        messages.error(self.request, "You do not have permission to access that.")
        return redirect("home:dashboard")

    def get(self, *args, **kwargs):
        obj = self.get_object()

        logger.info(
            "Generating XLSX report for %s %s by request of %s",
            obj.__class__.__name__,
            obj.id,
            self.request.user,
        )

        try:
            report_name = generate_report_name(obj)
            engine = reportwriter.Reportwriter(obj, template_loc=None)

            output = io.BytesIO()
            engine.generate_excel_xlsx(output)
            output.seek(0)
            response = HttpResponse(
                output.read(),
                content_type="application/application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            )
            response["Content-Disposition"] = f'attachment; filename="{report_name}.xlsx"'
            output.close()

            return response
        except Exception as error:
            logger.exception(
                "XLSX generation failed unexpectedly for %s %s and user %s",
                obj.__class__.__name__,
                obj.id,
                self.request.user,
            )
            messages.error(
                self.request,
                "Encountered an error generating the spreadsheet: {}".format(error),
                extra_tags="alert-danger",
            )
        return HttpResponseRedirect(reverse("reporting:report_detail", kwargs={"pk": obj.pk}))


class GenerateReportPPTX(RoleBasedAccessControlMixin, SingleObjectMixin, View):
    """Generate a PPTX report for an individual :model:`reporting.Report`."""

    model = Report

    def test_func(self):
        return verify_access(self.request.user, self.get_object().project)

    def handle_no_permission(self):
        messages.error(self.request, "You do not have permission to access that.")
        return redirect("home:dashboard")

    def get(self, *args, **kwargs):
        obj = self.get_object()

        logger.info(
            "Generating PPTX report for %s %s by request of %s",
            obj.__class__.__name__,
            obj.id,
            self.request.user,
        )

        try:
            report_name = generate_report_name(obj)

            # Get the template for this report
            if obj.pptx_template:
                report_template = obj.pptx_template
            else:
                report_config = ReportConfiguration.get_solo()
                report_template = report_config.default_pptx_template
                if not report_template:
                    raise MissingTemplate
            template_loc = report_template.document.path

            # Check template's linting status
            template_status = report_template.get_status()
            if template_status in ("error", "failed"):
                messages.error(
                    self.request,
                    "The selected report template has linting errors and cannot be used to render a PPTX document.",
                    extra_tags="alert-danger",
                )
                return HttpResponseRedirect(reverse("reporting:report_detail", kwargs={"pk": obj.pk}))

            # Template available and passes linting checks, so proceed with generation
            engine = reportwriter.Reportwriter(obj, template_loc)
            pptx = engine.generate_powerpoint_pptx()
            response = HttpResponse(
                content_type="application/application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
            response["Content-Disposition"] = f'attachment; filename="{report_name}.pptx"'
            pptx.save(response)

            return response
        except MissingTemplate:
            logger.error(
                "PPTX generation failed for %s %s and user %s because no template was configured",
                obj.__class__.__name__,
                obj.id,
                self.request.user,
            )
            messages.error(
                self.request,
                "You do not have a PowerPoint template selected and have not configured a default template.",
                extra_tags="alert-danger",
            )
        except ValueError as exception:
            logger.exception(
                "PPTX generation failed for %s %s and user %s because the template could not be loaded as a PPTX",
                obj.__class__.__name__,
                obj.id,
                self.request.user,
            )
            messages.error(
                self.request,
                f"Your selected template could not be loaded as a PowerPoint template: {exception}",
                extra_tags="alert-danger",
            )
        except PptxPackageNotFoundError:
            logger.exception(
                "PPTX generation failed for %s %s and user %s because the template file was missing",
                obj.__class__.__name__,
                obj.id,
                self.request.user,
            )
            messages.error(
                self.request,
                "Your selected PowerPoint template could not be found on the server – try uploading it again.",
                extra_tags="alert-danger",
            )
        except FileNotFoundError as error:
            logger.exception(
                "PPTX generation failed for %s %s and user %s because an evidence file was missing",
                obj.__class__.__name__,
                obj.id,
                self.request.user,
            )
            messages.error(
                self.request,
                "Halted document generation because an evidence file is missing: {}".format(error),
                extra_tags="alert-danger",
            )
        except UnrecognizedImageError as error:
            logger.exception(
                "PPTX generation failed for %s %s and user %s because of an unrecognized or corrupt image",
                obj.__class__.__name__,
                obj.id,
                self.request.user,
            )
            messages.error(
                self.request,
                "Encountered an error generating the document: {}".format(error).replace('"', "").replace("'", "`"),
                extra_tags="alert-danger",
            )
        except Exception as error:
            logger.exception(
                "PPTX generation failed unexpectedly for %s %s and user %s",
                obj.__class__.__name__,
                obj.id,
                self.request.user,
            )
            messages.error(
                self.request,
                "Encountered an error generating the document: {}".format(error).replace('"', "").replace("'", "`"),
                extra_tags="alert-danger",
            )

        return HttpResponseRedirect(reverse("reporting:report_detail", kwargs={"pk": obj.pk}))


class GenerateReportAll(RoleBasedAccessControlMixin, SingleObjectMixin, View):
    """Generate all report types for an individual :model:`reporting.Report`."""

    model = Report

    def test_func(self):
        return verify_access(self.request.user, self.get_object().project)

    def handle_no_permission(self):
        messages.error(self.request, "You do not have permission to access that.")
        return redirect("home:dashboard")

    def get(self, *args, **kwargs):
        obj = self.get_object()

        logger.info(
            "Generating PPTX report for %s %s by request of %s",
            obj.__class__.__name__,
            obj.id,
            self.request.user,
        )

        try:
            report_name = generate_report_name(obj)
            engine = reportwriter.Reportwriter(obj, template_loc=None)

            # Get the templates for Word and PowerPoint
            if obj.docx_template:
                docx_template = obj.docx_template
            else:
                report_config = ReportConfiguration.get_solo()
                docx_template = report_config.default_docx_template
                if not docx_template:
                    raise MissingTemplate
            docx_template = docx_template.document.path

            if obj.pptx_template:
                pptx_template = obj.pptx_template
            else:
                report_config = ReportConfiguration.get_solo()
                pptx_template = report_config.default_pptx_template
                if not pptx_template:
                    raise MissingTemplate
            pptx_template = pptx_template.document.path

            # Generate all types of reports
            json_doc, docx_doc, xlsx_doc, pptx_doc = engine.generate_all_reports(docx_template, pptx_template)

            # Create a zip file in memory and add the reports to it
            zip_buffer = io.BytesIO()
            with zipfile.ZipFile(zip_buffer, "a") as zf:
                zf.writestr(f"{report_name}.json", json_doc)
                zf.writestr(f"{report_name}.docx", docx_doc.getvalue())
                zf.writestr(f"{report_name}.xlsx", xlsx_doc.getvalue())
                zf.writestr(f"{report_name}.pptx", pptx_doc.getvalue())
            zip_buffer.seek(0)

            # Return the buffer in the HTTP response
            response = HttpResponse(content_type="application/x-zip-compressed")
            response["Content-Disposition"] = f'attachment; filename="{report_name}.zip"'
            response.write(zip_buffer.read())

            return response
        except MissingTemplate:
            messages.error(
                self.request,
                "You do not have a PowerPoint template selected and have not configured a default template.",
                extra_tags="alert-danger",
            )
        except ValueError as exception:
            messages.error(
                self.request,
                f"Your selected template could not be loaded as a PowerPoint template: {exception}",
                extra_tags="alert-danger",
            )
        except DocxPackageNotFoundError:
            messages.error(
                self.request,
                "Your selected Word template could not be found on the server – try uploading it again.",
                extra_tags="alert-danger",
            )
        except PptxPackageNotFoundError:
            messages.error(
                self.request,
                "Your selected PowerPoint template could not be found on the server – try uploading it again.",
                extra_tags="alert-danger",
            )
        except Exception as error:
            messages.error(
                self.request,
                "Encountered an error generating the document: {}".format(error),
                extra_tags="alert-danger",
            )

        return HttpResponseRedirect(reverse("reporting:report_detail", kwargs={"pk": obj.pk}))


# CBVs related to :model:`reporting.ReportFindingLink`


class ReportFindingLinkUpdate(RoleBasedAccessControlMixin, UpdateView):
    """
    Update an individual instance of :model:`reporting.ReportFindingLink`.

    **Context**

    ``cancel_link``
        Link for the form's Cancel button to return to report's detail page

    **Template**

    :template:`reporting/local_edit.html.html`
    """

    model = ReportFindingLink
    form_class = ReportFindingLinkUpdateForm
    template_name = "reporting/local_edit.html"
    success_url = reverse_lazy("reporting:reports")

    def test_func(self):
        return verify_access(self.request.user, self.get_object().report.project)

    def handle_no_permission(self):
        messages.error(self.request, "You do not have permission to access that.")
        return redirect("home:dashboard")

    def get_context_data(self, **kwargs):
        ctx = super().get_context_data(**kwargs)
        ctx["cancel_link"] = reverse("reporting:report_detail", kwargs={"pk": self.object.report.pk})
        return ctx

    def form_valid(self, form):
        if form.changed_data:
            changed_at = dateformat.format(timezone.now(), "H:i:s e")
            async_to_sync(channel_layer.group_send)(
                "finding_{}".format(self.object.id),
                {
                    "type": "message",
                    "message": {
                        "message": f"User {self.request.user.username} updated this finding at {changed_at}.",
                        "level": "warning",
                        "title": "Content Has Changed",
                    },
                },
            )

        # Send Websockets messages if assignment changed
        if "assigned_to" in form.changed_data:
            # Get the entries current values (those being changed)
            old_entry = ReportFindingLink.objects.get(pk=self.object.pk)
            old_assignee = old_entry.assigned_to
            # Notify new assignee over WebSockets
            if "assigned_to" in form.changed_data:
                # Only notify if the assignee is not the user who made the change
                if self.request.user != self.object.assigned_to:
                    try:
                        # Send a message to the assigned user
                        async_to_sync(channel_layer.group_send)(
                            "notify_{}".format(self.object.assigned_to),
                            {
                                "type": "message",
                                "message": {
                                    "message": "You have been assigned to this finding for {}:\n{}".format(
                                        self.object.report, self.object.title
                                    ),
                                    "level": "info",
                                    "title": "New Assignment",
                                },
                            },
                        )
                    except gaierror:
                        # WebSocket are unavailable (unit testing)
                        pass
                if self.request.user != old_assignee:
                    try:
                        # Send a message to the unassigned user
                        async_to_sync(channel_layer.group_send)(
                            "notify_{}".format(old_assignee),
                            {
                                "type": "message",
                                "message": {
                                    "message": "You have been unassigned from this finding for {}:\n{}".format(
                                        self.object.report, self.object.title
                                    ),
                                    "level": "info",
                                    "title": "Assignment Change",
                                },
                            },
                        )
                    except gaierror:
                        # WebSocket are unavailable (unit testing)
                        pass
        return super().form_valid(form)

    def get_form(self, form_class=None):
        form = super().get_form(form_class)
        user_primary_keys = ProjectAssignment.objects.filter(project=self.object.report.project).values_list(
            "operator", flat=True
        )
        form.fields["assigned_to"].queryset = User.objects.filter(id__in=user_primary_keys)
        return form

    def get_success_url(self):
        messages.success(
            self.request,
            "Successfully updated {}.".format(self.get_object().title),
            extra_tags="alert-success",
        )
        return reverse("reporting:report_detail", kwargs={"pk": self.object.report.id})


# CBVs related to :model:`reporting.Evidence`


class EvidenceDetailView(RoleBasedAccessControlMixin, DetailView):
    """
    Display an individual instance of :model:`reporting.Evidence`.

    **Template**

    :template:`reporting/evidence_detail.html`
    """

    model = Evidence

    def test_func(self):
        return verify_access(self.request.user, self.get_object().finding.report.project)

    def handle_no_permission(self):
        messages.error(self.request, "You do not have permission to access that.")
        return redirect("home:dashboard")

    def get_context_data(self, **kwargs):
        ctx = super().get_context_data(**kwargs)
        file_content = None
        if os.path.isfile(self.object.document.path):
            if (
                self.object.document.name.lower().endswith(".txt")
                or self.object.document.name.lower().endswith(".log")
                or self.object.document.name.lower().endswith(".md")
            ):
                filetype = "text"
                file_content = []
                temp = self.object.document.read().splitlines()
                for line in temp:
                    try:
                        file_content.append(line.decode("utf-8", errors="replace"))
                    except UnicodeError:
                        file_content.append(line)
            elif (
                self.object.document.name.lower().endswith(".jpg")
                or self.object.document.name.lower().endswith(".png")
                or self.object.document.name.lower().endswith(".jpeg")
            ):
                filetype = "image"
            else:
                filetype = "unknown"
        else:
            filetype = "text"
            file_content = ["FILE NOT FOUND"]

        ctx["filetype"] = filetype
        ctx["evidence"] = self.object
        ctx["file_content"] = file_content

        return ctx


class EvidenceCreate(RoleBasedAccessControlMixin, CreateView):
    """
    Create an individual :model:`reporting.Evidence` entry linked to an individual
    :model:`reporting.ReportFindingLink`.

    **Template**

    :template:`reporting/evidence_form.html`
    """

    model = Evidence
    form_class = EvidenceForm

    def test_func(self):
        return verify_access(self.request.user, self.finding_instance.report.project)

    def handle_no_permission(self):
        messages.error(self.request, "You do not have permission to access that.")
        return redirect("home:dashboard")

    def setup(self, request, *args, **kwargs):
        super().setup(request, *args, **kwargs)
        finding_pk = self.kwargs.get("pk")
        self.finding_instance = get_object_or_404(ReportFindingLink, pk=finding_pk)
        self.evidence_queryset = Evidence.objects.filter(finding=self.finding_instance.pk)

    def get_template_names(self):
        if "modal" in self.kwargs:
            modal = self.kwargs["modal"]
            if modal:
                return ["reporting/evidence_form_modal.html"]
            return ["reporting/evidence_form.html"]
        return ["reporting/evidence_form.html"]

    def get_form_kwargs(self):
        kwargs = super().get_form_kwargs()
        kwargs.update({"evidence_queryset": self.evidence_queryset})
        if "modal" in self.kwargs:
            kwargs.update({"is_modal": True})
        return kwargs

    def get_context_data(self, **kwargs):
        ctx = super().get_context_data(**kwargs)
        ctx["cancel_link"] = reverse("reporting:report_detail", kwargs={"pk": self.finding_instance.report.pk})
        if "modal" in self.kwargs:
            friendly_names = self.evidence_queryset.values_list("friendly_name", flat=True)
            used_friendly_names = []
            # Convert the queryset into a list to pass to JavaScript later
            for name in friendly_names:
                used_friendly_names.append(name)
            ctx["used_friendly_names"] = used_friendly_names

        return ctx

    def form_valid(self, form, **kwargs):
        obj = form.save(commit=False)
        obj.uploaded_by = self.request.user
        obj.finding = self.finding_instance
        obj.save()
        form.save_m2m()
        if os.path.isfile(obj.document.path):
            messages.success(
                self.request,
                "Evidence uploaded successfully.",
                extra_tags="alert-success",
            )
        else:
            messages.error(
                self.request,
                "Evidence file failed to upload!",
                extra_tags="alert-danger",
            )
        return HttpResponseRedirect(self.get_success_url())

    def get_success_url(self):
        if "modal" in self.kwargs:
            return reverse("reporting:upload_evidence_modal_success")
        return reverse("reporting:report_detail", args=(self.finding_instance.report.pk,))


class EvidenceUpdate(RoleBasedAccessControlMixin, UpdateView):
    """
    Update an individual instance of :model:`reporting.Evidence`.

    **Context**

    ``cancel_link``
        Link for the form's Cancel button to return to evidence's detail page

    **Template**

    :template:`reporting/evidence_form.html`
    """

    model = Evidence
    form_class = EvidenceForm

    def test_func(self):
        return verify_access(self.request.user, self.get_object().finding.report.project)

    def handle_no_permission(self):
        messages.error(self.request, "You do not have permission to access that.")
        return redirect("home:dashboard")

    def get_form_kwargs(self):
        kwargs = super().get_form_kwargs()
        evidence_queryset = Evidence.objects.filter(finding=self.object.finding.pk)
        kwargs.update({"evidence_queryset": evidence_queryset})
        return kwargs

    def get_context_data(self, **kwargs):
        ctx = super().get_context_data(**kwargs)
        ctx["cancel_link"] = reverse(
            "reporting:evidence_detail",
            kwargs={"pk": self.object.pk},
        )
        return ctx

    def get_success_url(self):
        messages.success(
            self.request,
            "Successfully updated {}.".format(self.get_object().friendly_name),
            extra_tags="alert-success",
        )
        return reverse("reporting:evidence_detail", kwargs={"pk": self.object.pk})


class EvidenceDelete(RoleBasedAccessControlMixin, DeleteView):
    """
    Delete an individual instance of :model:`reporting.Evidence`.

    **Context**

    ``object_type``
        String describing what is to be deleted
    ``object_to_be_deleted``
        To-be-deleted instance of :model:`reporting.Evidence`
    ``cancel_link``
        Link for the form's Cancel button to return to evidence's detail page

    **Template**

    :template:`confirm_delete.html`
    """

    model = Evidence
    template_name = "confirm_delete.html"

    def test_func(self):
        return verify_access(self.request.user, self.get_object().finding.report.project)

    def handle_no_permission(self):
        messages.error(self.request, "You do not have permission to access that.")
        return redirect("home:dashboard")

    def get_success_url(self):
        message = "Successfully deleted the evidence and associated file."
        if os.path.isfile(self.object.document.name):
            message = "Successfully deleted the evidence, but could not delete the associated file."
        messages.success(
            self.request,
            message,
            extra_tags="alert-success",
        )
        return reverse("reporting:report_detail", kwargs={"pk": self.object.finding.report.pk})

    def get_context_data(self, **kwargs):
        ctx = super().get_context_data(**kwargs)
        queryset = kwargs["object"]
        ctx["cancel_link"] = reverse("reporting:evidence_detail", kwargs={"pk": queryset.pk})
        ctx["object_type"] = "evidence file (and associated file on disk)"
        ctx["object_to_be_deleted"] = queryset.friendly_name
        return ctx


# CBVs related to :model:`reporting.Finding`


class FindingNoteCreate(RoleBasedAccessControlMixin, CreateView):
    """
    Create an individual instance of :model:`reporting.FindingNote`.

    **Context**

    ``cancel_link``
        Link for the form's Cancel button to return to finding's detail page

    **Template**

    :template:`note_form.html`
    """

    model = FindingNote
    form_class = FindingNoteForm
    template_name = "note_form.html"

    def get_context_data(self, **kwargs):
        ctx = super().get_context_data(**kwargs)
        finding_instance = get_object_or_404(Finding, pk=self.kwargs.get("pk"))
        ctx["cancel_link"] = reverse("reporting:finding_detail", kwargs={"pk": finding_instance.pk})
        return ctx

    def get_success_url(self):
        messages.success(
            self.request,
            "Successfully added your note to this finding.",
            extra_tags="alert-success",
        )
        return "{}#notes".format(reverse("reporting:finding_detail", kwargs={"pk": self.object.finding.id}))

    def form_valid(self, form, **kwargs):
        self.object = form.save(commit=False)
        self.object.operator = self.request.user
        self.object.finding_id = self.kwargs.get("pk")
        self.object.save()
        return super().form_valid(form)


class FindingNoteUpdate(RoleBasedAccessControlMixin, UpdateView):
    """
    Update an individual instance of :model:`reporting.FindingNote`.

    **Context**

    ``cancel_link``
        Link for the form's Cancel button to return to finding's detail page

    **Template**

    :template:`note_form.html`
    """

    model = FindingNote
    form_class = FindingNoteForm
    template_name = "note_form.html"

    def test_func(self):
        return self.get_object().operator.id == self.request.user.id

    def handle_no_permission(self):
        messages.error(self.request, "You do not have permission to access that.")
        return redirect("home:dashboard")

    def get_context_data(self, **kwargs):
        ctx = super().get_context_data(**kwargs)
        ctx["cancel_link"] = reverse("reporting:finding_detail", kwargs={"pk": self.get_object().finding.pk})
        return ctx

    def get_success_url(self):
        messages.success(self.request, "Successfully updated the note.", extra_tags="alert-success")
        return reverse("reporting:finding_detail", kwargs={"pk": self.get_object().finding.pk})


# CBVs related to :model:`reporting.LocalFindingNote`


class LocalFindingNoteCreate(RoleBasedAccessControlMixin, CreateView):
    """
    Create an individual instance of :model:`reporting.LocalFindingNote`.

    **Context**

    ``cancel_link``
        Link for the form's Cancel button to return to finding's detail page

    **Template**

    :template:`note_form.html`
    """

    model = LocalFindingNote
    form_class = LocalFindingNoteForm
    template_name = "note_form.html"

    def setup(self, request, *args, **kwargs):
        super().setup(request, *args, **kwargs)
        self.finding_instance = get_object_or_404(ReportFindingLink, pk=self.kwargs.get("pk"))

    def test_func(self):
        return verify_access(self.request.user, self.finding_instance.report.project)

    def handle_no_permission(self):
        messages.error(self.request, "You do not have permission to access that.")
        return redirect("home:dashboard")

    def get_context_data(self, **kwargs):
        ctx = super().get_context_data(**kwargs)
        ctx["cancel_link"] = reverse("reporting:local_edit", kwargs={"pk": self.finding_instance.pk})
        return ctx

    def get_success_url(self):
        messages.success(
            self.request,
            "Successfully added your note to this finding.",
            extra_tags="alert-success",
        )
        return reverse("reporting:local_edit", kwargs={"pk": self.object.finding.pk})

    def form_valid(self, form, **kwargs):
        self.object = form.save(commit=False)
        self.object.operator = self.request.user
        self.object.finding_id = self.kwargs.get("pk")
        self.object.save()
        return super().form_valid(form)


class LocalFindingNoteUpdate(RoleBasedAccessControlMixin, UpdateView):
    """
    Update an individual instance of :model:`reporting.LocalFindingNote`.

    **Context**

    ``cancel_link``
        Link for the form's Cancel button to return to finding's detail page

    **Template**

    :template:`note_form.html`
    """

    model = LocalFindingNote
    form_class = LocalFindingNoteForm
    template_name = "note_form.html"

    def test_func(self):
        return self.get_object().operator.id == self.request.user.id

    def handle_no_permission(self):
        messages.error(self.request, "You do not have permission to access that.")
        return redirect("home:dashboard")

    def get_context_data(self, **kwargs):
        ctx = super().get_context_data(**kwargs)
        note_instance = get_object_or_404(LocalFindingNote, pk=self.kwargs.get("pk"))
        ctx["cancel_link"] = reverse("reporting:local_edit", kwargs={"pk": note_instance.finding.id})
        return ctx

    def get_success_url(self):
        messages.success(self.request, "Successfully updated the note.", extra_tags="alert-success")
        return reverse("reporting:local_edit", kwargs={"pk": self.get_object().finding.pk})
