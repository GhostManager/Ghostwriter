import pptx
from io import BytesIO
from zipfile import ZipFile

from django.test import TestCase
from .test_rich_text_docx import clean_xml

from ghostwriter.modules.reportwriter.richtext.pptx import HtmlToPptx

PPTX_PREFIX = """<?xml version='1.0' encoding='UTF-8' standalone='yes'?>
<p:sld xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
    xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
    xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
    <p:cSld>
        <p:spTree>
            <p:nvGrpSpPr>
                <p:cNvPr id="1" name=""/>
                <p:cNvGrpSpPr/>
                <p:nvPr/>
            </p:nvGrpSpPr>
            <p:grpSpPr/>
            <p:sp>
                <p:nvSpPr>
                    <p:cNvPr id="2" name="Title 1"/>
                    <p:cNvSpPr>
                        <a:spLocks noGrp="1"/>
                    </p:cNvSpPr>
                    <p:nvPr>
                        <p:ph type="title"/>
                    </p:nvPr>
                </p:nvSpPr>
                <p:spPr/>
                <p:txBody>
                    <a:bodyPr/>
                    <a:lstStyle/>
                    <a:p/>
                </p:txBody>
            </p:sp>
            <p:sp>
                <p:nvSpPr>
                    <p:cNvPr id="3" name="Content Placeholder 2"/>
                    <p:cNvSpPr>
                        <a:spLocks noGrp="1"/>
                    </p:cNvSpPr>
                    <p:nvPr>
                        <p:ph idx="1"/>
                    </p:nvPr>
                </p:nvSpPr>
                <p:spPr/>
                <p:txBody>
                    <a:bodyPr/>
                    <a:lstStyle/>
"""
PPTX_SUFFIX = """
                </p:txBody>
            </p:sp>
        </p:spTree>
    </p:cSld>
    <p:clrMapOvr>
        <a:masterClrMapping/>
    </p:clrMapOvr>
</p:sld>
"""
SLD_LAYOUT_TITLE_AND_CONTENT = 1


def mk_test_pptx(name, input, expected_output, add_suffix=True):
    expected_output = clean_xml(PPTX_PREFIX + expected_output + (PPTX_SUFFIX if add_suffix else ""))

    def test_func(self):
        ppt = pptx.Presentation()
        slide = ppt.slides.add_slide(ppt.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT])
        shape = slide.shapes.placeholders[1]
        shape.text_frame.clear()
        HtmlToPptx.run(input, slide, shape)
        HtmlToPptx.delete_extra_paragraph(shape)

        out = BytesIO()
        ppt.part.save(out)

        # Uncomment to write generates pptx files for manual inspection
        # with open(name + ".pptx", "wb") as f:
        #     f.write(out.getvalue())

        with ZipFile(out) as zip:
            with zip.open("ppt/slides/slide1.xml") as file:
                contents = file.read()
        contents = clean_xml(contents)
        self.assertEqual(contents, expected_output)

    test_func.__name__ = name
    return test_func


class RichTextToPptxTests(TestCase):
    maxDiff = None

    test_paragraphs = mk_test_pptx(
        "test_paragraphs",
        "<p>Hello World!</p><p>This is a test!</p>",
        """
            <a:p><a:r><a:t>Hello World!</a:t></a:r></a:p>
            <a:p><a:r><a:t>This is a test!</a:t></a:r></a:p>
        """,
    )

    test_headers = mk_test_pptx(
        "test_headers",
        "<h1>Hello World!</h1><h2>Heading two</h2><h3>Heading three</h3><p>Paragraph</p>",
        """
            <a:p><a:r><a:rPr b="1" /><a:t>Hello World!</a:t></a:r></a:p>
            <a:p><a:r><a:rPr b="1" /><a:t>Heading two</a:t></a:r></a:p>
            <a:p><a:r><a:rPr b="1" /><a:t>Heading three</a:t></a:r></a:p>
            <a:p><a:r><a:t>Paragraph</a:t></a:r></a:p>
        """,
    )

    test_formatting = mk_test_pptx(
        "test_formatting",
        """
            <p>
                <b>Bold</b>
                <strong>Strong</strong>
                <i>Italic</i>
                <em>Emphasis</em>
                <u>Underline</u>
                <sub>Subscript</sub>
                <sup>Superscript</sup>
                <del>Strikethrough</del>
            </p>
        """,
        """
            <a:p>
                <a:r><a:t/> </a:r>
                <a:r><a:rPr b="1" /><a:t>Bold</a:t></a:r>
                <a:r><a:t> </a:t></a:r>
                <a:r><a:rPr b="1" /><a:t>Strong</a:t></a:r>
                <a:r><a:t> </a:t></a:r>
                <a:r><a:rPr i="1" /><a:t>Italic</a:t></a:r>
                <a:r><a:t> </a:t></a:r>
                <a:r><a:rPr i="1" /><a:t>Emphasis</a:t></a:r>
                <a:r><a:t> </a:t></a:r>
                <a:r><a:rPr u="sng" /><a:t>Underline</a:t></a:r>
                <a:r><a:t> </a:t></a:r>
                <a:r><a:rPr baseline="-25000" /><a:t>Subscript</a:t></a:r>
                <a:r><a:t> </a:t></a:r>
                <a:r><a:rPr baseline="30000" /><a:t>Superscript</a:t></a:r>
                <a:r><a:t> </a:t></a:r>
                <a:r><a:rPr strike="sngStrike" /><a:t>Strikethrough</a:t></a:r>
                <a:r><a:t/> </a:r>
            </a:p>
        """,
    )

    test_unordered_list = mk_test_pptx(
        "test_unordered_list",
        """
            <p>List test one:</p>
            <ul>
                <li>Item one</li>
                <li>Item two</li>
                <li>Item three:<ul>
                    <li>Subitem one</li>
                    <li>Subitem two</li>
                </ul></li>
            </ul>
            <p>List test two:</p>
            <ul>
                <li>Item one</li>
                <li>Item two</li>
            </ul>
        """,
        """
            <a:p><a:r><a:t>List test one:</a:t></a:r></a:p>
            <a:p><a:pPr lvl="1"/><a:r><a:t>Item one</a:t></a:r></a:p>
            <a:p><a:pPr lvl="1"/><a:r><a:t>Item two</a:t></a:r></a:p>
            <a:p><a:pPr lvl="1"/><a:r><a:t>Item three:</a:t></a:r></a:p>
            <a:p><a:pPr lvl="2"/><a:r><a:t>Subitem one</a:t></a:r></a:p>
            <a:p><a:pPr lvl="2"/><a:r><a:t>Subitem two</a:t></a:r></a:p>
            <a:p><a:r><a:t>List test two:</a:t></a:r></a:p>
            <a:p><a:pPr lvl="1"/><a:r><a:t>Item one</a:t></a:r></a:p>
            <a:p><a:pPr lvl="1"/><a:r><a:t>Item two</a:t></a:r></a:p>
        """,
    )

    test_ordered_list = mk_test_pptx(
        "test_ordered_list",
        """
            <p>List test one:</p>
            <ol>
                <li>Item one</li>
                <li>Item two</li>
                <li>Item three:<ol>
                    <li>Subitem one</li>
                    <li>Subitem two</li>
                </ol></li>
            </ol>
            <p>List test two:</p>
            <ol>
                <li>Item one</li>
                <li>Item two</li>
            </ol>
        """,
        """
            <a:p><a:r><a:t>List test one:</a:t></a:r></a:p>
            <a:p><a:pPr lvl="1"/><a:r><a:t>Item one</a:t></a:r></a:p>
            <a:p><a:pPr lvl="1"/><a:r><a:t>Item two</a:t></a:r></a:p>
            <a:p><a:pPr lvl="1"/><a:r><a:t>Item three:</a:t></a:r></a:p>
            <a:p><a:pPr lvl="2"/><a:r><a:t>Subitem one</a:t></a:r></a:p>
            <a:p><a:pPr lvl="2"/><a:r><a:t>Subitem two</a:t></a:r></a:p>
            <a:p><a:r><a:t>List test two:</a:t></a:r></a:p>
            <a:p><a:pPr lvl="1"/><a:r><a:t>Item one</a:t></a:r></a:p>
            <a:p><a:pPr lvl="1"/><a:r><a:t>Item two</a:t></a:r></a:p>
        """,
    )

    test_mixed_list = mk_test_pptx(
        "test_mixed_list",
        """
            <p>List test one:</p>
            <ul>
                <li>Item one</li>
                <li>Item two</li>
                <li>Item three:<ol>
                    <li>Subitem one</li>
                    <li>Subitem two</li>
                </ol></li>
            </ul>
            <p>List test two:</p>
            <ol>
                <li>Item one</li>
                <li>Item two</li>
            </ol>
        """,
        """
            <a:p><a:r><a:t>List test one:</a:t></a:r></a:p>
            <a:p><a:pPr lvl="1"/><a:r><a:t>Item one</a:t></a:r></a:p>
            <a:p><a:pPr lvl="1"/><a:r><a:t>Item two</a:t></a:r></a:p>
            <a:p><a:pPr lvl="1"/><a:r><a:t>Item three:</a:t></a:r></a:p>
            <a:p><a:pPr lvl="2"/><a:r><a:t>Subitem one</a:t></a:r></a:p>
            <a:p><a:pPr lvl="2"/><a:r><a:t>Subitem two</a:t></a:r></a:p>
            <a:p><a:r><a:t>List test two:</a:t></a:r></a:p>
            <a:p><a:pPr lvl="1"/><a:r><a:t>Item one</a:t></a:r></a:p>
            <a:p><a:pPr lvl="1"/><a:r><a:t>Item two</a:t></a:r></a:p>
        """,
    )

    test_blockquote = mk_test_pptx(
        "test_blockquote",
        """
        <p>A quote:</p>
        <blockquote>Alas poor <b>Yorik</b>, I knew him well</blockquote>
        """,
        """
            <a:p><a:r><a:t>A quote:</a:t></a:r></a:p>
            <a:p>
                <a:r><a:t>Alas poor </a:t></a:r>
                <a:r><a:rPr b="1" /><a:t>Yorik</a:t></a:r>
                <a:r><a:t>, I knew him well</a:t></a:r>
            </a:p>
        """,
    )

    test_pre = mk_test_pptx(
        "test_pre",
        """
        <p>Some code:</p>
        <pre>int main() {
    printf("hello world!\\n");
}</pre>
        """,
        """
            <a:p><a:r><a:t>Some code:</a:t></a:r></a:p>
            </p:txBody></p:sp>
            <p:sp>
                <p:nvSpPr>
                    <p:cNvPr id="4" name="TextBox 3"/>
                    <p:cNvSpPr txBox="1"/>
                    <p:nvPr/>
                </p:nvSpPr>
                <p:spPr>
                    <a:xfrm>
                        <a:off x="7315200" y="1508760"/>
                        <a:ext cx="4114800" cy="2743200"/>
                    </a:xfrm>
                    <a:prstGeom prst="rect">
                        <a:avLst/>
                    </a:prstGeom>
                    <a:noFill/>
                </p:spPr>
                <p:txBody>
                    <a:bodyPr wrap="none">
                        <a:spAutoFit/>
                    </a:bodyPr>
                    <a:lstStyle/>
                    <a:p/>
                    <a:p>
                        <a:pPr algn="l"/>
                        <a:r>
                            <a:rPr sz="1100">
                                <a:latin typeface="Courier New"/>
                            </a:rPr>
                            <a:t>int main() {
    printf("hello world!\\n");
}</a:t>
                        </a:r>
                    </a:p>
        """,
    )

    test_table = mk_test_pptx(
        "test_table",
        """
        <table>
            <tbody>
                <tr>
                    <td>Cell one</td>
                    <td>Cell two</td>
                    <td>Cell three</td>
                </tr>
                <tr>
                    <td>Cell four</td>
                    <td>Cell five</td>
                    <td>Cell six</td>
                </tr>
            </tbody>
        </table>
        """,
        """
                    <a:p/>
                </p:txBody>
            </p:sp>
            <p:graphicFrame>
                <p:nvGraphicFramePr>
                    <p:cNvPr id="4" name="Table 3"/>
                    <p:cNvGraphicFramePr>
                        <a:graphicFrameLocks noGrp="1"/>
                    </p:cNvGraphicFramePr>
                    <p:nvPr/>
                </p:nvGraphicFramePr>
                <p:xfrm>
                    <a:off x="9144000" y="4572000"/>
                    <a:ext cx="2743200" cy="914400"/>
                </p:xfrm>
                <a:graphic>
                    <a:graphicData uri="http://schemas.openxmlformats.org/drawingml/2006/table">
                        <a:tbl>
                            <a:tblPr firstRow="1" bandRow="1">
                                <a:tableStyleId>{5C22544A-7EE6-4342-B048-85BDC9FD1C3A}</a:tableStyleId>
                            </a:tblPr>
                            <a:tblGrid>
                                <a:gridCol w="914400"/>
                                <a:gridCol w="914400"/>
                                <a:gridCol w="914400"/>
                            </a:tblGrid>
                            <a:tr h="457200">
                                <a:tc>
                                    <a:txBody>
                                        <a:bodyPr/>
                                        <a:lstStyle/>
                                        <a:p><a:r><a:t>Cell one</a:t></a:r></a:p>
                                    </a:txBody>
                                    <a:tcPr/>
                                </a:tc>
                                <a:tc>
                                    <a:txBody>
                                        <a:bodyPr/>
                                        <a:lstStyle/>
                                        <a:p><a:r><a:t>Cell two</a:t></a:r></a:p>
                                    </a:txBody>
                                    <a:tcPr/>
                                </a:tc>
                                <a:tc>
                                    <a:txBody>
                                        <a:bodyPr/>
                                        <a:lstStyle/>
                                        <a:p><a:r><a:t>Cell three</a:t></a:r></a:p>
                                    </a:txBody>
                                    <a:tcPr/>
                                </a:tc>
                            </a:tr>
                            <a:tr h="457200">
                                <a:tc>
                                    <a:txBody>
                                        <a:bodyPr/>
                                        <a:lstStyle/>
                                        <a:p><a:r><a:t>Cell four</a:t></a:r></a:p>
                                    </a:txBody>
                                    <a:tcPr/>
                                </a:tc>
                                <a:tc>
                                    <a:txBody>
                                        <a:bodyPr/>
                                        <a:lstStyle/>
                                        <a:p><a:r><a:t>Cell five</a:t></a:r></a:p>
                                    </a:txBody>
                                    <a:tcPr/>
                                </a:tc>
                                <a:tc>
                                    <a:txBody>
                                        <a:bodyPr/>
                                        <a:lstStyle/>
                                        <a:p><a:r><a:t>Cell six</a:t></a:r></a:p>
                                    </a:txBody>
                                    <a:tcPr/>
                                </a:tc>
                            </a:tr>
                        </a:tbl>
                    </a:graphicData>
                </a:graphic>
            </p:graphicFrame>
        </p:spTree>
    </p:cSld>
    <p:clrMapOvr>
        <a:masterClrMapping/>
    </p:clrMapOvr>
</p:sld>
        """,
        add_suffix=False,
    )

    test_table_spans = mk_test_pptx(
        "test_table_spans",
        """
        <table>
            <tbody>
                <tr>
                    <td>Cell one</td>
                    <td colspan="2">Wide cell</td>
                </tr>
                <tr>
                    <td rowspan="2" colspan = "2">Big cell</td>
                    <td>Cell two</td>
                </tr>
                <tr>
                    <td>Cell three</td>
                </tr>
            </tbody>
        </table>
        """,
        """
<a:p/>
                </p:txBody>
            </p:sp>
            <p:graphicFrame>
                <p:nvGraphicFramePr>
                    <p:cNvPr id="4" name="Table 3"/>
                    <p:cNvGraphicFramePr>
                        <a:graphicFrameLocks noGrp="1"/>
                    </p:cNvGraphicFramePr>
                    <p:nvPr/>
                </p:nvGraphicFramePr>
                <p:xfrm>
                    <a:off x="9144000" y="4572000"/>
                    <a:ext cx="2743200" cy="914400"/>
                </p:xfrm>
                <a:graphic>
                    <a:graphicData uri="http://schemas.openxmlformats.org/drawingml/2006/table">
                        <a:tbl>
                            <a:tblPr firstRow="1" bandRow="1">
                                <a:tableStyleId>{5C22544A-7EE6-4342-B048-85BDC9FD1C3A}</a:tableStyleId>
                            </a:tblPr>
                            <a:tblGrid>
                                <a:gridCol w="914400"/>
                                <a:gridCol w="914400"/>
                                <a:gridCol w="914400"/>
                            </a:tblGrid>
                            <a:tr h="304800">
                                <a:tc>
                                    <a:txBody>
                                        <a:bodyPr/>
                                        <a:lstStyle/>
                                        <a:p><a:r><a:t>Cell one</a:t></a:r></a:p>
                                    </a:txBody>
                                    <a:tcPr/>
                                </a:tc>
                                <a:tc gridSpan="2">
                                    <a:txBody>
                                        <a:bodyPr/>
                                        <a:lstStyle/>
                                        <a:p><a:r><a:t>Wide cell</a:t></a:r>
                                        </a:p>
                                    </a:txBody>
                                    <a:tcPr/>
                                </a:tc>
                                <a:tc hMerge="1">
                                    <a:txBody>
                                        <a:bodyPr/>
                                        <a:lstStyle/>
                                        <a:p/>
                                    </a:txBody>
                                    <a:tcPr/>
                                </a:tc>
                            </a:tr>
                            <a:tr h="304800">
                                <a:tc rowSpan="2" gridSpan="2">
                                    <a:txBody>
                                        <a:bodyPr/>
                                        <a:lstStyle/>
                                        <a:p><a:r><a:t>Big cell</a:t></a:r></a:p>
                                    </a:txBody>
                                    <a:tcPr/>
                                </a:tc>
                                <a:tc rowSpan="2" hMerge="1">
                                    <a:txBody>
                                        <a:bodyPr/>
                                        <a:lstStyle/>
                                        <a:p/>
                                    </a:txBody>
                                    <a:tcPr/>
                                </a:tc>
                                <a:tc>
                                    <a:txBody>
                                        <a:bodyPr/>
                                        <a:lstStyle/>
                                        <a:p><a:r><a:t>Cell two</a:t></a:r></a:p>
                                    </a:txBody>
                                    <a:tcPr/>
                                </a:tc>
                            </a:tr>
                            <a:tr h="304800">
                                <a:tc gridSpan="2" vMerge="1">
                                    <a:txBody>
                                        <a:bodyPr/>
                                        <a:lstStyle/>
                                        <a:p/>
                                    </a:txBody>
                                    <a:tcPr/>
                                </a:tc>
                                <a:tc hMerge="1" vMerge="1">
                                    <a:txBody>
                                        <a:bodyPr/>
                                        <a:lstStyle/>
                                        <a:p/>
                                    </a:txBody>
                                    <a:tcPr/>
                                </a:tc>
                                <a:tc>
                                    <a:txBody>
                                        <a:bodyPr/>
                                        <a:lstStyle/>
                                        <a:p><a:r><a:t>Cell three</a:t></a:r></a:p>
                                    </a:txBody>
                                    <a:tcPr/>
                                </a:tc>
                            </a:tr>
                        </a:tbl>
                    </a:graphicData>
                </a:graphic>
            </p:graphicFrame>
        </p:spTree>
    </p:cSld>
    <p:clrMapOvr>
        <a:masterClrMapping/>
    </p:clrMapOvr>
</p:sld>
        """,
        add_suffix=False,
    )
